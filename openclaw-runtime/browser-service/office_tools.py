#!/usr/bin/env python3

import csv
import html
import io
import json
import os
import posixpath
import re
import secrets
import shutil
import subprocess
import sys
import tempfile
import time
import zipfile
from typing import Dict, Iterable, List, Optional, Tuple
from xml.etree import ElementTree as ET

try:
    from openpyxl import Workbook, load_workbook
except ImportError:  # pragma: no cover - runtime dependency may be absent in dev
    Workbook = None
    load_workbook = None

try:
    from pptx import Presentation as PptxPresentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches
except ImportError:  # pragma: no cover - runtime dependency may be absent in dev
    PptxPresentation = None
    CategoryChartData = None
    RGBColor = None
    MSO_AUTO_SHAPE_TYPE = None
    Inches = None

try:
    from docling.document_converter import DocumentConverter
except ImportError:  # pragma: no cover - runtime dependency may be absent in dev
    DocumentConverter = None

WORKSPACE_ROOT = os.environ.get("ENTROPIC_WORKSPACE_PATH", "/data/workspace")
DESKTOP_ACTION_QUEUE_DIR = os.environ.get("ENTROPIC_DESKTOP_ACTION_QUEUE_DIR", "/data/browser/desktop-actions")

AIO_SPEC = "agent-interpretable-object"
AIO_ROOT = "Agent-Interpretable-Object"
AIO_MANIFEST = f"{AIO_ROOT}/manifest.yaml"
AIO_CORE_KEYWORDS = f"{AIO_ROOT}/core/keywords.yaml"
AIO_CORE_ONTOLOGY = f"{AIO_ROOT}/core/ontology.yaml"
AIO_CORE_CANONICAL_INSTANCE = f"{AIO_ROOT}/core/canonical-instance.yaml"
AIO_FAMILY_TABULAR = f"{AIO_ROOT}/families/tabular-space.yaml"
AIO_FAMILY_LINEAR_TEXT = f"{AIO_ROOT}/families/linear-text.yaml"
AIO_FAMILY_SLIDE_SPACE = f"{AIO_ROOT}/families/slide-space.yaml"
AIO_FAMILY_GRAPHICS_SCENE = f"{AIO_ROOT}/families/graphics-scene.md"
AIO_KIND_SPREADSHEET = f"{AIO_ROOT}/kinds/spreadsheet.yaml"
AIO_KIND_DOCUMENT = f"{AIO_ROOT}/kinds/document.yaml"
AIO_KIND_PRESENTATION = f"{AIO_ROOT}/kinds/presentation.yaml"

XLSX_NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
DOCX_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
DOC_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

PANDOC_DOCUMENT_EXTENSIONS = {
    ".md": "markdown",
    ".markdown": "markdown",
    ".qmd": "markdown",
    ".txt": "markdown",
    ".html": "html",
    ".htm": "html",
    ".rst": "rst",
    ".org": "org",
    ".tex": "latex",
    ".latex": "latex",
    ".adoc": "asciidoc",
    ".asciidoc": "asciidoc",
}

PANDOC_OUTPUT_FORMATS = {
    ".md": "gfm",
    ".markdown": "markdown",
    ".qmd": "markdown",
    ".txt": "plain",
    ".html": "html5",
    ".htm": "html5",
    ".rst": "rst",
    ".org": "org",
    ".tex": "latex",
    ".latex": "latex",
    ".adoc": "asciidoc",
    ".asciidoc": "asciidoc",
}

DOCLING_DOCUMENT_EXTENSIONS = {
    ".pdf",
    ".png",
    ".jpg",
    ".jpeg",
    ".tif",
    ".tiff",
    ".bmp",
    ".webp",
}

DOCLING_INPUT_FORMATS = {
    ".pdf": "pdf",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".tif": "image",
    ".tiff": "image",
    ".bmp": "image",
    ".webp": "image",
}

DOCLING_CONVERTER: Optional[object] = None
BACKEND_PROJECTION_MAX_BYTES = 250_000

ET.register_namespace("", XLSX_NS)
ET.register_namespace("a", DRAWING_NS)
ET.register_namespace("p", PPTX_NS)
ET.register_namespace("r", DOC_REL_NS)
ET.register_namespace("w", DOCX_NS)


def xlsx_tag(name: str) -> str:
    return f"{{{XLSX_NS}}}{name}"


def docx_tag(name: str) -> str:
    return f"{{{DOCX_NS}}}{name}"


def pptx_tag(name: str) -> str:
    return f"{{{PPTX_NS}}}{name}"


def drawing_tag(name: str) -> str:
    return f"{{{DRAWING_NS}}}{name}"


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def normalize_workspace_root() -> str:
    return posixpath.normpath(WORKSPACE_ROOT)


def resolve_workspace_path(raw_path: str) -> str:
    if not isinstance(raw_path, str) or not raw_path.strip():
        raise ValueError("A workspace file path is required.")
    root = normalize_workspace_root()
    trimmed = raw_path.strip()
    if trimmed == root or trimmed.startswith(f"{root}/"):
        full_path = posixpath.normpath(trimmed)
    else:
        normalized_relative = posixpath.normpath(f"/{trimmed}").lstrip("/")
        full_path = posixpath.join(root, normalized_relative)
    if full_path != root and not full_path.startswith(f"{root}/"):
        raise ValueError("The requested path is outside /data/workspace.")
    return full_path


def workspace_relative_path(raw_path: str) -> str:
    full_path = resolve_workspace_path(raw_path)
    root = normalize_workspace_root()
    if full_path == root:
        return ""
    return full_path[len(root) + 1 :]


def path_metadata(path: str) -> Dict[str, object]:
    exists = os.path.exists(path)
    if not exists:
        return {
            "exists": False,
            "modifiedMs": None,
            "size": 0,
            "etag": "missing",
        }
    stats = os.stat(path)
    modified_ms = int(round(stats.st_mtime_ns / 1_000_000))
    return {
        "exists": True,
        "modifiedMs": modified_ms,
        "size": int(stats.st_size),
        "etag": f"{modified_ms}:{stats.st_size}",
    }


def assert_expected_etag(path: str, expected_etag: Optional[str]) -> None:
    if not expected_etag:
        return
    current = path_metadata(path)["etag"]
    if current != expected_etag:
        raise RuntimeError(
            "The file changed on disk while you were editing it. Reload the viewer before saving again."
        )


def ensure_parent_dir(path: str) -> None:
    os.makedirs(posixpath.dirname(path), exist_ok=True)


def atomic_write_bytes(path: str, data: bytes) -> None:
    ensure_parent_dir(path)
    temp_fd, temp_path = tempfile.mkstemp(prefix=".entropic-office-", dir=posixpath.dirname(path))
    try:
        with os.fdopen(temp_fd, "wb") as handle:
            handle.write(data)
        os.replace(temp_path, path)
    finally:
        try:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
        except OSError:
            pass


def atomic_write_text(path: str, text: str) -> None:
    atomic_write_bytes(path, text.encode("utf-8"))


def split_extension(path: str) -> str:
    return posixpath.splitext(path)[1].lower()


def basename_without_extension(path: str) -> str:
    return posixpath.splitext(posixpath.basename(path))[0] or "object"


def slugify_identifier(value: object, fallback: str) -> str:
    text = str(value or "").strip().lower()
    text = re.sub(r"[^a-z0-9]+", "-", text).strip("-")
    return text or fallback


def pandoc_available() -> bool:
    return bool(shutil.which("pandoc"))


def docling_available() -> bool:
    return DocumentConverter is not None


def text_document_format_name(path: str) -> str:
    extension = split_extension(path)
    return extension.lstrip(".") or "txt"


def pandoc_input_format_for_extension(extension: str) -> Optional[str]:
    return PANDOC_DOCUMENT_EXTENSIONS.get(extension.lower())


def pandoc_output_format_for_extension(extension: str) -> Optional[str]:
    return PANDOC_OUTPUT_FORMATS.get(extension.lower())


def docling_input_format_for_extension(extension: str) -> Optional[str]:
    return DOCLING_INPUT_FORMATS.get(extension.lower())


def run_subprocess_text(argv: List[str], stdin_text: Optional[str] = None) -> str:
    try:
        result = subprocess.run(
            argv,
            input=stdin_text,
            capture_output=True,
            check=False,
            text=True,
        )
    except FileNotFoundError as error:
        raise RuntimeError(f"Required command is not installed: `{argv[0]}`.") from error
    if result.returncode != 0:
        stderr = (result.stderr or "").strip()
        stdout = (result.stdout or "").strip()
        detail = stderr or stdout or f"exit status {result.returncode}"
        raise RuntimeError(f"`{' '.join(argv)}` failed: {detail}")
    return result.stdout


def run_subprocess_json(argv: List[str], stdin_text: Optional[str] = None) -> object:
    output = run_subprocess_text(argv, stdin_text=stdin_text)
    try:
        return json.loads(output)
    except json.JSONDecodeError as error:
        raise RuntimeError(f"`{' '.join(argv)}` did not return valid JSON.") from error


def read_utf8_text(path: str) -> str:
    with open(path, "r", encoding="utf-8") as handle:
        return handle.read()


def embedded_backend_payload(
    payload: object,
    max_bytes: int = BACKEND_PROJECTION_MAX_BYTES,
) -> Tuple[Optional[object], Dict[str, object]]:
    try:
        serialized = json.dumps(payload, ensure_ascii=False)
    except TypeError:
        serialized = str(payload)
    size_bytes = len(serialized.encode("utf-8"))
    summary: Dict[str, object] = {
        "size_bytes": size_bytes,
        "included": size_bytes <= max_bytes,
    }
    if isinstance(payload, dict):
        summary["keys"] = sorted(str(key) for key in payload.keys())[:20]
    if size_bytes <= max_bytes:
        return payload, summary
    summary["reason"] = "payload exceeds inline projection limit"
    summary["preview"] = serialized[:2000]
    return None, summary


def default_backend_name(kind: str, format_name: str) -> str:
    if format_name == "docx":
        return "native-docx"
    if format_name == "pptx":
        return "native-pptx"
    if format_name == "xlsx":
        return "native-xlsx"
    if format_name == "csv":
        return "native-csv"
    if kind == "document" and f".{format_name}" in PANDOC_DOCUMENT_EXTENSIONS:
        return "pandoc" if pandoc_available() else "native-text"
    if kind == "document" and f".{format_name}" in DOCLING_DOCUMENT_EXTENSIONS:
        if format_name == "pdf":
            return "docling" if docling_available() else "native-pdf"
        return "docling" if docling_available() else "native-image"
    return "native"


def aio_backend_descriptors(kind: str, format_name: str, active_backend: Optional[str]) -> List[Dict[str, object]]:
    active = active_backend or default_backend_name(kind, format_name)
    active_role = "canonical-read-only" if kind == "document" and active == "docling" else "canonical-read-write"
    backends: List[Dict[str, object]] = [
        {
            "name": active,
            "role": active_role,
            "active": True,
            "available": True,
        }
    ]
    if kind == "document":
        if active != "pandoc":
            backends.append(
                {
                    "name": "pandoc",
                    "role": "text-projection",
                    "active": False,
                    "available": pandoc_available(),
                }
            )
        if active != "docling":
            backends.append(
                {
                    "name": "docling",
                    "role": "analysis-adapter",
                    "active": False,
                    "available": docling_available(),
                }
            )
    return backends


def compact_optional_mapping(payload: Optional[Dict[str, object]]) -> Optional[Dict[str, object]]:
    if not isinstance(payload, dict):
        return None
    result = {key: value for key, value in payload.items() if value not in (None, [], {}, "")}
    return result or None


def merge_nested_dict(base: Optional[Dict[str, object]], override: Optional[Dict[str, object]]) -> Dict[str, object]:
    result: Dict[str, object] = dict(base or {})
    for key, value in (override or {}).items():
        existing = result.get(key)
        if isinstance(existing, dict) and isinstance(value, dict):
            result[key] = merge_nested_dict(existing, value)
        else:
            result[key] = value
    return result


def aio_definition_refs(kind: str) -> Dict[str, object]:
    families: List[str] = []
    kind_path = ""
    if kind == "spreadsheet":
        families = [AIO_FAMILY_TABULAR]
        kind_path = AIO_KIND_SPREADSHEET
    elif kind == "document":
        families = [AIO_FAMILY_LINEAR_TEXT]
        kind_path = AIO_KIND_DOCUMENT
    elif kind == "presentation":
        families = [AIO_FAMILY_SLIDE_SPACE, AIO_FAMILY_GRAPHICS_SCENE]
        kind_path = AIO_KIND_PRESENTATION
    return {
        "manifest": AIO_MANIFEST,
        "core": [AIO_CORE_KEYWORDS, AIO_CORE_ONTOLOGY, AIO_CORE_CANONICAL_INSTANCE],
        "families": families,
        "kind": kind_path,
    }


def aio_source_block(document: Dict[str, object], extra_notes: Optional[List[str]] = None) -> Dict[str, object]:
    notes = []
    warning = str(document.get("warning") or "").strip()
    if warning:
        notes.append(warning)
    for note in extra_notes or []:
        normalized = str(note or "").strip()
        if normalized:
            notes.append(normalized)
    return {
        "path": str(document.get("path") or ""),
        "exists": bool(document.get("exists")),
        "etag": str(document.get("etag") or "missing"),
        "modifiedMs": document.get("modifiedMs"),
        "size": int(document.get("size") or 0),
        "notes": notes,
    }


def aio_envelope(
    kind: str,
    document: Dict[str, object],
    object_payload: Dict[str, object],
    capability_set: str,
    extra_notes: Optional[List[str]] = None,
    analysis: Optional[Dict[str, object]] = None,
    annotations: Optional[List[Dict[str, object]]] = None,
    projections: Optional[Dict[str, object]] = None,
    artifacts: Optional[Dict[str, object]] = None,
    backend_name: Optional[str] = None,
) -> Dict[str, object]:
    format_name = str(document.get("format") or "").strip()
    payload = {
        "spec": AIO_SPEC,
        "kind": kind,
        "format": format_name,
        "definitions": aio_definition_refs(kind),
        "realization": {
            "name": "entropic-office",
            "mode": "automation",
            "capability_set": capability_set,
            "backends": aio_backend_descriptors(kind, format_name, backend_name),
        },
        "source": aio_source_block(document, extra_notes=extra_notes),
        "object": object_payload,
    }
    normalized_analysis = compact_optional_mapping(analysis)
    if normalized_analysis is not None:
        payload["analysis"] = normalized_analysis
    if annotations:
        payload["annotations"] = [item for item in annotations if isinstance(item, dict)]
    normalized_projections = compact_optional_mapping(projections)
    if normalized_projections is not None:
        payload["projections"] = normalized_projections
    normalized_artifacts = compact_optional_mapping(artifacts)
    if normalized_artifacts is not None:
        payload["artifacts"] = normalized_artifacts
    return payload


def extract_aio_source_path(payload: Dict[str, object]) -> str:
    source = payload.get("source")
    if isinstance(source, dict):
        return str(source.get("path") or "").strip()
    return ""


def extract_aio_source_etag(payload: Dict[str, object]) -> Optional[str]:
    source = payload.get("source")
    if not isinstance(source, dict):
        return None
    etag = str(source.get("etag") or "").strip()
    return etag or None


def ensure_aio_payload(payload: Dict[str, object]) -> Dict[str, object]:
    if not isinstance(payload, dict) or payload.get("spec") != AIO_SPEC:
        raise RuntimeError("Expected an agent-interpretable-object payload.")
    return payload


def col_index_to_label(index: int) -> str:
    if index < 1:
        raise ValueError("Column index must be positive.")
    label = []
    value = index
    while value > 0:
        value, remainder = divmod(value - 1, 26)
        label.append(chr(65 + remainder))
    return "".join(reversed(label))


def label_to_col_index(label: str) -> int:
    value = 0
    for char in label.upper():
        if not ("A" <= char <= "Z"):
            raise ValueError(f"Invalid column label `{label}`.")
        value = value * 26 + (ord(char) - 64)
    return value


def cell_ref(row: int, col: int) -> str:
    return f"{col_index_to_label(col)}{row}"


def parse_cell_ref(ref: str) -> Tuple[int, int]:
    match = re.fullmatch(r"([A-Za-z]+)(\d+)", ref or "")
    if not match:
        raise ValueError(f"Invalid cell reference `{ref}`.")
    return int(match.group(2)), label_to_col_index(match.group(1))


def is_numeric_value(value: str) -> bool:
    if not isinstance(value, str):
        return False
    trimmed = value.strip()
    if not trimmed:
        return False
    return re.fullmatch(r"[+-]?(?:\d+(?:\.\d+)?|\.\d+)", trimmed) is not None


def coerce_boolean_text(value: str) -> Optional[str]:
    if not isinstance(value, str):
        return None
    lowered = value.strip().lower()
    if lowered in {"true", "1", "yes"}:
        return "1"
    if lowered in {"false", "0", "no"}:
        return "0"
    return None


def safe_inline_text(parent: ET.Element, value: str) -> None:
    inline = ET.SubElement(parent, xlsx_tag("is"))
    text = ET.SubElement(inline, xlsx_tag("t"))
    text.text = value


def read_shared_strings(archive: zipfile.ZipFile) -> List[str]:
    if "xl/sharedStrings.xml" not in archive.namelist():
        return []
    root = ET.fromstring(archive.read("xl/sharedStrings.xml"))
    values: List[str] = []
    for item in root:
        if local_name(item.tag) != "si":
            continue
        text_parts = []
        for text_node in item.iter():
            if local_name(text_node.tag) == "t":
                text_parts.append(text_node.text or "")
        values.append("".join(text_parts))
    return values


def decode_xlsx_cell(cell: ET.Element, shared_strings: List[str]) -> Dict[str, object]:
    ref = cell.attrib.get("r", "")
    row, col = parse_cell_ref(ref)
    formula_node = cell.find(xlsx_tag("f"))
    value_node = cell.find(xlsx_tag("v"))
    cell_type = cell.attrib.get("t")
    formula = formula_node.text or "" if formula_node is not None and formula_node.text else None
    raw_value = value_node.text or "" if value_node is not None and value_node.text else ""
    if cell_type == "s":
        try:
            display_value = shared_strings[int(raw_value)]
        except (IndexError, ValueError):
            display_value = ""
        kind = "string"
    elif cell_type == "inlineStr":
        text_parts = []
        inline = cell.find(xlsx_tag("is"))
        if inline is not None:
            for node in inline.iter():
                if local_name(node.tag) == "t":
                    text_parts.append(node.text or "")
        display_value = "".join(text_parts)
        kind = "string"
    elif cell_type == "b":
        display_value = "TRUE" if raw_value == "1" else "FALSE"
        kind = "boolean"
    else:
        display_value = raw_value
        kind = "number" if raw_value and is_numeric_value(raw_value) else "string"
    return {
        "ref": ref,
        "row": row,
        "col": col,
        "formula": formula,
        "value": display_value,
        "display": display_value if formula is None else (raw_value or display_value),
        "kind": "formula" if formula else kind,
    }


def parse_xlsx_row_metadata(root: ET.Element) -> List[Dict[str, object]]:
    metadata: List[Dict[str, object]] = []
    sheet_data = root.find(xlsx_tag("sheetData"))
    if sheet_data is None:
        return metadata
    for row in sheet_data.findall(xlsx_tag("row")):
        row_index = int(row.attrib.get("r") or 0)
        if row_index <= 0:
            continue
        entry: Dict[str, object] = {"index": row_index}
        height = row.attrib.get("ht")
        if height:
            try:
                entry["height"] = float(height)
            except ValueError:
                pass
        if row.attrib.get("hidden") in {"1", "true", "True"}:
            entry["hidden"] = True
        if len(entry) > 1:
            metadata.append(entry)
    return metadata


def parse_xlsx_column_metadata(root: ET.Element) -> List[Dict[str, object]]:
    metadata: List[Dict[str, object]] = []
    cols = root.find(xlsx_tag("cols"))
    if cols is None:
        return metadata
    for column in cols.findall(xlsx_tag("col")):
        try:
            min_index = int(column.attrib.get("min") or 0)
            max_index = int(column.attrib.get("max") or min_index)
        except ValueError:
            continue
        if min_index <= 0 or max_index < min_index:
            continue
        width = column.attrib.get("width")
        hidden = column.attrib.get("hidden") in {"1", "true", "True"}
        for index in range(min_index, max_index + 1):
            entry: Dict[str, object] = {"index": index, "label": col_index_to_label(index)}
            if width:
                try:
                    entry["width"] = float(width)
                except ValueError:
                    pass
            if hidden:
                entry["hidden"] = True
            if len(entry) > 2 or entry.get("hidden"):
                metadata.append(entry)
    return metadata


def parse_xlsx_merged_ranges(root: ET.Element) -> List[Dict[str, object]]:
    merge_cells = root.find(xlsx_tag("mergeCells"))
    if merge_cells is None:
        return []
    merged_ranges: List[Dict[str, object]] = []
    for index, merge_cell in enumerate(merge_cells.findall(xlsx_tag("mergeCell")), start=1):
        ref = str(merge_cell.attrib.get("ref") or "").strip().upper()
        if not ref:
            continue
        merged_ranges.append(
            {
                "id": f"merged_range:{index}",
                "kind": "merged_range",
                "ref": ref,
            }
        )
    return merged_ranges


def parse_xlsx_freeze_pane(root: ET.Element) -> Optional[Dict[str, object]]:
    sheet_views = root.find(xlsx_tag("sheetViews"))
    if sheet_views is None:
        return None
    for sheet_view in sheet_views.findall(xlsx_tag("sheetView")):
        pane = sheet_view.find(xlsx_tag("pane"))
        if pane is None:
            continue
        state = str(pane.attrib.get("state") or "").strip()
        if state and state not in {"frozen", "frozenSplit"}:
            continue
        ref = str(pane.attrib.get("topLeftCell") or "").strip().upper()
        if not ref:
            try:
                x_split = int(float(str(pane.attrib.get("xSplit") or "0")))
                y_split = int(float(str(pane.attrib.get("ySplit") or "0")))
            except ValueError:
                x_split = 0
                y_split = 0
            if x_split > 0 or y_split > 0:
                ref = cell_ref(max(y_split + 1, 1), max(x_split + 1, 1))
        if ref:
            return {"ref": ref}
    return None


def normalize_sheet_dimension_rows(raw_rows: object) -> List[Dict[str, object]]:
    normalized: List[Dict[str, object]] = []
    if not isinstance(raw_rows, list):
        return normalized
    for row in raw_rows:
        if not isinstance(row, dict):
            continue
        index = int(row.get("index") or row.get("row") or 0)
        if index <= 0:
            continue
        entry: Dict[str, object] = {"index": index}
        if row.get("height") is not None:
            try:
                entry["height"] = float(row.get("height") or 0)
            except (TypeError, ValueError):
                pass
        if row.get("hidden") is not None:
            entry["hidden"] = bool(row.get("hidden"))
        if len(entry) > 1:
            normalized.append(entry)
    normalized.sort(key=lambda item: int(item.get("index") or 0))
    return normalized


def normalize_sheet_dimension_columns(raw_columns: object) -> List[Dict[str, object]]:
    normalized: List[Dict[str, object]] = []
    if not isinstance(raw_columns, list):
        return normalized
    for column in raw_columns:
        if not isinstance(column, dict):
            continue
        index = int(column.get("index") or 0)
        if index <= 0:
            label = str(column.get("label") or "").strip().upper()
            if label:
                try:
                    index = label_to_col_index(label)
                except ValueError:
                    index = 0
        if index <= 0:
            continue
        entry: Dict[str, object] = {
            "index": index,
            "label": col_index_to_label(index),
        }
        if column.get("width") is not None:
            try:
                entry["width"] = float(column.get("width") or 0)
            except (TypeError, ValueError):
                pass
        if column.get("hidden") is not None:
            entry["hidden"] = bool(column.get("hidden"))
        if len(entry) > 2 or entry.get("hidden"):
            normalized.append(entry)
    normalized.sort(key=lambda item: int(item.get("index") or 0))
    return normalized


def normalize_sheet_merged_ranges(raw_ranges: object) -> List[Dict[str, object]]:
    normalized: List[Dict[str, object]] = []
    if not isinstance(raw_ranges, list):
        return normalized
    for index, raw_range in enumerate(raw_ranges, start=1):
        if isinstance(raw_range, dict):
            ref = str(raw_range.get("ref") or "").strip().upper()
            range_id = str(raw_range.get("id") or f"merged_range:{index}")
        else:
            ref = str(raw_range or "").strip().upper()
            range_id = f"merged_range:{index}"
        if not ref:
            continue
        normalized.append({"id": range_id, "kind": "merged_range", "ref": ref})
    return normalized


def normalize_sheet_freeze_pane(raw_freeze: object) -> Optional[Dict[str, object]]:
    if isinstance(raw_freeze, dict):
        ref = str(raw_freeze.get("ref") or "").strip().upper()
    else:
        ref = str(raw_freeze or "").strip().upper()
    if not ref:
        return None
    return {"ref": ref}


def parse_xlsx_sheet(sheet_name: str, payload: bytes, shared_strings: List[str]) -> Dict[str, object]:
    root = ET.fromstring(payload)
    sheet_data = root.find(xlsx_tag("sheetData"))
    auto_filter = root.find(xlsx_tag("autoFilter"))
    cells = []
    max_row = 0
    max_col = 0
    if sheet_data is not None:
        for row in sheet_data.findall(xlsx_tag("row")):
            for cell in row.findall(xlsx_tag("c")):
                decoded = decode_xlsx_cell(cell, shared_strings)
                cells.append(decoded)
                max_row = max(max_row, int(decoded["row"]))
                max_col = max(max_col, int(decoded["col"]))
    return {
        "name": sheet_name,
        "cells": cells,
        "rowCount": max_row,
        "colCount": max_col,
        "rows": parse_xlsx_row_metadata(root),
        "columns": parse_xlsx_column_metadata(root),
        "mergedRanges": parse_xlsx_merged_ranges(root),
        "freezePane": parse_xlsx_freeze_pane(root),
        "filterRef": str((auto_filter.attrib.get("ref") if auto_filter is not None else "") or "").strip().upper() or None,
    }


def normalize_relationship_target(base_dir: str, target: str) -> str:
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(base_dir, target))


def read_xlsx_document(path: str) -> Dict[str, object]:
    metadata = path_metadata(path)
    if not metadata["exists"]:
        return {
            "kind": "spreadsheet",
            "format": "xlsx",
            "path": path,
            **metadata,
            "warning": None,
            "sheets": [{"name": "Sheet1", "cells": [], "rowCount": 0, "colCount": 0}],
        }
    with zipfile.ZipFile(path, "r") as archive:
        workbook = ET.fromstring(archive.read("xl/workbook.xml"))
        rels = ET.fromstring(archive.read("xl/_rels/workbook.xml.rels"))
        targets_by_id: Dict[str, str] = {}
        for rel in rels:
            if local_name(rel.tag) != "Relationship":
                continue
            rel_id = rel.attrib.get("Id")
            target = rel.attrib.get("Target")
            if rel_id and target:
                targets_by_id[rel_id] = normalize_relationship_target("xl", target)
        shared_strings = read_shared_strings(archive)
        sheets: List[Dict[str, object]] = []
        sheets_parent = workbook.find(xlsx_tag("sheets"))
        if sheets_parent is not None:
            for sheet in sheets_parent.findall(xlsx_tag("sheet")):
                rel_id = sheet.attrib.get(f"{{{DOC_REL_NS}}}id", "")
                target = targets_by_id.get(rel_id)
                if not target:
                    continue
                if target not in archive.namelist():
                    continue
                sheets.append(
                    parse_xlsx_sheet(
                        sheet.attrib.get("name", f"Sheet{len(sheets) + 1}"),
                        archive.read(target),
                        shared_strings,
                    )
                )
        if not sheets:
            sheets = [{"name": "Sheet1", "cells": [], "rowCount": 0, "colCount": 0}]
        return {
            "kind": "spreadsheet",
            "format": "xlsx",
            "path": path,
            **metadata,
            "warning": None,
            "sheets": sheets,
        }


def read_csv_document(path: str) -> Dict[str, object]:
    metadata = path_metadata(path)
    if not metadata["exists"]:
        return {
            "kind": "spreadsheet",
            "format": "csv",
            "path": path,
            **metadata,
            "warning": None,
            "sheets": [{"name": "Sheet1", "cells": [], "rowCount": 0, "colCount": 0}],
        }
    with open(path, "r", encoding="utf-8", newline="") as handle:
        reader = csv.reader(handle)
        rows = list(reader)
    cells = []
    max_col = 0
    for row_index, row in enumerate(rows, start=1):
        max_col = max(max_col, len(row))
        for col_index, value in enumerate(row, start=1):
            if value == "":
                continue
            cells.append(
                {
                    "ref": cell_ref(row_index, col_index),
                    "row": row_index,
                    "col": col_index,
                    "formula": value[1:] if value.startswith("=") else None,
                    "value": value,
                    "display": value,
                    "kind": "formula" if value.startswith("=") else "string",
                }
            )
    return {
        "kind": "spreadsheet",
        "format": "csv",
        "path": path,
        **metadata,
        "warning": None,
        "sheets": [{"name": "Sheet1", "cells": cells, "rowCount": len(rows), "colCount": max_col}],
    }


def legacy_minimal_xlsx(path: str) -> bool:
    if not os.path.exists(path):
        return False
    try:
        with zipfile.ZipFile(path, "r") as archive:
            entries = set(archive.namelist())
    except (OSError, zipfile.BadZipFile):
        return False
    required_minimal = {
        "[Content_Types].xml",
        "_rels/.rels",
        "xl/workbook.xml",
        "xl/_rels/workbook.xml.rels",
    }
    if not required_minimal.issubset(entries):
        return False
    has_sheet = any(name.startswith("xl/worksheets/") and name.endswith(".xml") for name in entries)
    if not has_sheet:
        return False
    # Older Entropic-generated workbooks omitted these standard parts, which ONLYOFFICE
    # rejects even though the file is readable by the lightweight inspector.
    missing_standard_parts = {
        "docProps/app.xml",
        "docProps/core.xml",
        "xl/styles.xml",
        "xl/theme/theme1.xml",
    }
    return bool(entries.isdisjoint(missing_standard_parts))


def normalize_sheet_payload(sheet: Dict[str, object]) -> Dict[str, object]:
    name = str(sheet.get("name") or "Sheet1").strip() or "Sheet1"
    cleaned_cells = []
    max_row = 0
    max_col = 0
    for raw_cell in sheet.get("cells") or []:
        if not isinstance(raw_cell, dict):
            continue
        raw_ref = str(raw_cell.get("ref") or "").strip().upper()
        if raw_ref:
            row, col = parse_cell_ref(raw_ref)
            ref = raw_ref
        else:
            row = int(raw_cell.get("row") or 0)
            col = int(raw_cell.get("col") or 0)
            if row <= 0 or col <= 0:
                continue
            ref = cell_ref(row, col)
        formula = str(raw_cell.get("formula") or "").strip()
        if formula.startswith("="):
            formula = formula[1:]
        value = str(raw_cell.get("value") or "")
        if formula:
            cleaned_cells.append(
                {
                    "ref": ref,
                    "row": row,
                    "col": col,
                    "formula": formula,
                    "value": value,
                    "display": str(raw_cell.get("display") or value),
                    "kind": "formula",
                }
            )
        elif value != "":
            kind = str(raw_cell.get("kind") or "")
            if kind not in {"number", "boolean", "string"}:
                if is_numeric_value(value):
                    kind = "number"
                elif coerce_boolean_text(value) is not None:
                    kind = "boolean"
                else:
                    kind = "string"
            cleaned_cells.append(
                {
                    "ref": ref,
                    "row": row,
                    "col": col,
                    "formula": None,
                    "value": value,
                    "display": str(raw_cell.get("display") or value),
                    "kind": kind,
                }
            )
        max_row = max(max_row, row)
        max_col = max(max_col, col)
    cleaned_cells.sort(key=lambda cell: (int(cell["row"]), int(cell["col"])))
    normalized_rows = normalize_sheet_dimension_rows(sheet.get("rows"))
    normalized_columns = normalize_sheet_dimension_columns(sheet.get("columns"))
    merged_ranges = normalize_sheet_merged_ranges(sheet.get("mergedRanges") or sheet.get("merged_ranges"))
    freeze_pane = normalize_sheet_freeze_pane(sheet.get("freezePane") or sheet.get("freeze_pane"))
    filter_ref = str(sheet.get("filterRef") or sheet.get("filter_ref") or "").strip().upper() or None
    max_row = max(
        max_row,
        max((int(entry.get("index") or 0) for entry in normalized_rows), default=0),
    )
    max_col = max(
        max_col,
        max((int(entry.get("index") or 0) for entry in normalized_columns), default=0),
    )
    normalized: Dict[str, object] = {
        "name": name,
        "cells": cleaned_cells,
        "rowCount": max_row,
        "colCount": max_col,
    }
    if "rows" in sheet:
        normalized["rows"] = normalized_rows
    elif normalized_rows:
        normalized["rows"] = normalized_rows
    if "columns" in sheet:
        normalized["columns"] = normalized_columns
    elif normalized_columns:
        normalized["columns"] = normalized_columns
    if "mergedRanges" in sheet or "merged_ranges" in sheet:
        normalized["mergedRanges"] = merged_ranges
    elif merged_ranges:
        normalized["mergedRanges"] = merged_ranges
    if "freezePane" in sheet or "freeze_pane" in sheet or freeze_pane:
        normalized["freezePane"] = freeze_pane
    if "filterRef" in sheet or "filter_ref" in sheet or filter_ref:
        normalized["filterRef"] = filter_ref
    return normalized


def openpyxl_available() -> bool:
    return Workbook is not None and load_workbook is not None


def coerce_openpyxl_value(cell: Dict[str, object]):
    formula = str(cell.get("formula") or "").strip()
    if formula:
        return f"={formula}"
    value = str(cell.get("value") or "")
    kind = str(cell.get("kind") or "string")
    if kind == "boolean":
        return coerce_boolean_text(value) == "1"
    if kind == "number" and is_numeric_value(value):
        trimmed = value.strip()
        if re.fullmatch(r"[+-]?\d+", trimmed):
            try:
                return int(trimmed)
            except ValueError:
                return trimmed
        try:
            return float(trimmed)
        except ValueError:
            return trimmed
    return value


def build_sheet_root(sheet: Dict[str, object], existing_bytes: Optional[bytes]) -> bytes:
    existing_rows: Dict[int, ET.Element] = {}
    existing_cells: Dict[str, ET.Element] = {}
    if existing_bytes:
        root = ET.fromstring(existing_bytes)
    else:
        root = ET.Element(xlsx_tag("worksheet"))
        ET.SubElement(root, xlsx_tag("sheetData"))
    sheet_data = root.find(xlsx_tag("sheetData"))
    if sheet_data is None:
        sheet_data = ET.SubElement(root, xlsx_tag("sheetData"))
    for row in sheet_data.findall(xlsx_tag("row")):
        row_index = int(row.attrib.get("r") or 0)
        if row_index > 0:
            existing_rows[row_index] = row
        for cell in row.findall(xlsx_tag("c")):
            cell_key = cell.attrib.get("r")
            if cell_key:
                existing_cells[cell_key] = cell
    for child in list(sheet_data):
        sheet_data.remove(child)

    rows_by_index: Dict[int, List[Dict[str, object]]] = {}
    for cell in sheet["cells"]:
        rows_by_index.setdefault(int(cell["row"]), []).append(cell)

    row_metadata_by_index = {
        int(entry.get("index") or 0): entry
        for entry in sheet.get("rows") or []
        if isinstance(entry, dict) and int(entry.get("index") or 0) > 0
    }
    all_row_indexes = sorted(set(rows_by_index) | set(row_metadata_by_index))

    for row_index in all_row_indexes:
        row_element = ET.Element(xlsx_tag("row"))
        if row_index in existing_rows:
            row_element.attrib.update(existing_rows[row_index].attrib)
        row_element.attrib["r"] = str(row_index)
        row_metadata = row_metadata_by_index.get(row_index)
        if row_metadata is None and "rows" in sheet:
            for key in ("ht", "hidden", "customHeight"):
                row_element.attrib.pop(key, None)
        elif row_metadata is not None:
            if row_metadata.get("height") is not None:
                row_element.attrib["ht"] = str(row_metadata.get("height"))
                row_element.attrib["customHeight"] = "1"
            elif "rows" in sheet:
                row_element.attrib.pop("ht", None)
                row_element.attrib.pop("customHeight", None)
            if row_metadata.get("hidden") is True:
                row_element.attrib["hidden"] = "1"
            elif "hidden" in row_metadata or "rows" in sheet:
                row_element.attrib.pop("hidden", None)
        for cell in rows_by_index.get(row_index, []):
            ref = str(cell["ref"])
            existing_cell = existing_cells.get(ref)
            cell_element = ET.Element(xlsx_tag("c"))
            if existing_cell is not None:
                for key, value in existing_cell.attrib.items():
                    if key not in {"r", "t"}:
                        cell_element.attrib[key] = value
            cell_element.attrib["r"] = ref
            formula = cell.get("formula")
            value = str(cell.get("value") or "")
            if formula:
                ET.SubElement(cell_element, xlsx_tag("f")).text = str(formula)
                cached = str(cell.get("display") or value)
                if cached != "":
                    boolean_value = coerce_boolean_text(cached)
                    if boolean_value is not None:
                        cell_element.attrib["t"] = "b"
                        ET.SubElement(cell_element, xlsx_tag("v")).text = boolean_value
                    elif is_numeric_value(cached):
                        ET.SubElement(cell_element, xlsx_tag("v")).text = cached.strip()
                    else:
                        cell_element.attrib["t"] = "str"
                        ET.SubElement(cell_element, xlsx_tag("v")).text = cached
            else:
                kind = str(cell.get("kind") or "string")
                if kind == "boolean":
                    boolean_value = coerce_boolean_text(value) or "0"
                    cell_element.attrib["t"] = "b"
                    ET.SubElement(cell_element, xlsx_tag("v")).text = boolean_value
                elif kind == "number" and is_numeric_value(value):
                        ET.SubElement(cell_element, xlsx_tag("v")).text = value.strip()
                else:
                    cell_element.attrib["t"] = "inlineStr"
                    safe_inline_text(cell_element, value)
            row_element.append(cell_element)
        sheet_data.append(row_element)

    column_metadata = [
        entry
        for entry in sheet.get("columns") or []
        if isinstance(entry, dict) and int(entry.get("index") or 0) > 0
    ]
    existing_cols = root.find(xlsx_tag("cols"))
    if existing_cols is not None:
        root.remove(existing_cols)
    if column_metadata:
        cols = ET.Element(xlsx_tag("cols"))
        for entry in column_metadata:
            index = int(entry.get("index") or 0)
            if index <= 0:
                continue
            column = ET.SubElement(cols, xlsx_tag("col"))
            column.attrib["min"] = str(index)
            column.attrib["max"] = str(index)
            if entry.get("width") is not None:
                column.attrib["width"] = str(entry.get("width"))
                column.attrib["customWidth"] = "1"
            if entry.get("hidden") is True:
                column.attrib["hidden"] = "1"
        insert_index = list(root).index(sheet_data) if sheet_data in list(root) else len(list(root))
        root.insert(insert_index, cols)

    existing_merge_cells = root.find(xlsx_tag("mergeCells"))
    if existing_merge_cells is not None:
        root.remove(existing_merge_cells)
    merged_ranges = [
        entry
        for entry in sheet.get("mergedRanges") or []
        if isinstance(entry, dict) and str(entry.get("ref") or "").strip()
    ]
    if merged_ranges or "mergedRanges" in sheet:
        if merged_ranges:
            merge_cells = ET.Element(xlsx_tag("mergeCells"))
            merge_cells.attrib["count"] = str(len(merged_ranges))
            for entry in merged_ranges:
                merge_cell = ET.SubElement(merge_cells, xlsx_tag("mergeCell"))
                merge_cell.attrib["ref"] = str(entry.get("ref") or "").strip().upper()
            root.append(merge_cells)

    existing_auto_filter = root.find(xlsx_tag("autoFilter"))
    if existing_auto_filter is not None:
        root.remove(existing_auto_filter)
    filter_ref = str(sheet.get("filterRef") or "").strip().upper()
    if filter_ref:
        auto_filter = ET.Element(xlsx_tag("autoFilter"))
        auto_filter.attrib["ref"] = filter_ref
        root.append(auto_filter)

    sheet_views = root.find(xlsx_tag("sheetViews"))
    if sheet_views is None and ("freezePane" in sheet):
        sheet_views = ET.Element(xlsx_tag("sheetViews"))
        sheet_view = ET.SubElement(sheet_views, xlsx_tag("sheetView"))
        sheet_view.attrib["workbookViewId"] = "0"
        insert_index = 1 if root.find(xlsx_tag("dimension")) is not None else 0
        root.insert(insert_index, sheet_views)
    if sheet_views is not None:
        sheet_view = sheet_views.find(xlsx_tag("sheetView"))
        if sheet_view is None:
            sheet_view = ET.SubElement(sheet_views, xlsx_tag("sheetView"))
            sheet_view.attrib["workbookViewId"] = "0"
        existing_pane = sheet_view.find(xlsx_tag("pane"))
        freeze_pane = sheet.get("freezePane")
        if existing_pane is not None:
            sheet_view.remove(existing_pane)
        if isinstance(freeze_pane, dict) and str(freeze_pane.get("ref") or "").strip():
            ref = str(freeze_pane.get("ref") or "").strip().upper()
            try:
                row, col = parse_cell_ref(ref)
            except ValueError:
                row, col = (1, 1)
            pane = ET.SubElement(sheet_view, xlsx_tag("pane"))
            if col > 1:
                pane.attrib["xSplit"] = str(col - 1)
            if row > 1:
                pane.attrib["ySplit"] = str(row - 1)
            pane.attrib["topLeftCell"] = ref
            pane.attrib["state"] = "frozen"

    dimension_ref = "A1"
    if sheet["cells"] or row_metadata_by_index or column_metadata:
        max_row = max(
            [int(cell["row"]) for cell in sheet["cells"]]
            + [int(entry.get("index") or 0) for entry in row_metadata_by_index.values()]
            + [1]
        )
        max_col = max(
            [int(cell["col"]) for cell in sheet["cells"]]
            + [int(entry.get("index") or 0) for entry in column_metadata]
            + [1]
        )
        dimension_ref = f"A1:{cell_ref(max_row, max_col)}"
    dimension = root.find(xlsx_tag("dimension"))
    if dimension is None:
        dimension = ET.Element(xlsx_tag("dimension"))
        root.insert(0, dimension)
    dimension.attrib["ref"] = dimension_ref
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def minimal_xlsx_entries(sheets: List[Dict[str, object]]) -> Dict[str, bytes]:
    workbook = ET.Element(xlsx_tag("workbook"))
    sheets_root = ET.SubElement(workbook, xlsx_tag("sheets"))
    for index, sheet in enumerate(sheets, start=1):
        sheet_el = ET.SubElement(sheets_root, xlsx_tag("sheet"))
        sheet_el.set("name", str(sheet["name"]))
        sheet_el.set("sheetId", str(index))
        sheet_el.set(f"{{{DOC_REL_NS}}}id", f"rId{index}")
    calc_pr = ET.SubElement(workbook, xlsx_tag("calcPr"))
    calc_pr.set("fullCalcOnLoad", "1")
    calc_pr.set("calcMode", "auto")

    rels = ET.Element(f"{{{PKG_REL_NS}}}Relationships")
    for index in range(1, len(sheets) + 1):
        relationship = ET.SubElement(rels, f"{{{PKG_REL_NS}}}Relationship")
        relationship.set("Id", f"rId{index}")
        relationship.set(
            "Type",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet",
        )
        relationship.set("Target", f"worksheets/sheet{index}.xml")

    content_types = ET.Element(f"{{{CONTENT_TYPES_NS}}}Types")
    default_rels = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Default")
    default_rels.set("Extension", "rels")
    default_rels.set("ContentType", "application/vnd.openxmlformats-package.relationships+xml")
    default_xml = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Default")
    default_xml.set("Extension", "xml")
    default_xml.set("ContentType", "application/xml")
    workbook_override = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Override")
    workbook_override.set("PartName", "/xl/workbook.xml")
    workbook_override.set(
        "ContentType",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml",
    )
    for index in range(1, len(sheets) + 1):
        sheet_override = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Override")
        sheet_override.set("PartName", f"/xl/worksheets/sheet{index}.xml")
        sheet_override.set(
            "ContentType",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml",
        )

    root_rels = ET.Element(f"{{{PKG_REL_NS}}}Relationships")
    workbook_rel = ET.SubElement(root_rels, f"{{{PKG_REL_NS}}}Relationship")
    workbook_rel.set("Id", "rId1")
    workbook_rel.set(
        "Type",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument",
    )
    workbook_rel.set("Target", "xl/workbook.xml")

    entries = {
        "[Content_Types].xml": ET.tostring(content_types, encoding="utf-8", xml_declaration=True),
        "_rels/.rels": ET.tostring(root_rels, encoding="utf-8", xml_declaration=True),
        "xl/workbook.xml": ET.tostring(workbook, encoding="utf-8", xml_declaration=True),
        "xl/_rels/workbook.xml.rels": ET.tostring(rels, encoding="utf-8", xml_declaration=True),
    }
    for index, sheet in enumerate(sheets, start=1):
        entries[f"xl/worksheets/sheet{index}.xml"] = build_sheet_root(sheet, None)
    return entries


def write_xlsx_document(path: str, sheets: List[Dict[str, object]], expected_etag: Optional[str]) -> Dict[str, object]:
    assert_expected_etag(path, expected_etag)
    normalized_sheets = [normalize_sheet_payload(sheet) for sheet in sheets] or [
        {"name": "Sheet1", "cells": [], "rowCount": 0, "colCount": 0}
    ]

    if openpyxl_available():
        if os.path.exists(path):
            workbook = load_workbook(path)
            if len(workbook.worksheets) != len(normalized_sheets):
                raise RuntimeError(
                    "Adding or removing worksheets is not supported yet. Edit the existing sheets or create a new workbook."
                )
        else:
            workbook = Workbook()
            while len(workbook.worksheets) < len(normalized_sheets):
                workbook.create_sheet()
            while len(workbook.worksheets) > len(normalized_sheets):
                workbook.remove(workbook.worksheets[-1])

        for index, sheet in enumerate(normalized_sheets):
            worksheet = workbook.worksheets[index]
            worksheet.title = str(sheet["name"])
            max_row = max(worksheet.max_row or 0, int(sheet.get("rowCount") or 0))
            max_col = max(worksheet.max_column or 0, int(sheet.get("colCount") or 0))
            for row in range(1, max_row + 1):
                for col in range(1, max_col + 1):
                    worksheet.cell(row=row, column=col).value = None
            for cell in sheet["cells"]:
                worksheet[str(cell["ref"])].value = coerce_openpyxl_value(cell)
            if "freezePane" in sheet:
                freeze_pane = sheet.get("freezePane")
                worksheet.freeze_panes = str(freeze_pane.get("ref") or "").strip().upper() if isinstance(freeze_pane, dict) else None
            if "filterRef" in sheet:
                worksheet.auto_filter.ref = str(sheet.get("filterRef") or "").strip().upper()
            if "mergedRanges" in sheet:
                for merged_range in list(worksheet.merged_cells.ranges):
                    try:
                        worksheet.unmerge_cells(str(merged_range))
                    except Exception:
                        continue
                for merged_range in sheet.get("mergedRanges") or []:
                    if not isinstance(merged_range, dict):
                        continue
                    ref = str(merged_range.get("ref") or "").strip().upper()
                    if not ref:
                        continue
                    try:
                        worksheet.merge_cells(ref)
                    except Exception:
                        continue
            if "rows" in sheet:
                for index_key, dimension in list(worksheet.row_dimensions.items()):
                    if index_key not in {int(entry.get("index") or 0) for entry in sheet.get("rows") or [] if isinstance(entry, dict)}:
                        dimension.height = None
                        dimension.hidden = False
                for row_meta in sheet.get("rows") or []:
                    if not isinstance(row_meta, dict):
                        continue
                    row_index = int(row_meta.get("index") or 0)
                    if row_index <= 0:
                        continue
                    dimension = worksheet.row_dimensions[row_index]
                    dimension.height = row_meta.get("height")
                    dimension.hidden = bool(row_meta.get("hidden")) if row_meta.get("hidden") is not None else False
            if "columns" in sheet:
                current_columns = {str(key).upper(): dimension for key, dimension in worksheet.column_dimensions.items()}
                desired_columns = {
                    str(col_index_to_label(int(entry.get("index") or 0))).upper(): entry
                    for entry in sheet.get("columns") or []
                    if isinstance(entry, dict) and int(entry.get("index") or 0) > 0
                }
                for label, dimension in current_columns.items():
                    if label not in desired_columns:
                        dimension.width = None
                        dimension.hidden = False
                for label, column_meta in desired_columns.items():
                    dimension = worksheet.column_dimensions[label]
                    dimension.width = column_meta.get("width")
                    dimension.hidden = bool(column_meta.get("hidden")) if column_meta.get("hidden") is not None else False

        temp_fd, temp_path = tempfile.mkstemp(prefix=".entropic-office-", dir=posixpath.dirname(path))
        os.close(temp_fd)
        try:
            workbook.save(temp_path)
            os.replace(temp_path, path)
        finally:
            try:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
            except OSError:
                pass
        return read_xlsx_document(path)

    if os.path.exists(path):
        with zipfile.ZipFile(path, "r") as archive:
            entries = {name: archive.read(name) for name in archive.namelist()}
        workbook_root = ET.fromstring(entries["xl/workbook.xml"])
        rels_root = ET.fromstring(entries["xl/_rels/workbook.xml.rels"])
        workbook_sheets = workbook_root.find(xlsx_tag("sheets"))
        if workbook_sheets is None:
            raise RuntimeError("The workbook is missing sheet metadata.")
        existing_names = workbook_sheets.findall(xlsx_tag("sheet"))
        if len(existing_names) != len(normalized_sheets):
            raise RuntimeError(
                "Adding or removing worksheets is not supported yet. Edit the existing sheets or create a new workbook."
            )
        rel_targets: Dict[str, str] = {}
        for rel in rels_root:
            if local_name(rel.tag) != "Relationship":
                continue
            rel_id = rel.attrib.get("Id")
            target = rel.attrib.get("Target")
            if rel_id and target:
                rel_targets[rel_id] = normalize_relationship_target("xl", target)
        for index, sheet in enumerate(existing_names):
            rel_id = sheet.attrib.get(f"{{{DOC_REL_NS}}}id", "")
            target = rel_targets.get(rel_id)
            if not target:
                continue
            sheet.set("name", str(normalized_sheets[index]["name"]))
            entries[target] = build_sheet_root(normalized_sheets[index], entries.get(target))
        entries["xl/workbook.xml"] = ET.tostring(workbook_root, encoding="utf-8", xml_declaration=True)
    else:
        entries = minimal_xlsx_entries(normalized_sheets)

    temp_fd, temp_path = tempfile.mkstemp(prefix=".entropic-office-", dir=posixpath.dirname(path))
    os.close(temp_fd)
    try:
        with zipfile.ZipFile(temp_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for name in sorted(entries):
                archive.writestr(name, entries[name])
        os.replace(temp_path, path)
    finally:
        try:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
        except OSError:
            pass
    return read_xlsx_document(path)


def write_csv_document(path: str, sheets: List[Dict[str, object]], expected_etag: Optional[str]) -> Dict[str, object]:
    assert_expected_etag(path, expected_etag)
    normalized_sheet = normalize_sheet_payload(sheets[0] if sheets else {"name": "Sheet1", "cells": []})
    max_row = max((int(cell["row"]) for cell in normalized_sheet["cells"]), default=0)
    max_col = max((int(cell["col"]) for cell in normalized_sheet["cells"]), default=0)
    grid = [["" for _ in range(max_col)] for _ in range(max_row)]
    for cell in normalized_sheet["cells"]:
        row = int(cell["row"]) - 1
        col = int(cell["col"]) - 1
        if row < 0 or col < 0:
            continue
        if cell.get("formula"):
            grid[row][col] = f"={cell['formula']}"
        else:
            grid[row][col] = str(cell.get("value") or "")
    buffer = io.StringIO()
    writer = csv.writer(buffer)
    for row in grid:
        writer.writerow(row)
    atomic_write_text(path, buffer.getvalue())
    return read_csv_document(path)


TEXT_DOCUMENT_AUTOMATION_WARNING = (
    "Structured text automation uses Pandoc when available and falls back to local parsing for "
    "simple text/markup inputs. The canonical AIO document stays semantic even when the source "
    "format is markdown, HTML, org, rst, or plain text."
)

DOCLING_DOCUMENT_WARNING = (
    "Docling-backed document automation extracts semantic blocks, layout context, and provenance "
    "from PDFs and images. These formats are currently read-only in entropic-office."
)


def get_docling_converter() -> object:
    global DOCLING_CONVERTER
    if DocumentConverter is None:
        raise RuntimeError("Docling is not installed in this runtime.")
    if DOCLING_CONVERTER is None:
        DOCLING_CONVERTER = DocumentConverter()
    return DOCLING_CONVERTER


def markdown_escape_text(text: str) -> str:
    escaped = text.replace("\\", "\\\\")
    for token in ("`", "*", "_", "[", "]", "(", ")", "<", ">"):
        escaped = escaped.replace(token, f"\\{token}")
    return escaped


def markdown_inline_items(text: str, block_id: str) -> List[Dict[str, object]]:
    inlines: List[Dict[str, object]] = []
    inline_index = 1
    position = 0
    while position < len(text):
        link_match = re.match(r"\[([^\]]+)\]\(([^)]+)\)", text[position:])
        if link_match:
            label = link_match.group(1)
            href = link_match.group(2).strip()
            inlines.append(
                {
                    "id": f"{block_id}:inline:{inline_index}",
                    "kind": "link",
                    "text": label,
                    "href": href,
                }
            )
            inline_index += 1
            position += len(link_match.group(0))
            continue
        matched = False
        for delimiter, mark in (
            ("**", "strong"),
            ("__", "strong"),
            ("~~", "strike"),
            ("`", "code"),
            ("*", "emphasis"),
            ("_", "emphasis"),
        ):
            if not text.startswith(delimiter, position):
                continue
            end = text.find(delimiter, position + len(delimiter))
            if end <= position + len(delimiter):
                continue
            segment = text[position + len(delimiter) : end]
            inline_index = add_text_inline(
                inlines,
                block_id,
                inline_index,
                "text",
                segment,
                marks=[mark],
            )
            position = end + len(delimiter)
            matched = True
            break
        if matched:
            continue
        next_special = len(text)
        for marker in ("[", "**", "__", "~~", "`", "*", "_"):
            candidate = text.find(marker, position + 1)
            if candidate != -1:
                next_special = min(next_special, candidate)
        chunk = text[position:next_special]
        inline_index = add_text_inline(inlines, block_id, inline_index, "text", chunk)
        position = next_special
    return inlines


def markdown_paragraph_block(block_id: str, text: str, kind: str = "paragraph") -> Dict[str, object]:
    block: Dict[str, object] = {
        "id": block_id,
        "kind": kind,
        "text": text,
    }
    inlines = markdown_inline_items(text, block_id)
    if inlines:
        block["inlines"] = inlines
    return block


def html_to_text_fallback(source_text: str) -> str:
    text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", "", source_text)
    text = re.sub(r"(?is)<head.*?>.*?</head>", "", text)
    text = re.sub(r"(?i)<br\s*/?>", "\n", text)
    text = re.sub(r"(?i)<li[^>]*>", "- ", text)
    text = re.sub(r"(?i)</(p|div|section|article|h[1-6]|li|tr|table|ul|ol|blockquote)>", "\n\n", text)
    text = re.sub(r"(?s)<[^>]+>", "", text)
    return html.unescape(text)


def markdown_table_separator(line: str) -> bool:
    cells = [cell.strip() for cell in line.strip().strip("|").split("|")]
    if not cells:
        return False
    return all(bool(re.fullmatch(r":?-{3,}:?", cell or "")) for cell in cells)


def markdown_table_cells(line: str) -> List[str]:
    stripped = line.strip().strip("|")
    if stripped == "":
        return []
    return [cell.strip() for cell in stripped.split("|")]


def parse_markdown_table(lines: List[str], start_index: int, block_id: str) -> Tuple[Dict[str, object], int]:
    raw_rows: List[List[str]] = [markdown_table_cells(lines[start_index])]
    cursor = start_index + 2
    while cursor < len(lines):
        if "|" not in lines[cursor]:
            break
        row = markdown_table_cells(lines[cursor])
        if not row:
            break
        raw_rows.append(row)
        cursor += 1
    rows: List[Dict[str, object]] = []
    max_cols = max((len(row) for row in raw_rows), default=0)
    for row_index, row in enumerate(raw_rows, start=1):
        cells: List[Dict[str, object]] = []
        for cell_index in range(max_cols):
            cell_text = row[cell_index] if cell_index < len(row) else ""
            cell_block_id = f"{block_id}:row:{row_index}:cell:{cell_index + 1}:block:1"
            cells.append(
                {
                    "id": f"{block_id}:row:{row_index}:cell:{cell_index + 1}",
                    "kind": "table_cell",
                    "blocks": [markdown_paragraph_block(cell_block_id, cell_text)],
                }
            )
        row_payload: Dict[str, object] = {
            "id": f"{block_id}:row:{row_index}",
            "kind": "table_row",
            "cells": cells,
        }
        if row_index == 1:
            row_payload["role"] = "header"
        rows.append(row_payload)
    return {
        "id": block_id,
        "kind": "table",
        "rows": rows,
        "cols": max_cols,
    }, cursor


def parse_markdown_list(lines: List[str], start_index: int, block_id: str) -> Tuple[Dict[str, object], int]:
    marker_re = re.compile(r"^(\s*)([-+*]|\d+\.)\s+(.*)$")
    cursor = start_index
    items: List[Dict[str, object]] = []
    ordered = False
    base_indent: Optional[int] = None
    while cursor < len(lines):
        match = marker_re.match(lines[cursor])
        if not match:
            break
        indent = len(match.group(1).replace("\t", "    "))
        marker = match.group(2)
        text = match.group(3).strip()
        if base_indent is None:
            base_indent = indent
            ordered = marker.endswith(".") and marker[:-1].isdigit()
        if indent != base_indent:
            break
        blocks: List[Dict[str, object]] = []
        if text:
            blocks.append(markdown_paragraph_block(f"{block_id}:item:{len(items) + 1}:block:1", text))
        items.append(
            {
                "id": f"{block_id}:item:{len(items) + 1}",
                "kind": "list_item",
                "level": max(indent // 2, 0),
                "blocks": blocks or [markdown_paragraph_block(f"{block_id}:item:{len(items) + 1}:block:1", "")],
            }
        )
        cursor += 1
    return {
        "id": block_id,
        "kind": "list",
        "ordered": ordered,
        "items": items,
    }, cursor


def parse_markdown_fallback(text: str, block_prefix: str = "block") -> List[Dict[str, object]]:
    lines = text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    blocks: List[Dict[str, object]] = []
    paragraph_lines: List[str] = []
    block_index = 1

    def next_block_id() -> str:
        nonlocal block_index
        value = f"{block_prefix}:{block_index}"
        block_index += 1
        return value

    def flush_paragraph() -> None:
        if not paragraph_lines:
            return
        paragraph_text = " ".join(line.strip() for line in paragraph_lines if line.strip()).strip()
        if paragraph_text:
            blocks.append(markdown_paragraph_block(next_block_id(), paragraph_text))
        paragraph_lines.clear()

    cursor = 0
    while cursor < len(lines):
        line = lines[cursor]
        stripped = line.strip()
        if stripped == "":
            flush_paragraph()
            cursor += 1
            continue
        heading_match = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if heading_match:
            flush_paragraph()
            block_id = next_block_id()
            heading_text = heading_match.group(2).strip()
            heading = markdown_paragraph_block(block_id, heading_text, kind="heading")
            heading["level"] = len(heading_match.group(1))
            blocks.append(heading)
            cursor += 1
            continue
        if stripped.startswith(("```", "~~~")):
            flush_paragraph()
            fence = stripped[:3]
            language = stripped[3:].strip()
            cursor += 1
            code_lines: List[str] = []
            while cursor < len(lines) and not lines[cursor].strip().startswith(fence):
                code_lines.append(lines[cursor])
                cursor += 1
            if cursor < len(lines):
                cursor += 1
            code_block: Dict[str, object] = {
                "id": next_block_id(),
                "kind": "code_block",
                "text": "\n".join(code_lines),
            }
            if language:
                code_block["language_ref"] = language
            blocks.append(code_block)
            continue
        if stripped in {"---", "***", "___"}:
            flush_paragraph()
            blocks.append({"id": next_block_id(), "kind": "thematic_break"})
            cursor += 1
            continue
        if stripped.startswith(">"):
            flush_paragraph()
            quote_lines: List[str] = []
            while cursor < len(lines) and lines[cursor].strip().startswith(">"):
                quote_lines.append(re.sub(r"^\s*>\s?", "", lines[cursor]))
                cursor += 1
            blocks.append(
                {
                    "id": next_block_id(),
                    "kind": "quote",
                    "blocks": parse_markdown_fallback("\n".join(quote_lines), block_prefix=f"{block_prefix}:quote"),
                }
            )
            continue
        if "|" in line and cursor + 1 < len(lines) and markdown_table_separator(lines[cursor + 1]):
            flush_paragraph()
            table_block, cursor = parse_markdown_table(lines, cursor, next_block_id())
            blocks.append(table_block)
            continue
        if re.match(r"^\s*([-+*]|\d+\.)\s+", line):
            flush_paragraph()
            list_block, cursor = parse_markdown_list(lines, cursor, next_block_id())
            blocks.append(list_block)
            continue
        paragraph_lines.append(line)
        cursor += 1
    flush_paragraph()
    return blocks


def pandoc_attr_payload(attr: object) -> Dict[str, object]:
    if not isinstance(attr, list) or len(attr) != 3:
        return {}
    identifier = str(attr[0] or "").strip()
    classes = [str(value) for value in attr[1] if str(value or "").strip()]
    attributes = {
        str(item[0]): str(item[1])
        for item in attr[2]
        if isinstance(item, list) and len(item) == 2 and str(item[0] or "").strip()
    }
    payload: Dict[str, object] = {}
    if identifier:
        payload["identifier"] = identifier
    if classes:
        payload["classes"] = classes
    if attributes:
        payload["attributes"] = attributes
    return payload


def pandoc_plain_text(node: object) -> str:
    if node is None:
        return ""
    if isinstance(node, str):
        return node
    if isinstance(node, list):
        return "".join(pandoc_plain_text(item) for item in node)
    if not isinstance(node, dict):
        return ""
    tag = str(node.get("t") or "")
    content = node.get("c")
    if tag == "Str":
        return str(content or "")
    if tag == "Space":
        return " "
    if tag in {"SoftBreak", "LineBreak"}:
        return "\n"
    if tag in {"Emph", "Strong", "Underline", "Strikeout", "Superscript", "Subscript", "SmallCaps", "Cite"}:
        return pandoc_plain_text(content)
    if tag == "Code":
        if isinstance(content, list) and len(content) >= 2:
            return str(content[1] or "")
        return ""
    if tag == "Math":
        if isinstance(content, list) and len(content) >= 2:
            return str(content[1] or "")
        return ""
    if tag == "Span":
        if isinstance(content, list) and len(content) >= 2:
            return pandoc_plain_text(content[1])
        return ""
    if tag == "Link":
        if isinstance(content, list):
            if len(content) >= 3:
                return pandoc_plain_text(content[1])
            if len(content) >= 2:
                return pandoc_plain_text(content[0])
        return ""
    if tag == "Image":
        if isinstance(content, list):
            if len(content) >= 3:
                return pandoc_plain_text(content[1])
            if len(content) >= 2:
                return pandoc_plain_text(content[0])
        return ""
    if tag == "Note":
        return pandoc_plain_text(content)
    if isinstance(content, list):
        return pandoc_plain_text(content)
    return ""


def pandoc_link_parts(content: object) -> Tuple[object, str]:
    if not isinstance(content, list):
        return [], ""
    if len(content) >= 3:
        target = content[2]
        href = target[0] if isinstance(target, list) and target else ""
        return content[1], str(href or "")
    if len(content) >= 2:
        target = content[1]
        href = target[0] if isinstance(target, list) and target else ""
        return content[0], str(href or "")
    return [], ""


def pandoc_inlines_to_document_inlines(inline_nodes: object, block_id: str) -> List[Dict[str, object]]:
    inlines: List[Dict[str, object]] = []
    inline_index = 1

    def append_inline(
        kind: str,
        text: str = "",
        marks: Optional[List[str]] = None,
        href: Optional[str] = None,
        nested_inlines: Optional[List[Dict[str, object]]] = None,
        src: Optional[str] = None,
    ) -> None:
        nonlocal inline_index
        payload: Dict[str, object] = {
            "id": f"{block_id}:inline:{inline_index}",
            "kind": kind,
        }
        if kind != "line_break":
            payload["text"] = text
        if marks:
            payload["marks"] = sorted(set(mark for mark in marks if mark))
        if href:
            payload["href"] = href
        if nested_inlines:
            payload["label_inlines"] = nested_inlines
        if src:
            payload["src"] = src
        inlines.append(payload)
        inline_index += 1

    def visit(node: object, active_marks: Optional[List[str]] = None) -> None:
        marks = list(active_marks or [])
        if isinstance(node, list):
            for child in node:
                visit(child, marks)
            return
        if not isinstance(node, dict):
            text = pandoc_plain_text(node)
            if text:
                append_inline("text", text=text, marks=marks)
            return
        tag = str(node.get("t") or "")
        content = node.get("c")
        if tag == "Str":
            append_inline("text", text=str(content or ""), marks=marks)
            return
        if tag == "Space":
            append_inline("text", text=" ", marks=marks)
            return
        if tag == "SoftBreak":
            append_inline("text", text=" ", marks=marks)
            return
        if tag == "LineBreak":
            append_inline("line_break")
            return
        if tag in {
            "Emph",
            "Strong",
            "Underline",
            "Strikeout",
            "Superscript",
            "Subscript",
            "SmallCaps",
        }:
            mark_name = {
                "Emph": "emphasis",
                "Strong": "strong",
                "Underline": "underline",
                "Strikeout": "strike",
                "Superscript": "superscript",
                "Subscript": "subscript",
                "SmallCaps": "small_caps",
            }[tag]
            visit(content, marks + [mark_name])
            return
        if tag == "Span":
            if isinstance(content, list) and len(content) >= 2:
                visit(content[1], marks)
            return
        if tag == "Cite":
            if isinstance(content, list) and len(content) >= 2:
                visit(content[1], marks)
            return
        if tag == "Code":
            code_text = str(content[1] or "") if isinstance(content, list) and len(content) >= 2 else ""
            append_inline("text", text=code_text, marks=marks + ["code"])
            return
        if tag == "Math":
            math_text = str(content[1] or "") if isinstance(content, list) and len(content) >= 2 else ""
            append_inline("text", text=math_text, marks=marks + ["math"])
            return
        if tag == "Link":
            label_nodes, href = pandoc_link_parts(content)
            label_inlines = pandoc_inlines_to_document_inlines(label_nodes, f"{block_id}:inline:{inline_index}:label")
            label_text = inlines_text(label_inlines) or pandoc_plain_text(label_nodes)
            append_inline("link", text=label_text, href=href, nested_inlines=label_inlines or None)
            return
        if tag == "Image":
            label_nodes, href = pandoc_link_parts(content)
            alt_text = pandoc_plain_text(label_nodes)
            append_inline("image", text=alt_text, src=href)
            return
        if tag == "Note":
            note_text = pandoc_plain_text(content).strip()
            if note_text:
                append_inline("text", text=f"[{note_text}]", marks=marks)
            return
        text = pandoc_plain_text(node)
        if text:
            append_inline("text", text=text, marks=marks)

    visit(inline_nodes)
    return inlines


def pandoc_cell_blocks(cell: object, block_id: str) -> List[Dict[str, object]]:
    if isinstance(cell, dict) and str(cell.get("t") or "") == "Cell":
        content = cell.get("c")
        if isinstance(content, list) and content:
            return pandoc_blocks_to_document_blocks(content[-1], block_prefix=block_id)
    if isinstance(cell, list) and cell:
        last_item = cell[-1]
        if isinstance(last_item, list):
            return pandoc_blocks_to_document_blocks(last_item, block_prefix=block_id)
    return [markdown_paragraph_block(f"{block_id}:1", pandoc_plain_text(cell))]


def pandoc_section_rows(section: object) -> List[object]:
    if isinstance(section, dict):
        tag = str(section.get("t") or "")
        content = section.get("c")
        if tag in {"TableHead", "TableFoot"} and isinstance(content, list) and len(content) >= 2:
            return content[1] if isinstance(content[1], list) else []
        if tag == "TableBody" and isinstance(content, list) and len(content) >= 4:
            return content[3] if isinstance(content[3], list) else []
    if isinstance(section, list) and section:
        last_item = section[-1]
        if isinstance(last_item, list):
            return last_item
    return []


def pandoc_table_block(content: object, block_id: str) -> Optional[Dict[str, object]]:
    if not isinstance(content, list) or len(content) < 5:
        return None
    caption_text = pandoc_plain_text(content[1]) if len(content) >= 2 else ""
    row_sets: List[Tuple[str, List[object]]] = [("header", pandoc_section_rows(content[3]))]
    if len(content) >= 5 and isinstance(content[4], list):
        for body in content[4]:
            row_sets.append(("body", pandoc_section_rows(body)))
    if len(content) >= 6:
        row_sets.append(("footer", pandoc_section_rows(content[5])))
    rows: List[Dict[str, object]] = []
    row_index = 1
    max_cols = 0
    for row_role, raw_rows in row_sets:
        for raw_row in raw_rows:
            cells_raw = []
            if isinstance(raw_row, dict) and str(raw_row.get("t") or "") == "Row":
                row_content = raw_row.get("c")
                if isinstance(row_content, list) and len(row_content) >= 2 and isinstance(row_content[1], list):
                    cells_raw = row_content[1]
            elif isinstance(raw_row, list) and raw_row and isinstance(raw_row[-1], list):
                cells_raw = raw_row[-1]
            cells: List[Dict[str, object]] = []
            for cell_index, raw_cell in enumerate(cells_raw, start=1):
                cells.append(
                    {
                        "id": f"{block_id}:row:{row_index}:cell:{cell_index}",
                        "kind": "table_cell",
                        "blocks": pandoc_cell_blocks(
                            raw_cell,
                            f"{block_id}:row:{row_index}:cell:{cell_index}:block",
                        ),
                    }
                )
            row_payload: Dict[str, object] = {
                "id": f"{block_id}:row:{row_index}",
                "kind": "table_row",
                "cells": cells,
            }
            if row_role != "body":
                row_payload["role"] = row_role
            rows.append(row_payload)
            max_cols = max(max_cols, len(cells))
            row_index += 1
    if not rows:
        return None
    block: Dict[str, object] = {
        "id": block_id,
        "kind": "table",
        "rows": rows,
        "cols": max_cols,
    }
    if caption_text:
        block["caption"] = caption_text
    return block


def pandoc_blocks_to_document_blocks(block_nodes: object, block_prefix: str = "block") -> List[Dict[str, object]]:
    blocks: List[Dict[str, object]] = []
    counter = 1
    for node in block_nodes or []:
        if not isinstance(node, dict):
            text = pandoc_plain_text(node).strip()
            if text:
                blocks.append(markdown_paragraph_block(f"{block_prefix}:{counter}", text))
                counter += 1
            continue
        tag = str(node.get("t") or "")
        content = node.get("c")
        block_id = f"{block_prefix}:{counter}"
        counter += 1
        if tag in {"Para", "Plain"}:
            inlines = pandoc_inlines_to_document_inlines(content, block_id)
            block: Dict[str, object] = {
                "id": block_id,
                "kind": "paragraph",
                "text": inlines_text(inlines),
            }
            if inlines:
                block["inlines"] = inlines
            blocks.append(block)
            continue
        if tag == "Header" and isinstance(content, list) and len(content) >= 3:
            level = int(content[0] or 1)
            inlines = pandoc_inlines_to_document_inlines(content[2], block_id)
            block = {
                "id": block_id,
                "kind": "heading",
                "level": level,
                "text": inlines_text(inlines),
            }
            if inlines:
                block["inlines"] = inlines
            attrs = pandoc_attr_payload(content[1])
            if attrs:
                block["attrs"] = attrs
            blocks.append(block)
            continue
        if tag == "BulletList" and isinstance(content, list):
            items = []
            for item_index, item in enumerate(content, start=1):
                items.append(
                    {
                        "id": f"{block_id}:item:{item_index}",
                        "kind": "list_item",
                        "level": 0,
                        "blocks": pandoc_blocks_to_document_blocks(item, block_prefix=f"{block_id}:item:{item_index}:block"),
                    }
                )
            blocks.append(
                {
                    "id": block_id,
                    "kind": "list",
                    "ordered": False,
                    "items": items,
                }
            )
            continue
        if tag == "OrderedList" and isinstance(content, list) and len(content) >= 2:
            list_attrs = content[0] if isinstance(content[0], list) else []
            items_payload = []
            for item_index, item in enumerate(content[1], start=1):
                items_payload.append(
                    {
                        "id": f"{block_id}:item:{item_index}",
                        "kind": "list_item",
                        "level": 0,
                        "blocks": pandoc_blocks_to_document_blocks(item, block_prefix=f"{block_id}:item:{item_index}:block"),
                    }
                )
            ordered_list: Dict[str, object] = {
                "id": block_id,
                "kind": "list",
                "ordered": True,
                "items": items_payload,
            }
            if isinstance(list_attrs, list) and list_attrs:
                ordered_list["start"] = int(list_attrs[0] or 1)
                if len(list_attrs) >= 2:
                    ordered_list["list_style"] = str(list_attrs[1])
            blocks.append(ordered_list)
            continue
        if tag == "BlockQuote":
            blocks.append(
                {
                    "id": block_id,
                    "kind": "quote",
                    "blocks": pandoc_blocks_to_document_blocks(content, block_prefix=f"{block_id}:block"),
                }
            )
            continue
        if tag == "CodeBlock":
            attrs = content[0] if isinstance(content, list) and len(content) >= 1 else []
            code_text = str(content[1] or "") if isinstance(content, list) and len(content) >= 2 else ""
            code_block: Dict[str, object] = {
                "id": block_id,
                "kind": "code_block",
                "text": code_text,
            }
            attr_payload = pandoc_attr_payload(attrs)
            if attr_payload.get("classes"):
                code_block["language_ref"] = str(attr_payload["classes"][0])
            if attr_payload:
                code_block["attrs"] = attr_payload
            blocks.append(code_block)
            continue
        if tag == "HorizontalRule":
            blocks.append({"id": block_id, "kind": "thematic_break"})
            continue
        if tag == "LineBlock" and isinstance(content, list):
            line_inlines: List[Dict[str, object]] = []
            line_index = 1
            for content_index, line in enumerate(content):
                for inline in pandoc_inlines_to_document_inlines(line, f"{block_id}:line:{content_index + 1}"):
                    inline_payload = dict(inline)
                    inline_payload["id"] = f"{block_id}:inline:{line_index}"
                    line_inlines.append(inline_payload)
                    line_index += 1
                if content_index < len(content) - 1:
                    line_inlines.append({"id": f"{block_id}:inline:{line_index}", "kind": "line_break"})
                    line_index += 1
            paragraph: Dict[str, object] = {
                "id": block_id,
                "kind": "paragraph",
                "text": inlines_text(line_inlines),
                "inlines": line_inlines,
            }
            blocks.append(paragraph)
            continue
        if tag == "Div" and isinstance(content, list) and len(content) >= 2:
            section: Dict[str, object] = {
                "id": block_id,
                "kind": "section",
                "blocks": pandoc_blocks_to_document_blocks(content[1], block_prefix=f"{block_id}:block"),
            }
            attrs = pandoc_attr_payload(content[0])
            if attrs:
                section["attrs"] = attrs
            blocks.append(section)
            continue
        if tag == "Table":
            table_block = pandoc_table_block(content, block_id)
            if table_block is not None:
                blocks.append(table_block)
                continue
        if tag == "RawBlock" and isinstance(content, list) and len(content) >= 2:
            blocks.append(
                {
                    "id": block_id,
                    "kind": "raw_block",
                    "format": str(content[0] or ""),
                    "text": str(content[1] or ""),
                }
            )
            continue
        if tag == "Null":
            continue
        fallback_text = pandoc_plain_text(node).strip()
        if fallback_text:
            blocks.append(markdown_paragraph_block(block_id, fallback_text))
    return blocks


def pandoc_projection_payload(input_format: str, ast: Dict[str, object]) -> Dict[str, object]:
    projection: Dict[str, object] = {
        "backend:pandoc": {
            "kind": "backend-ast",
            "backend": "pandoc",
            "available": True,
            "input_format": input_format,
        }
    }
    embedded, summary = embedded_backend_payload(ast)
    if embedded is not None:
        projection["backend:pandoc"]["ast"] = embedded
    projection["backend:pandoc"]["summary"] = summary
    return projection


def pandoc_ast_from_path(path: str, input_format: str) -> Dict[str, object]:
    payload = run_subprocess_json(["pandoc", "--from", input_format, "--to", "json", path])
    if not isinstance(payload, dict):
        raise RuntimeError("Pandoc JSON output was not an object.")
    return payload


def pandoc_ast_from_text(source_text: str, input_format: str) -> Dict[str, object]:
    payload = run_subprocess_json(["pandoc", "--from", input_format, "--to", "json"], stdin_text=source_text)
    if not isinstance(payload, dict):
        raise RuntimeError("Pandoc JSON output was not an object.")
    return payload


def parse_text_source_blocks(
    source_text: str,
    input_format: str,
    path: Optional[str] = None,
) -> Tuple[List[Dict[str, object]], Optional[Dict[str, object]], str]:
    if pandoc_available():
        ast = pandoc_ast_from_path(path, input_format) if path else pandoc_ast_from_text(source_text, input_format)
        return pandoc_blocks_to_document_blocks(ast.get("blocks") or []), ast, "pandoc"
    if input_format == "html":
        return parse_markdown_fallback(html_to_text_fallback(source_text)), None, "native-text"
    return parse_markdown_fallback(source_text), None, "native-text"


def text_export_artifact(value: str) -> Dict[str, object]:
    embedded, summary = embedded_backend_payload(value)
    artifact: Dict[str, object] = dict(summary)
    if embedded is not None:
        artifact["value"] = embedded
    return artifact


def docling_collection_count(payload: object) -> int:
    if isinstance(payload, dict):
        return len(payload)
    if isinstance(payload, list):
        return len(payload)
    return 0


def docling_page_summaries(document_dict: Dict[str, object]) -> List[Dict[str, object]]:
    pages_raw = document_dict.get("pages")
    page_items = list(pages_raw.values()) if isinstance(pages_raw, dict) else list(pages_raw or [])
    summaries: List[Dict[str, object]] = []
    for page_index, page in enumerate(page_items, start=1):
        if not isinstance(page, dict):
            continue
        summary: Dict[str, object] = {
            "page_no": int(page.get("page_no") or page.get("page") or page_index),
        }
        for key in ("width", "height", "dpi", "rotation"):
            if page.get(key) is not None:
                summary[key] = page.get(key)
        size = page.get("size")
        if isinstance(size, dict):
            for key in ("width", "height"):
                if size.get(key) is not None:
                    summary[key] = size.get(key)
        assembled = page.get("assembled")
        if isinstance(assembled, dict) and isinstance(assembled.get("elements"), list):
            summary["element_count"] = len(assembled["elements"])
        summaries.append(summary)
    return summaries


def docling_summary_from_dict(document_dict: Dict[str, object]) -> Dict[str, object]:
    return {
        "page_count": docling_collection_count(document_dict.get("pages")),
        "text_item_count": docling_collection_count(document_dict.get("texts")),
        "table_count": docling_collection_count(document_dict.get("tables")),
        "picture_count": docling_collection_count(document_dict.get("pictures")),
        "group_count": docling_collection_count(document_dict.get("groups")),
        "page_summaries": docling_page_summaries(document_dict),
    }


def build_document_record(
    path: str,
    format_name: str,
    metadata: Dict[str, object],
    warning: str,
    backend_name: str,
    blocks: List[Dict[str, object]],
    analysis: Optional[Dict[str, object]] = None,
    annotations: Optional[List[Dict[str, object]]] = None,
    projections: Optional[Dict[str, object]] = None,
    artifacts: Optional[Dict[str, object]] = None,
) -> Dict[str, object]:
    final_analysis = merge_nested_dict(document_analysis_payload(blocks, backend_name), analysis)
    final_projections = merge_nested_dict(document_projection_payload(blocks), projections)
    final_artifacts = merge_nested_dict(document_artifact_payload(format_name, backend_name, 1), artifacts)
    final_annotations = [item for item in annotations or [] if isinstance(item, dict)] or collect_document_annotations(blocks)
    return {
        "kind": "document",
        "format": format_name,
        "path": path,
        **metadata,
        "backend_name": backend_name,
        "warning": warning,
        "blocks": blocks,
        "paragraphs": document_blocks_to_paragraphs(blocks),
        "analysis": final_analysis,
        "annotations": final_annotations,
        "projections": final_projections,
        "artifacts": final_artifacts,
    }


def read_pandoc_document(path: str) -> Dict[str, object]:
    metadata = path_metadata(path)
    format_name = text_document_format_name(path)
    backend_name = default_backend_name("document", format_name)
    if not metadata["exists"]:
        return build_document_record(
            path,
            format_name,
            metadata,
            TEXT_DOCUMENT_AUTOMATION_WARNING,
            backend_name,
            [],
            analysis={"source_format": pandoc_input_format_for_extension(split_extension(path)) or format_name},
        )
    source_text = read_utf8_text(path)
    input_format = pandoc_input_format_for_extension(split_extension(path)) or "markdown"
    blocks, pandoc_ast, parse_backend = parse_text_source_blocks(source_text, input_format, path=path)
    analysis: Dict[str, object] = {
        "source_backend": parse_backend,
        "source_format": input_format,
        "text_stats": {
            "character_count": len(source_text),
            "line_count": len(source_text.splitlines()),
        },
    }
    if pandoc_ast is not None and isinstance(pandoc_ast.get("meta"), dict) and pandoc_ast.get("meta"):
        analysis["metadata"] = pandoc_ast.get("meta")
    projections = pandoc_projection_payload(input_format, pandoc_ast) if pandoc_ast is not None else {}
    artifacts = {
        "source_text": text_export_artifact(source_text),
    }
    return build_document_record(
        path,
        format_name,
        metadata,
        TEXT_DOCUMENT_AUTOMATION_WARNING,
        backend_name,
        blocks,
        analysis=analysis,
        projections=projections,
        artifacts=artifacts,
    )


def read_docling_document(path: str) -> Dict[str, object]:
    metadata = path_metadata(path)
    if not metadata["exists"]:
        raise RuntimeError("PDF and image inspection requires an existing source file.")
    if not docling_available():
        raise RuntimeError("Docling is not installed in this runtime.")
    extension = split_extension(path)
    format_name = text_document_format_name(path)
    converter = get_docling_converter()
    try:
        result = converter.convert(path)
    except Exception as error:
        raise RuntimeError(f"Docling failed to inspect `{posixpath.basename(path)}`: {error}") from error
    document = result.document
    markdown_text = document.export_to_markdown()
    html_text = document.export_to_html()
    plain_text = document.export_to_text()
    document_dict = document.export_to_dict()
    blocks, pandoc_ast, parse_backend = parse_text_source_blocks(markdown_text, "markdown")
    analysis: Dict[str, object] = {
        "source_backend": "docling",
        "source_format": docling_input_format_for_extension(extension) or format_name,
        "layout": docling_summary_from_dict(document_dict),
        "conversion": {
            "status": str(getattr(result, "status", "success")),
            "parse_backend": parse_backend,
        },
    }
    errors = getattr(result, "errors", None)
    if errors:
        analysis["conversion"]["errors"] = [str(item) for item in errors][:20]
    projections = {
        "backend:docling": {
            "kind": "analysis-adapter",
            "backend": "docling",
            "available": True,
        }
    }
    embedded_docling, docling_summary = embedded_backend_payload(document_dict)
    if embedded_docling is not None:
        projections["backend:docling"]["document"] = embedded_docling
    projections["backend:docling"]["summary"] = merge_nested_dict(docling_summary, docling_summary_from_dict(document_dict))
    if pandoc_ast is not None:
        projections = merge_nested_dict(projections, pandoc_projection_payload("markdown", pandoc_ast))
    artifacts = {
        "backend:docling": {
            "format": format_name,
            "input_format": docling_input_format_for_extension(extension) or format_name,
            "exports": {
                "markdown": text_export_artifact(markdown_text),
                "html": text_export_artifact(html_text),
                "text": text_export_artifact(plain_text),
            },
        }
    }
    return build_document_record(
        path,
        format_name,
        metadata,
        DOCLING_DOCUMENT_WARNING,
        "docling",
        blocks,
        analysis=analysis,
        projections=projections,
        artifacts=artifacts,
    )


def inline_markdown_text(inline: Dict[str, object]) -> str:
    kind = str(inline.get("kind") or "text")
    if kind == "line_break":
        return "  \n"
    if kind == "image":
        alt_text = markdown_escape_text(str(inline.get("text") or ""))
        src = str(inline.get("src") or inline.get("href") or "").strip()
        return f"![{alt_text}]({src})" if src else alt_text
    nested = inline.get("label_inlines")
    if isinstance(nested, list) and nested:
        text = "".join(inline_markdown_text(item) for item in nested if isinstance(item, dict))
    else:
        text = markdown_escape_text(str(inline.get("text") or ""))
    for mark in ("code", "strong", "emphasis", "underline", "strike", "superscript", "subscript", "math"):
        if mark not in inline.get("marks", []):
            continue
        if mark == "code":
            text = f"`{text}`"
        elif mark == "strong":
            text = f"**{text}**"
        elif mark == "emphasis":
            text = f"*{text}*"
        elif mark == "underline":
            text = f"<u>{text}</u>"
        elif mark == "strike":
            text = f"~~{text}~~"
        elif mark == "superscript":
            text = f"^{text}^"
        elif mark == "subscript":
            text = f"~{text}~"
        elif mark == "math":
            text = f"${text}$"
    if kind == "link":
        href = str(inline.get("href") or "").strip()
        if href:
            return f"[{text}]({href})"
    return text


def block_markdown_text(block: Dict[str, object], depth: int = 0) -> str:
    kind = str(block.get("kind") or "")
    indent = "  " * depth
    inlines = block.get("inlines")
    text = (
        "".join(inline_markdown_text(item) for item in inlines if isinstance(item, dict))
        if isinstance(inlines, list) and inlines
        else markdown_escape_text(str(block.get("text") or ""))
    )
    if kind == "heading":
        return f"{'#' * max(1, int(block.get('level') or 1))} {text}".strip()
    if kind == "code_block":
        language = str(block.get("language_ref") or "").strip()
        fence = f"```{language}".rstrip()
        return f"{fence}\n{str(block.get('text') or '')}\n```"
    if kind == "quote":
        nested = document_blocks_to_markdown(block.get("blocks") or [], depth=depth).strip()
        if not nested:
            return "> "
        return "\n".join(f"> {line}" if line else ">" for line in nested.splitlines())
    if kind == "list":
        lines: List[str] = []
        items = [item for item in block.get("items") or [] if isinstance(item, dict)]
        start = int(block.get("start") or 1)
        for item_index, item in enumerate(items, start=0):
            marker = f"{start + item_index}. " if block.get("ordered") else "- "
            item_body = document_blocks_to_markdown(item.get("blocks") or [], depth=depth + 1).strip()
            item_lines = item_body.splitlines() if item_body else [""]
            lines.append(f"{indent}{marker}{item_lines[0]}")
            continuation_prefix = f"{indent}{' ' * len(marker)}"
            for extra_line in item_lines[1:]:
                lines.append(f"{continuation_prefix}{extra_line}")
        return "\n".join(lines)
    if kind == "table":
        rows = [row for row in block.get("rows") or [] if isinstance(row, dict)]
        if not rows:
            return ""
        rendered_rows: List[str] = []
        for row_index, row in enumerate(rows, start=1):
            cells = []
            for cell in row.get("cells") or []:
                if not isinstance(cell, dict):
                    continue
                cell_text = " ".join(
                    document_block_text(nested).strip()
                    for nested in cell.get("blocks") or []
                    if isinstance(nested, dict) and document_block_text(nested).strip()
                )
                cells.append(cell_text)
            rendered_rows.append(f"| {' | '.join(cells)} |")
            if row_index == 1:
                rendered_rows.append(f"| {' | '.join('---' for _ in cells)} |")
        return "\n".join(rendered_rows)
    if kind == "thematic_break":
        return "---"
    if kind == "section":
        return document_blocks_to_markdown(block.get("blocks") or [], depth=depth)
    if kind == "raw_block":
        return str(block.get("text") or "")
    if kind == "image":
        alt_text = markdown_escape_text(str(block.get("text") or ""))
        src = str(block.get("src") or "").strip()
        return f"![{alt_text}]({src})" if src else alt_text
    return text


def document_blocks_to_markdown(blocks: List[Dict[str, object]], depth: int = 0) -> str:
    rendered = [
        block_markdown_text(block, depth=depth)
        for block in blocks
        if isinstance(block, dict) and block_markdown_text(block, depth=depth).strip()
    ]
    return "\n\n".join(rendered).strip()


def document_blocks_to_plain_text(blocks: List[Dict[str, object]]) -> str:
    return "\n\n".join(document_blocks_to_paragraphs(blocks)).strip()


def minimal_html_document(path: str, body_html: str) -> str:
    title = html.escape(basename_without_extension(path))
    return (
        "<!DOCTYPE html>\n"
        "<html>\n"
        "<head>\n"
        "  <meta charset=\"utf-8\" />\n"
        f"  <title>{title}</title>\n"
        "</head>\n"
        "<body>\n"
        f"{body_html}\n"
        "</body>\n"
        "</html>\n"
    )


def write_pandoc_document(path: str, blocks: List[Dict[str, object]], expected_etag: Optional[str]) -> Dict[str, object]:
    assert_expected_etag(path, expected_etag)
    extension = split_extension(path)
    output_format = pandoc_output_format_for_extension(extension)
    if output_format is None:
        raise RuntimeError("Unsupported text document format.")
    markdown_text = document_blocks_to_markdown(blocks).strip()
    if markdown_text:
        markdown_text += "\n"
    if extension in {".md", ".markdown", ".qmd"}:
        atomic_write_text(path, markdown_text)
        return attach_aio_projection(read_pandoc_document(path))
    if extension == ".txt":
        plain_text = document_blocks_to_plain_text(blocks)
        if plain_text:
            plain_text += "\n"
        atomic_write_text(path, plain_text)
        return attach_aio_projection(read_pandoc_document(path))
    if not pandoc_available():
        if extension in {".html", ".htm"}:
            plain_text = document_blocks_to_plain_text(blocks)
            body_html = f"<pre>{html.escape(plain_text)}</pre>"
            atomic_write_text(path, minimal_html_document(path, body_html))
            return attach_aio_projection(read_pandoc_document(path))
        raise RuntimeError("Pandoc is required to save this text document format.")
    rendered = run_subprocess_text(["pandoc", "--from", "markdown", "--to", output_format], stdin_text=markdown_text)
    if extension in {".html", ".htm"}:
        rendered = minimal_html_document(path, rendered.strip())
    atomic_write_text(path, rendered if rendered.endswith("\n") else f"{rendered}\n")
    return attach_aio_projection(read_pandoc_document(path))


DOCUMENT_AUTOMATION_WARNING = (
    "Structured document automation reads semantic blocks, basic list and table structure, "
    "common inline marks, hyperlinks, and outline projections. Native DOCX write now preserves "
    "headings, paragraphs, lists, tables, common marks, and hyperlinks, while exact numbering, "
    "advanced layout, and comments remain only partially modeled."
)


def docx_attr(name: str) -> str:
    return f"{{{DOCX_NS}}}{name}"


def docx_attr_value(node: Optional[ET.Element], name: str, default: str = "") -> str:
    if node is None:
        return default
    return str(node.attrib.get(docx_attr(name), node.attrib.get(name, default)) or default)


def docx_toggle_enabled(node: Optional[ET.Element]) -> bool:
    if node is None:
        return False
    value = docx_attr_value(node, "val", "true").strip().lower()
    return value not in {"0", "false", "off", "no"}


def docx_relationship_targets(archive: zipfile.ZipFile, path: str) -> Dict[str, str]:
    if path not in archive.namelist():
        return {}
    root = ET.fromstring(archive.read(path))
    targets: Dict[str, str] = {}
    for rel in root:
        if local_name(rel.tag) != "Relationship":
            continue
        rel_id = rel.attrib.get("Id")
        target = rel.attrib.get("Target")
        if not rel_id or not target:
            continue
        if rel.attrib.get("TargetMode") == "External":
            targets[rel_id] = target
        else:
            targets[rel_id] = normalize_relationship_target("word", target)
    return targets


def read_docx_style_names(archive: zipfile.ZipFile) -> Dict[str, str]:
    if "word/styles.xml" not in archive.namelist():
        return {}
    root = ET.fromstring(archive.read("word/styles.xml"))
    names: Dict[str, str] = {}
    for style in root.findall(docx_tag("style")):
        style_id = docx_attr_value(style, "styleId")
        name_node = style.find(docx_tag("name"))
        name = docx_attr_value(name_node, "val")
        if style_id:
            names[style_id] = name or style_id
    return names


def read_docx_numbering_formats(archive: zipfile.ZipFile) -> Dict[str, Dict[str, str]]:
    if "word/numbering.xml" not in archive.namelist():
        return {}
    root = ET.fromstring(archive.read("word/numbering.xml"))
    abstract_formats: Dict[str, Dict[str, str]] = {}
    for abstract in root.findall(docx_tag("abstractNum")):
        abstract_id = docx_attr_value(abstract, "abstractNumId")
        if not abstract_id:
            continue
        level_formats: Dict[str, str] = {}
        for level in abstract.findall(docx_tag("lvl")):
            level_id = docx_attr_value(level, "ilvl", "0")
            num_fmt = docx_attr_value(level.find(docx_tag("numFmt")), "val", "bullet")
            level_formats[level_id] = num_fmt or "bullet"
        abstract_formats[abstract_id] = level_formats
    numbering: Dict[str, Dict[str, str]] = {}
    for num in root.findall(docx_tag("num")):
        num_id = docx_attr_value(num, "numId")
        abstract_ref = docx_attr_value(num.find(docx_tag("abstractNumId")), "val")
        if num_id and abstract_ref:
            numbering[num_id] = abstract_formats.get(abstract_ref, {})
    return numbering


def docx_heading_level(style_id: str, style_name: str) -> Optional[int]:
    for candidate in (style_name, style_id):
        text = str(candidate or "").strip().lower()
        if not text:
            continue
        match = re.search(r"heading[\s_-]*([1-9])", text)
        if match:
            return int(match.group(1))
    return None


def docx_run_marks(run: ET.Element) -> List[str]:
    marks: List[str] = []
    properties = run.find(docx_tag("rPr"))
    if properties is None:
        return marks
    for tag_name, mark in (("b", "strong"), ("i", "emphasis"), ("u", "underline"), ("strike", "strike")):
        if docx_toggle_enabled(properties.find(docx_tag(tag_name))):
            marks.append(mark)
    vertical_align = properties.find(docx_tag("vertAlign"))
    align_value = docx_attr_value(vertical_align, "val").strip().lower()
    if align_value == "superscript":
        marks.append("superscript")
    elif align_value == "subscript":
        marks.append("subscript")
    return marks


def docx_run_text(run: ET.Element) -> str:
    parts: List[str] = []
    for child in run:
        name = local_name(child.tag)
        if name == "t":
            parts.append(child.text or "")
        elif name == "tab":
            parts.append("\t")
        elif name in {"br", "cr"}:
            parts.append("\n")
    return "".join(parts)


def inlines_text(inlines: List[Dict[str, object]]) -> str:
    parts: List[str] = []
    for inline in inlines:
        if str(inline.get("kind") or "") == "line_break":
            parts.append("\n")
        else:
            parts.append(str(inline.get("text") or ""))
    return "".join(parts)


def add_text_inline(
    inlines: List[Dict[str, object]],
    block_id: str,
    inline_index: int,
    kind: str,
    text: str,
    marks: Optional[List[str]] = None,
    href: Optional[str] = None,
) -> int:
    if kind != "line_break" and text == "":
        return inline_index
    payload: Dict[str, object] = {
        "id": f"{block_id}:inline:{inline_index}",
        "kind": kind,
    }
    if kind != "line_break":
        payload["text"] = text
    if marks:
        payload["marks"] = sorted(set(marks))
    if href:
        payload["href"] = href
    inlines.append(payload)
    return inline_index + 1


def paragraph_inlines(paragraph: ET.Element, relationship_targets: Dict[str, str], block_id: str) -> List[Dict[str, object]]:
    inlines: List[Dict[str, object]] = []
    inline_index = 1
    for child in paragraph:
        child_name = local_name(child.tag)
        if child_name == "r":
            marks = docx_run_marks(child)
            text = docx_run_text(child)
            if "\n" not in text:
                inline_index = add_text_inline(inlines, block_id, inline_index, "text", text, marks=marks)
                continue
            segments = text.split("\n")
            for segment_index, segment in enumerate(segments):
                if segment:
                    inline_index = add_text_inline(inlines, block_id, inline_index, "text", segment, marks=marks)
                if segment_index < len(segments) - 1:
                    inline_index = add_text_inline(inlines, block_id, inline_index, "line_break", "")
        elif child_name == "hyperlink":
            rel_id = child.attrib.get(f"{{{DOC_REL_NS}}}id")
            href = relationship_targets.get(rel_id or "", "")
            link_text_parts: List[str] = []
            link_marks: List[str] = []
            for run in child.findall(docx_tag("r")):
                link_text_parts.append(docx_run_text(run))
                link_marks.extend(docx_run_marks(run))
            link_text = "".join(link_text_parts)
            if link_text:
                inline_index = add_text_inline(
                    inlines,
                    block_id,
                    inline_index,
                    "link" if href else "text",
                    link_text,
                    marks=link_marks,
                    href=href or None,
                )
    return inlines


def paragraph_properties(
    paragraph: ET.Element,
    style_names: Dict[str, str],
    numbering_formats: Dict[str, Dict[str, str]],
) -> Dict[str, object]:
    properties: Dict[str, object] = {}
    paragraph_properties = paragraph.find(docx_tag("pPr"))
    if paragraph_properties is None:
        return properties
    style_id = docx_attr_value(paragraph_properties.find(docx_tag("pStyle")), "val")
    if style_id:
        properties["style_id"] = style_id
        properties["style_name"] = style_names.get(style_id, style_id)
        heading_level = docx_heading_level(style_id, str(properties.get("style_name") or ""))
        if heading_level is not None:
            properties["heading_level"] = heading_level
    numbering = paragraph_properties.find(docx_tag("numPr"))
    if numbering is not None:
        num_id = docx_attr_value(numbering.find(docx_tag("numId")), "val")
        level = int(docx_attr_value(numbering.find(docx_tag("ilvl")), "val", "0") or 0)
        if num_id:
            number_style = numbering_formats.get(num_id, {}).get(str(level), "bullet")
            properties["num_id"] = num_id
            properties["list_level"] = max(level, 0)
            properties["list_style"] = number_style
            properties["ordered"] = number_style not in {"bullet", "none"}
    return properties


def parse_docx_paragraph_block(
    paragraph: ET.Element,
    style_names: Dict[str, str],
    numbering_formats: Dict[str, Dict[str, str]],
    relationship_targets: Dict[str, str],
    block_id: str,
) -> Dict[str, object]:
    props = paragraph_properties(paragraph, style_names, numbering_formats)
    inlines = paragraph_inlines(paragraph, relationship_targets, block_id)
    text = inlines_text(inlines)
    if "num_id" in props:
        item_block: Dict[str, object] = {
            "id": block_id,
            "kind": "list_item",
            "level": int(props.get("list_level") or 0),
            "ordered": bool(props.get("ordered")),
            "list_style": str(props.get("list_style") or ""),
            "num_id": str(props.get("num_id") or ""),
            "blocks": [
                {
                    "id": f"{block_id}:paragraph",
                    "kind": "paragraph",
                    "text": text,
                    "inlines": inlines,
                }
            ],
        }
        return item_block
    block: Dict[str, object] = {
        "id": block_id,
        "kind": "heading" if "heading_level" in props else "paragraph",
        "text": text,
    }
    if inlines:
        block["inlines"] = inlines
    if "heading_level" in props:
        block["level"] = int(props["heading_level"])
    if props.get("style_name"):
        block["style_ref"] = str(props["style_name"])
    return block


def group_document_blocks(blocks: List[Dict[str, object]]) -> List[Dict[str, object]]:
    grouped: List[Dict[str, object]] = []
    current_list: Optional[Dict[str, object]] = None
    current_key: Optional[Tuple[bool, str, str]] = None
    for block in blocks:
        if not isinstance(block, dict):
            continue
        if str(block.get("kind") or "") != "list_item":
            if current_list is not None:
                grouped.append(current_list)
                current_list = None
                current_key = None
            grouped.append(block)
            continue
        key = (
            bool(block.get("ordered")),
            str(block.get("num_id") or ""),
            str(block.get("list_style") or ""),
        )
        item_payload: Dict[str, object] = {
            "id": f"{block.get('id')}:item",
            "kind": "list_item",
            "level": int(block.get("level") or 0),
            "blocks": block.get("blocks") or [],
        }
        if current_list is None or key != current_key:
            if current_list is not None:
                grouped.append(current_list)
            current_list = {
                "id": f"{block.get('id')}:list",
                "kind": "list",
                "ordered": bool(block.get("ordered")),
                "items": [item_payload],
            }
            if block.get("list_style"):
                current_list["list_style"] = str(block.get("list_style"))
            current_key = key
        else:
            current_list.setdefault("items", []).append(item_payload)
    if current_list is not None:
        grouped.append(current_list)
    return grouped


def parse_docx_table_block(
    table: ET.Element,
    style_names: Dict[str, str],
    numbering_formats: Dict[str, Dict[str, str]],
    relationship_targets: Dict[str, str],
    block_id: str,
) -> Dict[str, object]:
    rows: List[Dict[str, object]] = []
    max_cols = 0
    for row_index, row in enumerate(table.findall(docx_tag("tr")), start=1):
        cells: List[Dict[str, object]] = []
        for cell_index, cell in enumerate(row.findall(docx_tag("tc")), start=1):
            raw_blocks: List[Dict[str, object]] = []
            nested_counter = 1
            for child in cell:
                child_name = local_name(child.tag)
                nested_id = f"{block_id}:row:{row_index}:cell:{cell_index}:block:{nested_counter}"
                if child_name == "p":
                    raw_blocks.append(
                        parse_docx_paragraph_block(
                            child,
                            style_names,
                            numbering_formats,
                            relationship_targets,
                            nested_id,
                        )
                    )
                    nested_counter += 1
                elif child_name == "tbl":
                    raw_blocks.append(
                        parse_docx_table_block(
                            child,
                            style_names,
                            numbering_formats,
                            relationship_targets,
                            nested_id,
                        )
                    )
                    nested_counter += 1
            cell_payload: Dict[str, object] = {
                "id": f"{block_id}:row:{row_index}:cell:{cell_index}",
                "kind": "table_cell",
                "blocks": group_document_blocks(raw_blocks)
                or [{"id": f"{block_id}:row:{row_index}:cell:{cell_index}:block:1", "kind": "paragraph", "text": ""}],
            }
            cells.append(cell_payload)
        max_cols = max(max_cols, len(cells))
        rows.append(
            {
                "id": f"{block_id}:row:{row_index}",
                "kind": "table_row",
                "cells": cells,
            }
        )
    return {
        "id": block_id,
        "kind": "table",
        "rows": rows,
        "cols": max_cols,
    }


def walk_document_blocks(blocks: List[Dict[str, object]]) -> Iterable[Dict[str, object]]:
    for block in blocks:
        if not isinstance(block, dict):
            continue
        yield block
        if str(block.get("kind") or "") in {"quote", "section"}:
            yield from walk_document_blocks(block.get("blocks") or [])
        if str(block.get("kind") or "") == "list":
            for item in block.get("items") or []:
                if isinstance(item, dict):
                    yield from walk_document_blocks(item.get("blocks") or [])
        if str(block.get("kind") or "") == "table":
            for row in block.get("rows") or []:
                if not isinstance(row, dict):
                    continue
                for cell in row.get("cells") or []:
                    if isinstance(cell, dict):
                        yield from walk_document_blocks(cell.get("blocks") or [])


def document_block_text(block: Dict[str, object]) -> str:
    if not isinstance(block, dict):
        return ""
    kind = str(block.get("kind") or "")
    text = str(block.get("text") or "")
    if text:
        return text
    inlines = block.get("inlines")
    if isinstance(inlines, list) and inlines:
        return inlines_text(inlines)
    if kind in {"quote", "section"}:
        return "\n".join(
            document_block_text(nested)
            for nested in block.get("blocks") or []
            if isinstance(nested, dict) and document_block_text(nested)
        )
    if kind == "list":
        parts = []
        for item in block.get("items") or []:
            if isinstance(item, dict):
                for nested in item.get("blocks") or []:
                    if isinstance(nested, dict):
                        nested_text = document_block_text(nested)
                        if nested_text:
                            parts.append(nested_text)
        return "\n".join(parts)
    if kind == "table":
        row_texts = []
        for row in block.get("rows") or []:
            if not isinstance(row, dict):
                continue
            cell_texts = []
            for cell in row.get("cells") or []:
                if not isinstance(cell, dict):
                    continue
                cell_parts = [
                    document_block_text(nested)
                    for nested in cell.get("blocks") or []
                    if isinstance(nested, dict) and document_block_text(nested)
                ]
                cell_texts.append(" ".join(cell_parts).strip())
            if cell_texts:
                row_texts.append(" | ".join(cell_texts).strip())
        return "\n".join(value for value in row_texts if value)
    if kind in {"raw_block", "image"}:
        return str(block.get("text") or "")
    return ""


def document_blocks_to_paragraphs(blocks: List[Dict[str, object]]) -> List[str]:
    paragraphs: List[str] = []
    for block in blocks:
        if not isinstance(block, dict):
            continue
        kind = str(block.get("kind") or "")
        if kind in {"quote", "section"}:
            paragraphs.extend(
                document_blocks_to_paragraphs(
                    [nested for nested in block.get("blocks") or [] if isinstance(nested, dict)]
                )
            )
            continue
        if kind == "list":
            for item in block.get("items") or []:
                if not isinstance(item, dict):
                    continue
                for nested in item.get("blocks") or []:
                    if not isinstance(nested, dict):
                        continue
                    nested_text = document_block_text(nested)
                    if nested_text:
                        paragraphs.append(nested_text)
            continue
        if kind == "table":
            text = document_block_text(block)
            if text:
                paragraphs.extend([line for line in text.splitlines() if line.strip()])
            continue
        if kind == "thematic_break":
            continue
        text = document_block_text(block)
        if text:
            paragraphs.append(text)
    return paragraphs


def collect_document_annotations(blocks: List[Dict[str, object]]) -> List[Dict[str, object]]:
    annotations: List[Dict[str, object]] = []
    annotation_index = 1
    for block in walk_document_blocks(blocks):
        for inline in block.get("inlines") or []:
            if not isinstance(inline, dict):
                continue
            href = str(inline.get("href") or "").strip()
            if not href:
                continue
            annotations.append(
                {
                    "id": f"annotation:link:{annotation_index}",
                    "kind": "hyperlink",
                    "targets": [str(inline.get("id") or "")],
                    "href": href,
                    "text": str(inline.get("text") or ""),
                }
            )
            annotation_index += 1
    return annotations


def document_outline(blocks: List[Dict[str, object]]) -> List[Dict[str, object]]:
    outline: List[Dict[str, object]] = []
    for block in walk_document_blocks(blocks):
        if str(block.get("kind") or "") != "heading":
            continue
        outline.append(
            {
                "level": int(block.get("level") or 1),
                "text": document_block_text(block),
                "block_id": str(block.get("id") or ""),
            }
        )
    return outline


def document_analysis_payload(blocks: List[Dict[str, object]], source_backend: str) -> Dict[str, object]:
    all_blocks = list(walk_document_blocks(blocks))
    return {
        "source_backend": source_backend,
        "stats": {
            "block_count": len(all_blocks),
            "heading_count": sum(1 for block in all_blocks if str(block.get("kind") or "") == "heading"),
            "list_count": sum(1 for block in all_blocks if str(block.get("kind") or "") == "list"),
            "table_count": sum(1 for block in all_blocks if str(block.get("kind") or "") == "table"),
            "inline_count": sum(len(block.get("inlines") or []) for block in all_blocks),
        },
        "outline": document_outline(blocks),
    }


def document_projection_payload(blocks: List[Dict[str, object]]) -> Dict[str, object]:
    return {
        "outline-view": {
            "kind": "outline",
            "items": document_outline(blocks),
        },
        "reading-view": {
            "kind": "reading-order",
            "block_ids": [str(block.get("id") or "") for block in blocks if isinstance(block, dict)],
        },
        "backend:pandoc": {
            "kind": "backend-ast",
            "backend": "pandoc",
            "available": pandoc_available(),
        },
        "backend:docling": {
            "kind": "analysis-adapter",
            "backend": "docling",
            "available": docling_available(),
        },
    }


def document_artifact_payload(format_name: str, source_backend: str, package_part_count: int) -> Dict[str, object]:
    return {
        "native": {
            "backend": source_backend,
            "format": format_name,
            "role": "package-metadata",
            "package_part_count": package_part_count,
        }
    }


def read_docx_document(path: str) -> Dict[str, object]:
    metadata = path_metadata(path)
    if not metadata["exists"]:
        return {
            "kind": "document",
            "format": "docx",
            "path": path,
            **metadata,
            "warning": DOCUMENT_AUTOMATION_WARNING,
            "blocks": [],
            "annotations": [],
            "analysis": {
                "source_backend": "native-docx",
                "stats": {"block_count": 0, "heading_count": 0, "list_count": 0, "table_count": 0, "inline_count": 0},
                "outline": [],
            },
            "projections": document_projection_payload([]),
            "artifacts": document_artifact_payload("docx", "native-docx", 0),
            "paragraphs": [],
        }
    with zipfile.ZipFile(path, "r") as archive:
        style_names = read_docx_style_names(archive)
        numbering_formats = read_docx_numbering_formats(archive)
        relationship_targets = docx_relationship_targets(archive, "word/_rels/document.xml.rels")
        root = ET.fromstring(archive.read("word/document.xml"))
        package_part_count = len(archive.namelist())
    body = root.find(docx_tag("body"))
    blocks: List[Dict[str, object]] = []
    if body is not None:
        block_index = 1
        for child in body:
            child_name = local_name(child.tag)
            block_id = f"block:{block_index}"
            if child_name == "p":
                blocks.append(
                    parse_docx_paragraph_block(
                        child,
                        style_names,
                        numbering_formats,
                        relationship_targets,
                        block_id,
                    )
                )
                block_index += 1
            elif child_name == "tbl":
                blocks.append(
                    parse_docx_table_block(
                        child,
                        style_names,
                        numbering_formats,
                        relationship_targets,
                        block_id,
                    )
                )
                block_index += 1
    grouped_blocks = group_document_blocks(blocks)
    paragraphs = document_blocks_to_paragraphs(grouped_blocks)
    return {
        "kind": "document",
        "format": "docx",
        "path": path,
        **metadata,
        "warning": DOCUMENT_AUTOMATION_WARNING,
        "blocks": grouped_blocks,
        "annotations": collect_document_annotations(grouped_blocks),
        "analysis": document_analysis_payload(grouped_blocks, "native-docx"),
        "projections": document_projection_payload(grouped_blocks),
        "artifacts": document_artifact_payload("docx", "native-docx", package_part_count),
        "paragraphs": paragraphs,
    }


def build_doc_paragraph(text: str) -> ET.Element:
    paragraph = ET.Element(docx_tag("p"))
    run = ET.SubElement(paragraph, docx_tag("r"))
    text_node = ET.SubElement(run, docx_tag("t"))
    text_node.text = text
    return paragraph


def docx_append_text(run: ET.Element, text: str) -> None:
    text_node = ET.SubElement(run, docx_tag("t"))
    if text.startswith(" ") or text.endswith(" ") or "  " in text:
        text_node.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    text_node.text = text


def docx_run_for_text(text: str, marks: Optional[List[str]] = None) -> ET.Element:
    run = ET.Element(docx_tag("r"))
    mark_set = set(marks or [])
    if mark_set:
        properties = ET.SubElement(run, docx_tag("rPr"))
        if "strong" in mark_set:
            ET.SubElement(properties, docx_tag("b"))
        if "emphasis" in mark_set:
            ET.SubElement(properties, docx_tag("i"))
        if "underline" in mark_set:
            underline = ET.SubElement(properties, docx_tag("u"))
            underline.set(docx_attr("val"), "single")
        if "strike" in mark_set:
            ET.SubElement(properties, docx_tag("strike"))
        if "superscript" in mark_set:
            align = ET.SubElement(properties, docx_tag("vertAlign"))
            align.set(docx_attr("val"), "superscript")
        if "subscript" in mark_set:
            align = ET.SubElement(properties, docx_tag("vertAlign"))
            align.set(docx_attr("val"), "subscript")
    docx_append_text(run, text)
    return run


def ensure_docx_relationship(
    relationships_root: ET.Element,
    rel_type: str,
    target: str,
    target_mode: Optional[str] = None,
) -> str:
    for relationship in relationships_root:
        if local_name(relationship.tag) != "Relationship":
            continue
        if relationship.attrib.get("Type") != rel_type:
            continue
        if relationship.attrib.get("Target") != target:
            continue
        if (relationship.attrib.get("TargetMode") or "") != (target_mode or ""):
            continue
        rel_id = relationship.attrib.get("Id")
        if rel_id:
            return rel_id
    existing_ids = {
        str(relationship.attrib.get("Id") or "")
        for relationship in relationships_root
        if local_name(relationship.tag) == "Relationship"
    }
    index = 1
    while f"rId{index}" in existing_ids:
        index += 1
    rel_id = f"rId{index}"
    relationship = ET.SubElement(relationships_root, f"{{{PKG_REL_NS}}}Relationship")
    relationship.set("Id", rel_id)
    relationship.set("Type", rel_type)
    relationship.set("Target", target)
    if target_mode:
        relationship.set("TargetMode", target_mode)
    return rel_id


def append_docx_inline_nodes(paragraph: ET.Element, block: Dict[str, object], relationships_root: ET.Element) -> None:
    inlines = block.get("inlines")
    if not isinstance(inlines, list) or not inlines:
        text = str(block.get("text") or "")
        if text:
            paragraph.append(docx_run_for_text(text))
        return
    for inline in inlines:
        if not isinstance(inline, dict):
            continue
        kind = str(inline.get("kind") or "text")
        if kind == "line_break":
            run = ET.SubElement(paragraph, docx_tag("r"))
            ET.SubElement(run, docx_tag("br"))
            continue
        text = str(inline.get("text") or "")
        if not text:
            continue
        marks = [str(mark) for mark in inline.get("marks") or [] if str(mark).strip()]
        href = str(inline.get("href") or "").strip()
        if href:
            hyperlink = ET.SubElement(paragraph, docx_tag("hyperlink"))
            if href.startswith("#"):
                hyperlink.set(docx_attr("anchor"), href[1:])
            else:
                rel_id = ensure_docx_relationship(
                    relationships_root,
                    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
                    href,
                    target_mode="External",
                )
                hyperlink.set(f"{{{DOC_REL_NS}}}id", rel_id)
            hyperlink.append(docx_run_for_text(text, marks=marks))
            continue
        paragraph.append(docx_run_for_text(text, marks=marks))


def build_docx_paragraph_from_block(
    block: Dict[str, object],
    relationships_root: ET.Element,
    list_level: Optional[int] = None,
    ordered: bool = False,
) -> ET.Element:
    paragraph = ET.Element(docx_tag("p"))
    properties = ET.SubElement(paragraph, docx_tag("pPr"))
    kind = str(block.get("kind") or "paragraph")
    if kind == "heading":
        style = ET.SubElement(properties, docx_tag("pStyle"))
        level = min(max(int(block.get("level") or 1), 1), 6)
        style.set(docx_attr("val"), f"Heading{level}")
    elif str(block.get("style_ref") or "").strip():
        style = ET.SubElement(properties, docx_tag("pStyle"))
        style.set(docx_attr("val"), str(block.get("style_ref") or "").replace(" ", ""))
    if list_level is not None:
        numbering = ET.SubElement(properties, docx_tag("numPr"))
        ilvl = ET.SubElement(numbering, docx_tag("ilvl"))
        ilvl.set(docx_attr("val"), str(max(list_level, 0)))
        num_id = ET.SubElement(numbering, docx_tag("numId"))
        num_id.set(docx_attr("val"), "2" if ordered else "1")
    append_docx_inline_nodes(paragraph, block, relationships_root)
    if len(paragraph) == 1:
        paragraph.append(docx_run_for_text(str(block.get("text") or "")))
    return paragraph


def build_docx_table_from_block(block: Dict[str, object], relationships_root: ET.Element) -> ET.Element:
    table = ET.Element(docx_tag("tbl"))
    table_properties = ET.SubElement(table, docx_tag("tblPr"))
    table_width = ET.SubElement(table_properties, docx_tag("tblW"))
    table_width.set(docx_attr("w"), "0")
    table_width.set(docx_attr("type"), "auto")
    for row in block.get("rows") or []:
        if not isinstance(row, dict):
            continue
        table_row = ET.SubElement(table, docx_tag("tr"))
        for cell in row.get("cells") or []:
            if not isinstance(cell, dict):
                continue
            table_cell = ET.SubElement(table_row, docx_tag("tc"))
            ET.SubElement(table_cell, docx_tag("tcPr"))
            emitted = False
            for nested in cell.get("blocks") or []:
                if not isinstance(nested, dict):
                    continue
                for element in build_docx_block_elements(nested, relationships_root):
                    table_cell.append(element)
                    emitted = True
            if not emitted:
                table_cell.append(build_doc_paragraph(""))
    return table


def build_docx_list_elements(block: Dict[str, object], relationships_root: ET.Element) -> List[ET.Element]:
    elements: List[ET.Element] = []
    ordered = bool(block.get("ordered"))
    for item in block.get("items") or []:
        if not isinstance(item, dict):
            continue
        item_level = max(int(item.get("level") or 0), 0)
        nested_blocks = [nested for nested in item.get("blocks") or [] if isinstance(nested, dict)]
        if not nested_blocks:
            elements.append(
                build_docx_paragraph_from_block(
                    {"kind": "paragraph", "text": ""},
                    relationships_root,
                    list_level=item_level,
                    ordered=ordered,
                )
            )
            continue
        first_block = True
        for nested in nested_blocks:
            if first_block and str(nested.get("kind") or "") in {"paragraph", "heading", "code_block", "quote"}:
                elements.append(
                    build_docx_paragraph_from_block(
                        nested,
                        relationships_root,
                        list_level=item_level,
                        ordered=ordered,
                    )
                )
                first_block = False
                continue
            elements.extend(build_docx_block_elements(nested, relationships_root))
            first_block = False
    return elements


def build_docx_block_elements(block: Dict[str, object], relationships_root: ET.Element) -> List[ET.Element]:
    kind = str(block.get("kind") or "paragraph")
    if kind in {"paragraph", "heading", "quote", "code_block"}:
        return [build_docx_paragraph_from_block(block, relationships_root)]
    if kind == "list":
        return build_docx_list_elements(block, relationships_root)
    if kind == "table":
        return [build_docx_table_from_block(block, relationships_root)]
    return [build_docx_paragraph_from_block({"kind": "paragraph", "text": document_block_text(block)}, relationships_root)]


def docx_styles_xml_bytes() -> bytes:
    styles = ET.Element(docx_tag("styles"))
    definitions = [
        ("Normal", "Normal", True),
        ("Heading1", "heading 1", False),
        ("Heading2", "heading 2", False),
        ("Heading3", "heading 3", False),
        ("Heading4", "heading 4", False),
        ("Heading5", "heading 5", False),
        ("Heading6", "heading 6", False),
        ("Quote", "Quote", False),
    ]
    for style_id, name, is_default in definitions:
        style = ET.SubElement(styles, docx_tag("style"))
        style.set(docx_attr("type"), "paragraph")
        style.set(docx_attr("styleId"), style_id)
        if is_default:
            style.set(docx_attr("default"), "1")
        style_name = ET.SubElement(style, docx_tag("name"))
        style_name.set(docx_attr("val"), name)
        ET.SubElement(style, docx_tag("qFormat"))
    return ET.tostring(styles, encoding="utf-8", xml_declaration=True)


def docx_numbering_xml_bytes() -> bytes:
    numbering = ET.Element(docx_tag("numbering"))
    for abstract_id, num_fmt, lvl_text in (
        ("0", "bullet", "*"),
        ("1", "decimal", "%1."),
    ):
        abstract = ET.SubElement(numbering, docx_tag("abstractNum"))
        abstract.set(docx_attr("abstractNumId"), abstract_id)
        for level in range(0, 9):
            lvl = ET.SubElement(abstract, docx_tag("lvl"))
            lvl.set(docx_attr("ilvl"), str(level))
            start = ET.SubElement(lvl, docx_tag("start"))
            start.set(docx_attr("val"), "1")
            fmt = ET.SubElement(lvl, docx_tag("numFmt"))
            fmt.set(docx_attr("val"), num_fmt)
            text = ET.SubElement(lvl, docx_tag("lvlText"))
            text.set(docx_attr("val"), lvl_text if num_fmt == "bullet" else f"%{level + 1}.")
    for num_id, abstract_id in (("1", "0"), ("2", "1")):
        num = ET.SubElement(numbering, docx_tag("num"))
        num.set(docx_attr("numId"), num_id)
        abstract_ref = ET.SubElement(num, docx_tag("abstractNumId"))
        abstract_ref.set(docx_attr("val"), abstract_id)
    return ET.tostring(numbering, encoding="utf-8", xml_declaration=True)


def ensure_content_type_override(content_types_root: ET.Element, part_name: str, content_type: str) -> None:
    for override in content_types_root.findall(f"{{{CONTENT_TYPES_NS}}}Override"):
        if override.attrib.get("PartName") == part_name:
            override.attrib["ContentType"] = content_type
            return
    override = ET.SubElement(content_types_root, f"{{{CONTENT_TYPES_NS}}}Override")
    override.set("PartName", part_name)
    override.set("ContentType", content_type)


def minimal_docx_entries(blocks: List[Dict[str, object]]) -> Dict[str, bytes]:
    document = ET.Element(docx_tag("document"))
    body = ET.SubElement(document, docx_tag("body"))
    relationships = ET.Element(f"{{{PKG_REL_NS}}}Relationships")
    ensure_docx_relationship(
        relationships,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles",
        "styles.xml",
    )
    ensure_docx_relationship(
        relationships,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/numbering",
        "numbering.xml",
    )
    for block in blocks:
        if not isinstance(block, dict):
            continue
        for element in build_docx_block_elements(block, relationships):
            body.append(element)
    ET.SubElement(body, docx_tag("sectPr"))

    content_types = ET.Element(f"{{{CONTENT_TYPES_NS}}}Types")
    default_rels = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Default")
    default_rels.set("Extension", "rels")
    default_rels.set("ContentType", "application/vnd.openxmlformats-package.relationships+xml")
    default_xml = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Default")
    default_xml.set("Extension", "xml")
    default_xml.set("ContentType", "application/xml")
    document_override = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Override")
    document_override.set("PartName", "/word/document.xml")
    document_override.set(
        "ContentType",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
    )
    styles_override = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Override")
    styles_override.set("PartName", "/word/styles.xml")
    styles_override.set(
        "ContentType",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml",
    )
    numbering_override = ET.SubElement(content_types, f"{{{CONTENT_TYPES_NS}}}Override")
    numbering_override.set("PartName", "/word/numbering.xml")
    numbering_override.set(
        "ContentType",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.numbering+xml",
    )

    rels = ET.Element(f"{{{PKG_REL_NS}}}Relationships")
    rel = ET.SubElement(rels, f"{{{PKG_REL_NS}}}Relationship")
    rel.set("Id", "rId1")
    rel.set(
        "Type",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument",
    )
    rel.set("Target", "word/document.xml")

    return {
        "[Content_Types].xml": ET.tostring(content_types, encoding="utf-8", xml_declaration=True),
        "_rels/.rels": ET.tostring(rels, encoding="utf-8", xml_declaration=True),
        "word/_rels/document.xml.rels": ET.tostring(relationships, encoding="utf-8", xml_declaration=True),
        "word/document.xml": ET.tostring(document, encoding="utf-8", xml_declaration=True),
        "word/styles.xml": docx_styles_xml_bytes(),
        "word/numbering.xml": docx_numbering_xml_bytes(),
    }


def write_docx_document(path: str, blocks: List[Dict[str, object]], expected_etag: Optional[str]) -> Dict[str, object]:
    assert_expected_etag(path, expected_etag)
    normalized_blocks = [block for block in blocks if isinstance(block, dict)]
    if not normalized_blocks:
        normalized_blocks = [{"id": "block:1", "kind": "paragraph", "text": ""}]
    if os.path.exists(path):
        with zipfile.ZipFile(path, "r") as archive:
            entries = {name: archive.read(name) for name in archive.namelist()}
        root = ET.fromstring(entries["word/document.xml"])
        relationships = ET.fromstring(
            entries.get("word/_rels/document.xml.rels")
            or ET.tostring(ET.Element(f"{{{PKG_REL_NS}}}Relationships"), encoding="utf-8", xml_declaration=True)
        )
        body = root.find(docx_tag("body"))
        if body is None:
            body = ET.SubElement(root, docx_tag("body"))
        section = None
        if len(body) > 0 and local_name(body[-1].tag) == "sectPr":
            section = body[-1]
        body[:] = []
        for block in normalized_blocks:
            for element in build_docx_block_elements(block, relationships):
                body.append(element)
        if section is None:
            section = ET.Element(docx_tag("sectPr"))
        body.append(section)
        entries["word/_rels/document.xml.rels"] = ET.tostring(
            relationships,
            encoding="utf-8",
            xml_declaration=True,
        )
        entries["word/document.xml"] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
        if "word/styles.xml" not in entries:
            entries["word/styles.xml"] = docx_styles_xml_bytes()
        entries["word/numbering.xml"] = docx_numbering_xml_bytes()
        content_types = ET.fromstring(entries["[Content_Types].xml"])
        ensure_content_type_override(
            content_types,
            "/word/document.xml",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
        )
        ensure_content_type_override(
            content_types,
            "/word/styles.xml",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml",
        )
        ensure_content_type_override(
            content_types,
            "/word/numbering.xml",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.numbering+xml",
        )
        entries["[Content_Types].xml"] = ET.tostring(content_types, encoding="utf-8", xml_declaration=True)
        relationships = ET.fromstring(entries["word/_rels/document.xml.rels"])
        ensure_docx_relationship(
            relationships,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles",
            "styles.xml",
        )
        ensure_docx_relationship(
            relationships,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/numbering",
            "numbering.xml",
        )
        entries["word/_rels/document.xml.rels"] = ET.tostring(
            relationships,
            encoding="utf-8",
            xml_declaration=True,
        )
    else:
        entries = minimal_docx_entries(normalized_blocks)

    temp_fd, temp_path = tempfile.mkstemp(prefix=".entropic-office-", dir=posixpath.dirname(path))
    os.close(temp_fd)
    try:
        with zipfile.ZipFile(temp_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for name in sorted(entries):
                archive.writestr(name, entries[name])
        os.replace(temp_path, path)
    finally:
        try:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
        except OSError:
            pass
    return read_docx_document(path)


PRESENTATION_AUTOMATION_WARNING = (
    "Structured presentation automation edits slide text, notes, layout references, object ordering, table cell text, "
    "image placement, basic object frames, portable shape style, image crop, and chart title/data on existing charts. "
    "Themes, deep chart formatting, and animation semantics are preserved when possible but are not yet fully modeled."
)

PPTX_TITLE_PLACEHOLDER_TYPES = {"TITLE", "CENTER_TITLE"}
PPTX_BODY_PLACEHOLDER_TYPES = {"BODY", "SUBTITLE", "OBJECT"}
PPTX_NOTES_IGNORED_PLACEHOLDER_TYPES = {
    "HEADER",
    "DATE",
    "FOOTER",
    "SLIDE_NUMBER",
    "HDR",
    "DT",
    "FTR",
    "SLDNUM",
}


def python_pptx_available() -> bool:
    return PptxPresentation is not None and Inches is not None


def emu_frame_dict(x: object, y: object, w: object, h: object) -> Dict[str, object]:
    return {
        "x": int(x or 0),
        "y": int(y or 0),
        "w": int(w or 0),
        "h": int(h or 0),
        "unit": "emu",
    }


def normalize_presentation_frame(frame: object) -> Optional[Dict[str, object]]:
    if not isinstance(frame, dict):
        return None
    if not any(key in frame for key in {"x", "y", "w", "h", "left", "top", "width", "height"}):
        return None
    return emu_frame_dict(
        frame.get("x", frame.get("left", 0)),
        frame.get("y", frame.get("top", 0)),
        frame.get("w", frame.get("width", 0)),
        frame.get("h", frame.get("height", 0)),
    )


def normalize_layout_ref(raw: object) -> Optional[str]:
    text = str(raw or "").strip()
    if not text:
        return None
    if text.startswith("layout:"):
        return text
    return f"layout:{slugify_identifier(text, 'layout')}"


def normalize_master_ref(raw: object) -> Optional[str]:
    text = str(raw or "").strip()
    if not text:
        return None
    if text.startswith("master:"):
        return text
    return f"master:{slugify_identifier(text, 'master')}"


def body_items_from_outline(items: object) -> List[Dict[str, object]]:
    normalized: List[Dict[str, object]] = []
    if not isinstance(items, list):
        return normalized
    for item in items:
        if isinstance(item, dict):
            text = str(item.get("text") or "").strip()
            if not text:
                continue
            level = int(item.get("level") or 0)
        else:
            text = str(item or "").strip()
            if not text:
                continue
            level = 0
        normalized.append({"text": text, "level": max(level, 0)})
    return normalized


def normalize_presentation_body(body: object) -> Dict[str, object]:
    if isinstance(body, dict):
        kind = str(body.get("kind") or "").strip()
        if kind == "outline" or isinstance(body.get("items"), list):
            items = body_items_from_outline(body.get("items"))
            return {"kind": "outline", "items": items}
        bullets = body.get("bullets")
        if isinstance(bullets, list):
            return {
                "kind": "bullets",
                "bullets": [str(value) for value in bullets if str(value or "").strip()],
            }
        text = str(body.get("text") or "").strip()
        return {"kind": "text", "text": text}
    if isinstance(body, (list, tuple)):
        bullets = [str(value) for value in body if str(value or "").strip()]
        return {"kind": "bullets", "bullets": bullets}
    text = str(body or "").strip()
    return {"kind": "text", "text": text}


def presentation_paragraphs_from_body(body: object) -> List[Dict[str, object]]:
    normalized = normalize_presentation_body(body)
    kind = str(normalized.get("kind") or "").strip()
    if kind == "outline":
        return body_items_from_outline(normalized.get("items"))
    if kind == "bullets":
        return [{"text": value, "level": 0} for value in normalized.get("bullets") or []]
    text = str(normalized.get("text") or "").strip()
    return [{"text": line, "level": 0} for line in text.splitlines() if line.strip()]


def presentation_body_from_paragraphs(paragraphs: List[Dict[str, object]]) -> Dict[str, object]:
    cleaned = body_items_from_outline(paragraphs)
    if not cleaned:
        return {"kind": "text", "text": ""}
    if len(cleaned) == 1 and int(cleaned[0].get("level") or 0) == 0:
        return {"kind": "text", "text": str(cleaned[0].get("text") or "")}
    return {"kind": "outline", "items": cleaned}


def normalize_presentation_notes(notes: object) -> List[str]:
    if isinstance(notes, list):
        return [str(value) for value in notes if str(value or "").strip()]
    text = str(notes or "").strip()
    return [text] if text else []


def presentation_text_value(body: object) -> str:
    paragraphs = presentation_paragraphs_from_body(body)
    return "\n".join(str(item.get("text") or "") for item in paragraphs if str(item.get("text") or "").strip())


def presentation_body_lines(body: Dict[str, object]) -> List[str]:
    return [str(item.get("text") or "") for item in presentation_paragraphs_from_body(body)]


def normalize_string_matrix(raw: object) -> List[List[str]]:
    rows: List[List[str]] = []
    if not isinstance(raw, list):
        return rows
    for raw_row in raw:
        if not isinstance(raw_row, list):
            continue
        rows.append([str(value or "") for value in raw_row])
    return rows


def normalize_hex_color(raw: object) -> Optional[str]:
    text = str(raw or "").strip()
    if not text:
        return None
    if text.startswith("#"):
        text = text[1:]
    if re.fullmatch(r"[0-9A-Fa-f]{6}", text):
        return f"#{text.upper()}"
    return None


def normalize_presentation_paint(raw: object) -> Optional[Dict[str, object]]:
    if not isinstance(raw, dict):
        return None
    result: Dict[str, object] = {}
    paint_type = str(raw.get("type") or raw.get("kind") or "").strip()
    if paint_type:
        result["type"] = paint_type
    color = normalize_hex_color(raw.get("color") or raw.get("rgb"))
    if color:
        result["color"] = color
    theme_color = str(raw.get("theme_color") or raw.get("themeColor") or "").strip()
    if theme_color:
        result["theme_color"] = theme_color
    if raw.get("width") is not None:
        try:
            result["width"] = int(float(str(raw.get("width") or 0)))
        except (TypeError, ValueError):
            pass
    if raw.get("transparency") is not None:
        try:
            result["transparency"] = float(str(raw.get("transparency") or 0))
        except (TypeError, ValueError):
            pass
    return result or None


def normalize_presentation_style(raw: object) -> Optional[Dict[str, object]]:
    if not isinstance(raw, dict):
        return None
    result: Dict[str, object] = {}
    fill = normalize_presentation_paint(raw.get("fill"))
    if fill:
        result["fill"] = fill
    stroke = normalize_presentation_paint(raw.get("stroke") or raw.get("line"))
    if stroke:
        result["stroke"] = stroke
    return result or None


def normalize_presentation_crop(raw: object) -> Optional[Dict[str, object]]:
    if not isinstance(raw, dict):
        return None
    result: Dict[str, object] = {}
    for edge in ("left", "top", "right", "bottom"):
        if raw.get(edge) is None:
            continue
        try:
            result[edge] = float(str(raw.get(edge) or 0))
        except (TypeError, ValueError):
            continue
    return result or None


def rgb_hex_from_pptx_color(color_value: object) -> Optional[str]:
    text = str(color_value or "").strip()
    return normalize_hex_color(text)


def pptx_color_payload(color_format: object) -> Optional[Dict[str, object]]:
    if color_format is None:
        return None
    result: Dict[str, object] = {}
    try:
        rgb_value = getattr(color_format, "rgb", None)
    except Exception:
        rgb_value = None
    color = rgb_hex_from_pptx_color(rgb_value)
    if color:
        result["color"] = color
    try:
        theme_value = getattr(color_format, "theme_color", None)
    except Exception:
        theme_value = None
    theme_name = str(getattr(theme_value, "name", theme_value) or "").strip()
    if theme_name:
        result["theme_color"] = theme_name
    return result or None


def pptx_fill_payload(shape) -> Optional[Dict[str, object]]:
    try:
        fill = shape.fill
    except Exception:
        return None
    if fill is None:
        return None
    fill_type = str(getattr(getattr(fill, "type", None), "name", "") or "").strip().lower()
    result: Dict[str, object] = {}
    if fill_type:
        result["type"] = fill_type
    color = None
    try:
        color = pptx_color_payload(fill.fore_color)
    except Exception:
        color = None
    if color:
        result.update(color)
    return result or None


def pptx_line_payload(shape) -> Optional[Dict[str, object]]:
    try:
        line = shape.line
    except Exception:
        return None
    if line is None:
        return None
    result = pptx_color_payload(getattr(line, "color", None)) or {}
    try:
        if getattr(line, "width", None) is not None:
            result["width"] = int(line.width)
    except Exception:
        pass
    return result or None


def pptx_style_payload(shape) -> Optional[Dict[str, object]]:
    result: Dict[str, object] = {}
    fill = pptx_fill_payload(shape)
    if fill:
        result["fill"] = fill
    stroke = pptx_line_payload(shape)
    if stroke:
        result["stroke"] = stroke
    return result or None


def pptx_shape_rotation(shape) -> Optional[float]:
    try:
        rotation = getattr(shape, "rotation", None)
    except Exception:
        return None
    if rotation is None:
        return None
    try:
        return float(rotation)
    except (TypeError, ValueError):
        return None


def pptx_picture_crop_payload(shape) -> Optional[Dict[str, object]]:
    result: Dict[str, object] = {}
    for source, key in (
        ("crop_left", "left"),
        ("crop_top", "top"),
        ("crop_right", "right"),
        ("crop_bottom", "bottom"),
    ):
        try:
            value = getattr(shape, source, None)
        except Exception:
            value = None
        if value is None:
            continue
        try:
            result[key] = float(value)
        except (TypeError, ValueError):
            continue
    return result or None


def pptx_placeholder_index(shape) -> Optional[int]:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        return int(shape.placeholder_format.idx)
    except Exception:
        return None


def pptx_auto_shape_name(shape) -> Optional[str]:
    shape_type_name = str(getattr(getattr(shape, "shape_type", None), "name", "") or "").upper()
    if shape_type_name != "AUTO_SHAPE":
        return None
    try:
        auto_shape = getattr(shape, "auto_shape_type", None)
    except Exception:
        auto_shape = None
    if auto_shape is None:
        return None
    return str(getattr(auto_shape, "name", auto_shape) or "").strip().upper() or None


def xml_solid_color_payload(node: Optional[ET.Element]) -> Optional[Dict[str, object]]:
    if node is None:
        return None
    result: Dict[str, object] = {}
    srgb = node.find(drawing_tag("srgbClr"))
    if srgb is not None:
        color = normalize_hex_color(srgb.attrib.get("val"))
        if color:
            result["color"] = color
    scheme = node.find(drawing_tag("schemeClr"))
    if scheme is not None:
        theme_color = str(scheme.attrib.get("val") or "").strip()
        if theme_color:
            result["theme_color"] = theme_color
    return result or None


def xml_shape_style(shape: ET.Element) -> Optional[Dict[str, object]]:
    style: Dict[str, object] = {}
    sp_pr = shape.find(pptx_tag("spPr"))
    if sp_pr is None:
        return None
    solid_fill = xml_solid_color_payload(sp_pr.find(drawing_tag("solidFill")))
    if solid_fill:
        style["fill"] = {"type": "solid", **solid_fill}
    line = sp_pr.find(drawing_tag("ln"))
    if line is not None:
        stroke = xml_solid_color_payload(line.find(drawing_tag("solidFill"))) or {}
        width = line.attrib.get("w")
        if width:
            try:
                stroke["width"] = int(width)
            except ValueError:
                pass
        if stroke:
            style["stroke"] = stroke
    return style or None


def xml_shape_rotation(shape: ET.Element) -> Optional[float]:
    xfrm = shape.find(f".//{drawing_tag('xfrm')}")
    if xfrm is None:
        return None
    raw_rotation = xfrm.attrib.get("rot")
    if not raw_rotation:
        return None
    try:
        return round(int(raw_rotation) / 60000.0, 3)
    except ValueError:
        return None


def xml_shape_crop(shape: ET.Element) -> Optional[Dict[str, object]]:
    src_rect = shape.find(f".//{drawing_tag('srcRect')}")
    if src_rect is None:
        return None
    crop: Dict[str, object] = {}
    for attribute, edge in (("l", "left"), ("t", "top"), ("r", "right"), ("b", "bottom")):
        raw_value = src_rect.attrib.get(attribute)
        if raw_value is None:
            continue
        try:
            crop[edge] = int(raw_value) / 100000.0
        except ValueError:
            continue
    return crop or None


def xml_auto_shape_name(shape: ET.Element) -> Optional[str]:
    sp_pr = shape.find(pptx_tag("spPr"))
    if sp_pr is None:
        return None
    geometry = sp_pr.find(drawing_tag("prstGeom"))
    if geometry is None:
        return None
    return str(geometry.attrib.get("prst") or "").strip().upper() or None


def xml_placeholder_index(shape: ET.Element) -> Optional[int]:
    for path in (
        f"{pptx_tag('nvSpPr')}/{pptx_tag('nvPr')}/{pptx_tag('ph')}",
        f"{pptx_tag('nvPicPr')}/{pptx_tag('nvPr')}/{pptx_tag('ph')}",
        f"{pptx_tag('nvGraphicFramePr')}/{pptx_tag('nvPr')}/{pptx_tag('ph')}",
    ):
        placeholder = shape.find(path)
        if placeholder is None:
            continue
        raw_index = placeholder.attrib.get("idx")
        if raw_index is None:
            continue
        try:
            return int(raw_index)
        except ValueError:
            return None
    return None


def normalize_chart_series(raw: object) -> List[Dict[str, object]]:
    normalized: List[Dict[str, object]] = []
    if not isinstance(raw, list):
        return normalized
    for index, item in enumerate(raw, start=1):
        if not isinstance(item, dict):
            continue
        entry: Dict[str, object] = {
            "id": str(item.get("id") or f"series:{index}"),
        }
        name = str(item.get("name") or "").strip()
        if name:
            entry["name"] = name
        values = item.get("values")
        if isinstance(values, list):
            entry["values"] = [str(value) if isinstance(value, bool) else value for value in values]
        if len(entry) > 1:
            normalized.append(entry)
    return normalized


def normalize_presentation_transition(raw: object) -> Optional[Dict[str, object]]:
    if not isinstance(raw, dict):
        return None
    result: Dict[str, object] = {}
    effect = str(raw.get("effect") or "").strip()
    if effect:
        result["effect"] = effect
    speed = str(raw.get("speed") or "").strip()
    if speed:
        result["speed"] = speed
    if raw.get("advance_on_click") is not None:
        result["advance_on_click"] = bool(raw.get("advance_on_click"))
    if raw.get("advance_after_ms") is not None:
        try:
            result["advance_after_ms"] = int(raw.get("advance_after_ms") or 0)
        except (TypeError, ValueError):
            pass
    if raw.get("duration_ms") is not None:
        try:
            result["duration_ms"] = int(raw.get("duration_ms") or 0)
        except (TypeError, ValueError):
            pass
    return result or None


def normalize_presentation_timeline(raw: object) -> Optional[Dict[str, object]]:
    if not isinstance(raw, dict):
        return None
    result: Dict[str, object] = {}
    if raw.get("has_main_sequence") is not None:
        result["has_main_sequence"] = bool(raw.get("has_main_sequence"))
    if raw.get("effect_count") is not None:
        try:
            result["effect_count"] = int(raw.get("effect_count") or 0)
        except (TypeError, ValueError):
            pass
    return result or None


def normalize_presentation_object_payload(
    raw_object: object,
    slide_index: int,
    object_index: int,
) -> Optional[Dict[str, object]]:
    if not isinstance(raw_object, dict):
        return None
    kind = str(raw_object.get("kind") or "").strip() or "shape"
    result: Dict[str, object] = {
        "id": str(raw_object.get("id") or f"object:{slide_index}:{object_index}"),
        "kind": kind,
    }
    name = str(raw_object.get("name") or "").strip()
    if name:
        result["name"] = name
    shape_type = str(raw_object.get("shape_type") or "").strip().upper()
    if shape_type:
        result["shape_type"] = shape_type
    auto_shape = str(raw_object.get("auto_shape") or raw_object.get("autoShape") or "").strip().upper()
    if auto_shape:
        result["auto_shape"] = auto_shape
    placeholder_kind = str(raw_object.get("placeholder_kind") or "").strip().upper()
    if placeholder_kind:
        result["placeholder_kind"] = placeholder_kind
    if raw_object.get("placeholder_index") is not None:
        try:
            result["placeholder_index"] = int(raw_object.get("placeholder_index") or 0)
        except (TypeError, ValueError):
            pass
    frame = normalize_presentation_frame(raw_object.get("frame"))
    if frame:
        result["frame"] = frame
    if raw_object.get("rotation") is not None:
        try:
            result["rotation"] = float(str(raw_object.get("rotation") or 0))
        except (TypeError, ValueError):
            pass
    if raw_object.get("z_index") is not None:
        try:
            result["z_index"] = int(raw_object.get("z_index") or 0)
        except (TypeError, ValueError):
            pass
    if "text" in raw_object:
        result["text"] = str(raw_object.get("text") or "")
    if "body" in raw_object or "bullets" in raw_object or "items" in raw_object:
        result["body"] = normalize_presentation_body(raw_object.get("body", raw_object))
    image_ref = str(raw_object.get("image_ref") or "").strip()
    if image_ref:
        result["image_ref"] = image_ref
    image_name = str(raw_object.get("image_name") or "").strip()
    if image_name:
        result["image_name"] = image_name
    alt_text = str(raw_object.get("alt_text") or "").strip()
    if alt_text:
        result["alt_text"] = alt_text
    style = normalize_presentation_style(raw_object.get("style"))
    if style:
        result["style"] = style
    crop = normalize_presentation_crop(raw_object.get("crop"))
    if crop:
        result["crop"] = crop
    chart_kind = str(raw_object.get("chart_kind") or "").strip()
    if chart_kind:
        result["chart_kind"] = chart_kind
    chart_title = str(raw_object.get("title") or "").strip() if kind == "chart" else ""
    if chart_title:
        result["title"] = chart_title
    categories = raw_object.get("categories")
    if isinstance(categories, list):
        result["categories"] = [str(value or "") for value in categories]
    series = normalize_chart_series(raw_object.get("series"))
    if series:
        result["series"] = series
    table = raw_object.get("table")
    if isinstance(table, dict):
        next_table: Dict[str, object] = {}
        if table.get("rows") is not None:
            next_table["rows"] = int(table.get("rows") or 0)
        if table.get("cols") is not None:
            next_table["cols"] = int(table.get("cols") or 0)
        cells = normalize_string_matrix(table.get("cells"))
        if cells:
            next_table["cells"] = cells
            next_table["rows"] = max(int(next_table.get("rows") or 0), len(cells))
            next_table["cols"] = max(
                int(next_table.get("cols") or 0),
                max((len(row) for row in cells), default=0),
            )
        if next_table:
            result["table"] = next_table
    return result


def derive_presentation_title(objects: List[Dict[str, object]]) -> Optional[str]:
    for obj in objects:
        if str(obj.get("kind") or "") == "title" or str(obj.get("placeholder_kind") or "") in PPTX_TITLE_PLACEHOLDER_TYPES:
            if "text" in obj:
                return str(obj.get("text") or "")
            if "body" in obj:
                return presentation_text_value(obj.get("body"))
    return None


def derive_presentation_body(objects: List[Dict[str, object]]) -> Optional[Dict[str, object]]:
    for obj in objects:
        if str(obj.get("placeholder_kind") or "") in PPTX_BODY_PLACEHOLDER_TYPES:
            if "body" in obj:
                return normalize_presentation_body(obj.get("body"))
            if "text" in obj:
                return {"kind": "text", "text": str(obj.get("text") or "")}
    for obj in objects:
        if str(obj.get("kind") or "") == "text_box":
            if "body" in obj:
                return normalize_presentation_body(obj.get("body"))
            if "text" in obj:
                return {"kind": "text", "text": str(obj.get("text") or "")}
    return None


def normalize_presentation_slide_payload(
    raw_slide: object,
    index: int,
    derive_summary: bool = False,
) -> Dict[str, object]:
    result: Dict[str, object] = {
        "id": f"slide:{index}",
        "kind": "slide",
        "index": index,
    }
    if not isinstance(raw_slide, dict):
        return result
    result["id"] = str(raw_slide.get("id") or f"slide:{index}")
    layout_ref = normalize_layout_ref(raw_slide.get("layout_ref") or raw_slide.get("layout"))
    if layout_ref:
        result["layout_ref"] = layout_ref
    master_ref = normalize_master_ref(raw_slide.get("master_ref"))
    if master_ref:
        result["master_ref"] = master_ref
    if "title" in raw_slide:
        result["title"] = str(raw_slide.get("title") or "")
    if "body" in raw_slide:
        result["body"] = normalize_presentation_body(raw_slide.get("body"))
    if "notes" in raw_slide:
        result["notes"] = normalize_presentation_notes(raw_slide.get("notes"))
    transition = normalize_presentation_transition(raw_slide.get("transition"))
    if transition:
        result["transition"] = transition
    timeline = normalize_presentation_timeline(raw_slide.get("timeline"))
    if timeline:
        result["timeline"] = timeline
    remove_object_ids = [
        str(value).strip()
        for value in (raw_slide.get("remove_object_ids") or [])
        if str(value or "").strip()
    ]
    if remove_object_ids:
        result["remove_object_ids"] = remove_object_ids
    if raw_slide.get("reorder_objects") is True:
        result["reorder_objects"] = True
    if raw_slide.get("prune_missing_objects") is True:
        result["prune_missing_objects"] = True
    objects = []
    for object_index, raw_object in enumerate(raw_slide.get("objects") or [], start=1):
        normalized = normalize_presentation_object_payload(raw_object, index, object_index)
        if normalized is not None:
            objects.append(normalized)
    if objects:
        result["objects"] = objects
    if derive_summary and objects:
        if "title" not in result:
            derived_title = derive_presentation_title(objects)
            if derived_title is not None:
                result["title"] = derived_title
        if "body" not in result:
            derived_body = derive_presentation_body(objects)
            if derived_body is not None:
                result["body"] = derived_body
    return result


def drawing_paragraph_texts(node: ET.Element) -> List[str]:
    paragraphs: List[str] = []
    for paragraph in node.findall(f".//{drawing_tag('p')}"):
        text_parts = []
        for text_node in paragraph.iter():
            if local_name(text_node.tag) == "t":
                text_parts.append(text_node.text or "")
        text = "".join(text_parts).strip()
        if text:
            paragraphs.append(text)
    return paragraphs


def xml_transition_payload(slide_root: ET.Element) -> Optional[Dict[str, object]]:
    transition = slide_root.find(pptx_tag("transition"))
    if transition is None:
        return None
    payload: Dict[str, object] = {}
    speed = str(transition.attrib.get("spd") or "").strip()
    if speed:
        payload["speed"] = speed
    adv_click = transition.attrib.get("advClick")
    if adv_click is not None:
        payload["advance_on_click"] = adv_click not in {"0", "false", "False"}
    adv_time = transition.attrib.get("advTm")
    if adv_time:
        try:
            payload["advance_after_ms"] = int(adv_time)
        except ValueError:
            pass
    for child in list(transition):
        name = local_name(child.tag)
        if name != "extLst":
            payload["effect"] = name
            break
    return payload or None


def xml_timeline_payload(slide_root: ET.Element) -> Optional[Dict[str, object]]:
    timing = slide_root.find(pptx_tag("timing"))
    if timing is None:
        return None
    ctn_count = sum(1 for node in timing.iter() if local_name(node.tag) == "cTn")
    return {
        "has_main_sequence": ctn_count > 0,
        "effect_count": max(ctn_count - 1, 0),
    }


def xml_placeholder_type(shape: ET.Element) -> Optional[str]:
    placeholder = shape.find(
        f"{pptx_tag('nvSpPr')}/{pptx_tag('nvPr')}/{pptx_tag('ph')}"
    )
    if placeholder is None:
        return None
    raw_type = str(placeholder.attrib.get("type") or "").strip()
    return (raw_type or "body").upper()


def xml_shape_frame(shape: ET.Element) -> Optional[Dict[str, object]]:
    xfrm = shape.find(f".//{drawing_tag('xfrm')}")
    if xfrm is None:
        return None
    off = xfrm.find(drawing_tag("off"))
    ext = xfrm.find(drawing_tag("ext"))
    if off is None or ext is None:
        return None
    return emu_frame_dict(
        off.attrib.get("x", 0),
        off.attrib.get("y", 0),
        ext.attrib.get("cx", 0),
        ext.attrib.get("cy", 0),
    )


def xml_shape_non_visual(shape: ET.Element) -> Tuple[str, str]:
    paths = (
        (pptx_tag("nvSpPr"), pptx_tag("cNvPr")),
        (pptx_tag("nvPicPr"), pptx_tag("cNvPr")),
        (pptx_tag("nvGraphicFramePr"), pptx_tag("cNvPr")),
    )
    for first, second in paths:
        node = shape.find(f"{first}/{second}")
        if node is not None:
            return (
                str(node.attrib.get("id") or "").strip(),
                str(node.attrib.get("name") or "").strip(),
            )
    return ("", "")


def xml_shape_object(shape: ET.Element, fallback_id: str) -> Optional[Dict[str, object]]:
    object_id_value, name = xml_shape_non_visual(shape)
    object_id = f"shape:{object_id_value or fallback_id}"
    frame = xml_shape_frame(shape)
    placeholder_kind = xml_placeholder_type(shape)
    placeholder_index = xml_placeholder_index(shape)
    rotation = xml_shape_rotation(shape)
    style = xml_shape_style(shape)

    if local_name(shape.tag) == "sp":
        body = presentation_body_from_paragraphs(
            [{"text": paragraph, "level": 0} for paragraph in drawing_paragraph_texts(shape)]
        )
        if presentation_text_value(body) or name:
            auto_shape = xml_auto_shape_name(shape)
            if placeholder_kind in PPTX_TITLE_PLACEHOLDER_TYPES:
                kind = "title"
            elif auto_shape and placeholder_kind not in PPTX_BODY_PLACEHOLDER_TYPES:
                kind = "shape"
            else:
                kind = "text_box"
            result: Dict[str, object] = {"id": object_id, "kind": kind}
            if name:
                result["name"] = name
            result["shape_type"] = "AUTO_SHAPE" if auto_shape else "TEXT_BOX"
            if auto_shape:
                result["auto_shape"] = auto_shape
            if placeholder_kind:
                result["placeholder_kind"] = placeholder_kind
            if placeholder_index is not None:
                result["placeholder_index"] = placeholder_index
            if frame:
                result["frame"] = frame
            if rotation is not None:
                result["rotation"] = rotation
            if style:
                result["style"] = style
            if kind == "title":
                result["text"] = presentation_text_value(body)
            else:
                result["body"] = body
            return result
        return None

    if local_name(shape.tag) == "pic":
        result = {"id": object_id, "kind": "image"}
        if name:
            result["name"] = name
            result["image_name"] = name
        if frame:
            result["frame"] = frame
        if rotation is not None:
            result["rotation"] = rotation
        if style:
            result["style"] = style
        crop = xml_shape_crop(shape)
        if crop:
            result["crop"] = crop
        return result

    if local_name(shape.tag) == "graphicFrame":
        result: Dict[str, object] = {"id": object_id, "kind": "shape"}
        if name:
            result["name"] = name
        if frame:
            result["frame"] = frame
        if rotation is not None:
            result["rotation"] = rotation
        if shape.find(f".//{drawing_tag('tbl')}") is not None:
            rows = len(shape.findall(f".//{drawing_tag('tr')}"))
            cols = 0
            cells: List[List[str]] = []
            first_row = shape.find(f".//{drawing_tag('tr')}")
            if first_row is not None:
                cols = len(first_row.findall(drawing_tag("tc")))
            for row in shape.findall(f".//{drawing_tag('tr')}"):
                cell_row: List[str] = []
                for cell in row.findall(drawing_tag("tc")):
                    cell_row.append("\n".join(drawing_paragraph_texts(cell)))
                if cell_row:
                    cells.append(cell_row)
            result["kind"] = "table"
            result["table"] = {"rows": rows, "cols": cols}
            if cells:
                result["table"]["cells"] = cells
            return result
        for node in shape.iter():
            if local_name(node.tag) == "chart":
                result["kind"] = "chart"
                return result
    return None


def read_pptx_notes_xml(archive: zipfile.ZipFile, slide_path: str) -> List[str]:
    rels_path = posixpath.join(
        posixpath.dirname(slide_path),
        "_rels",
        f"{posixpath.basename(slide_path)}.rels",
    )
    if rels_path not in archive.namelist():
        return []
    rels_root = ET.fromstring(archive.read(rels_path))
    notes_target = None
    for rel in rels_root:
        if local_name(rel.tag) != "Relationship":
            continue
        rel_type = str(rel.attrib.get("Type") or "")
        if rel_type.endswith("/notesSlide"):
            target = rel.attrib.get("Target")
            if target:
                notes_target = normalize_relationship_target(posixpath.dirname(slide_path), target)
                break
    if not notes_target or notes_target not in archive.namelist():
        return []
    notes_root = ET.fromstring(archive.read(notes_target))
    notes: List[str] = []
    for shape in notes_root.findall(f".//{pptx_tag('sp')}"):
        placeholder_type = xml_placeholder_type(shape)
        if placeholder_type in PPTX_NOTES_IGNORED_PLACEHOLDER_TYPES:
            continue
        notes.extend(drawing_paragraph_texts(shape))
    return notes


def read_pptx_document_xml(path: str, metadata: Dict[str, object]) -> Dict[str, object]:
    with zipfile.ZipFile(path, "r") as archive:
        presentation_xml = "ppt/presentation.xml"
        presentation_rels_xml = "ppt/_rels/presentation.xml.rels"
        if presentation_xml not in archive.namelist() or presentation_rels_xml not in archive.namelist():
            raise RuntimeError("The presentation is missing required PPTX metadata.")
        presentation_root = ET.fromstring(archive.read(presentation_xml))
        presentation_rels_root = ET.fromstring(archive.read(presentation_rels_xml))
        slide_targets: Dict[str, str] = {}
        for rel in presentation_rels_root:
            if local_name(rel.tag) != "Relationship":
                continue
            rel_id = rel.attrib.get("Id")
            target = rel.attrib.get("Target")
            if rel_id and target:
                slide_targets[rel_id] = normalize_relationship_target("ppt", target)

        slides: List[Dict[str, object]] = []
        slide_id_list = presentation_root.find(pptx_tag("sldIdLst"))
        if slide_id_list is not None:
            for index, slide_ref in enumerate(slide_id_list.findall(pptx_tag("sldId")), start=1):
                slide_stable_id = str(slide_ref.attrib.get("id") or index)
                rel_id = slide_ref.attrib.get(f"{{{DOC_REL_NS}}}id")
                slide_path = slide_targets.get(rel_id or "")
                if not slide_path or slide_path not in archive.namelist():
                    continue
                slide_root = ET.fromstring(archive.read(slide_path))
                objects: List[Dict[str, object]] = []
                shape_tree = slide_root.find(f"{pptx_tag('cSld')}/{pptx_tag('spTree')}")
                raw_shapes = list(shape_tree) if shape_tree is not None else []
                object_index = 1
                for shape in raw_shapes:
                    if local_name(shape.tag) not in {"sp", "pic", "graphicFrame"}:
                        continue
                    obj = xml_shape_object(shape, f"{slide_stable_id}:{object_index}")
                    if obj is None:
                        continue
                    obj["z_index"] = object_index
                    objects.append(obj)
                    object_index += 1
                slides.append(
                    normalize_presentation_slide_payload(
                        {
                            "id": f"slide:{slide_stable_id}",
                            "objects": objects,
                            "notes": read_pptx_notes_xml(archive, slide_path),
                            "transition": xml_transition_payload(slide_root),
                            "timeline": xml_timeline_payload(slide_root),
                        },
                        index,
                        derive_summary=True,
                    )
                )

    return {
        "kind": "presentation",
        "format": "pptx",
        "path": path,
        **metadata,
        "warning": PRESENTATION_AUTOMATION_WARNING,
        "slides": slides,
    }


def pptx_placeholder_name(shape) -> str:
    if not getattr(shape, "is_placeholder", False):
        return ""
    try:
        placeholder_type = shape.placeholder_format.type
    except Exception:
        return ""
    name = getattr(placeholder_type, "name", None)
    return str(name or placeholder_type).upper()


def presentation_shape_paragraphs(shape) -> List[Dict[str, object]]:
    if not getattr(shape, "has_text_frame", False):
        return []
    paragraphs: List[Dict[str, object]] = []
    for paragraph in shape.text_frame.paragraphs:
        text = str(getattr(paragraph, "text", "") or "").strip()
        if not text:
            continue
        paragraphs.append(
            {
                "text": text,
                "level": max(int(getattr(paragraph, "level", 0) or 0), 0),
            }
        )
    return paragraphs


def presentation_table_cells(table) -> List[List[str]]:
    rows: List[List[str]] = []
    for row in getattr(table, "rows", []) or []:
        next_row: List[str] = []
        for cell in getattr(row, "cells", []) or []:
            next_row.append(str(getattr(cell, "text", "") or ""))
        if next_row:
            rows.append(next_row)
    return rows


def chart_series_payload(chart) -> List[Dict[str, object]]:
    payload: List[Dict[str, object]] = []
    try:
        series_iter = list(chart.series)
    except Exception:
        return payload
    for index, series in enumerate(series_iter, start=1):
        entry: Dict[str, object] = {"id": f"series:{index}"}
        name = str(getattr(series, "name", "") or "").strip()
        if name:
            entry["name"] = name
        try:
            values = list(getattr(series, "values"))
        except Exception:
            values = []
        if values:
            entry["values"] = values
        if len(entry) > 1:
            payload.append(entry)
    return payload


def chart_categories_payload(chart) -> List[str]:
    try:
        plots = list(chart.plots)
    except Exception:
        return []
    for plot in plots:
        try:
            categories = list(getattr(plot, "categories"))
        except Exception:
            categories = []
        if categories:
            return [str(value) for value in categories]
    return []


def presentation_shape_object(shape, z_index: Optional[int] = None) -> Optional[Dict[str, object]]:
    object_id = f"shape:{getattr(shape, 'shape_id', '') or shape.name}"
    placeholder_kind = pptx_placeholder_name(shape)
    placeholder_index = pptx_placeholder_index(shape)
    result: Dict[str, object] = {
        "id": object_id,
        "kind": "shape",
    }
    shape_type_name = str(getattr(getattr(shape, "shape_type", None), "name", "") or "").upper()
    if shape_type_name:
        result["shape_type"] = shape_type_name
    name = str(getattr(shape, "name", "") or "").strip()
    if name:
        result["name"] = name
    if placeholder_kind:
        result["placeholder_kind"] = placeholder_kind
    if placeholder_index is not None:
        result["placeholder_index"] = placeholder_index
    frame = emu_frame_dict(
        getattr(shape, "left", 0),
        getattr(shape, "top", 0),
        getattr(shape, "width", 0),
        getattr(shape, "height", 0),
    )
    result["frame"] = frame
    rotation = pptx_shape_rotation(shape)
    if rotation is not None:
        result["rotation"] = rotation
    if z_index is not None:
        result["z_index"] = z_index
    auto_shape = pptx_auto_shape_name(shape)
    if auto_shape:
        result["auto_shape"] = auto_shape
    style = pptx_style_payload(shape)
    if style:
        result["style"] = style

    if getattr(shape, "has_text_frame", False):
        paragraphs = presentation_shape_paragraphs(shape)
        body = presentation_body_from_paragraphs(paragraphs)
        if placeholder_kind in PPTX_TITLE_PLACEHOLDER_TYPES:
            result["kind"] = "title"
            result["text"] = presentation_text_value(body)
        elif auto_shape and placeholder_kind not in PPTX_BODY_PLACEHOLDER_TYPES:
            result["kind"] = "shape"
            result["body"] = body
        else:
            result["kind"] = "text_box"
            result["body"] = body
            text = presentation_text_value(body)
            if text:
                result["text"] = text
        return result

    if getattr(shape, "has_chart", False):
        result["kind"] = "chart"
        chart = getattr(shape, "chart", None)
        chart_type = getattr(chart, "chart_type", None)
        if chart_type is not None:
            result["chart_kind"] = str(getattr(chart_type, "name", chart_type))
        chart_title = ""
        if chart is not None:
            try:
                has_title = bool(getattr(chart, "has_title", False))
            except Exception:
                has_title = False
            if has_title:
                try:
                    chart_title = str(chart.chart_title.text_frame.text or "").strip()
                except Exception:
                    chart_title = ""
        if chart_title:
            result["title"] = chart_title
        categories = chart_categories_payload(chart) if chart is not None else []
        if categories:
            result["categories"] = categories
        series = chart_series_payload(chart) if chart is not None else []
        if series:
            result["series"] = series
        return result

    if getattr(shape, "has_table", False):
        result["kind"] = "table"
        table = getattr(shape, "table", None)
        if table is not None:
            result["table"] = {
                "rows": len(getattr(table, "rows", []) or []),
                "cols": len(getattr(table, "columns", []) or []),
            }
            cells = presentation_table_cells(table)
            if cells:
                result["table"]["cells"] = cells
        return result

    if shape_type_name == "PICTURE":
        result["kind"] = "image"
        try:
            filename = str(getattr(shape.image, "filename", "") or "").strip()
            if filename:
                result["image_name"] = filename
        except Exception:
            pass
        alt_text = str(getattr(shape, "alt_text", "") or "").strip()
        if alt_text:
            result["alt_text"] = alt_text
        crop = pptx_picture_crop_payload(shape)
        if crop:
            result["crop"] = crop
    return result


def slide_layout_ref(slide) -> Optional[str]:
    layout = getattr(slide, "slide_layout", None)
    if layout is None:
        return None
    return normalize_layout_ref(getattr(layout, "name", None))


def slide_master_ref(slide) -> Optional[str]:
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    return normalize_master_ref(getattr(master, "name", None))


def presentation_notes_from_slide(slide) -> List[str]:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return []
    notes: List[str] = []
    for shape in getattr(notes_slide, "shapes", []):
        if not getattr(shape, "has_text_frame", False):
            continue
        placeholder_kind = pptx_placeholder_name(shape)
        if placeholder_kind in PPTX_NOTES_IGNORED_PLACEHOLDER_TYPES:
            continue
        notes.extend(
            [str(item.get("text") or "") for item in presentation_shape_paragraphs(shape)]
        )
    return [note for note in notes if note.strip()]


def augment_presentation_document_from_xml(path: str, document: Dict[str, object]) -> Dict[str, object]:
    slides = document.get("slides")
    if not isinstance(slides, list) or not os.path.exists(path):
        return document
    sidecar_by_id: Dict[str, Dict[str, object]] = {}
    try:
        with zipfile.ZipFile(path, "r") as archive:
            presentation_xml = "ppt/presentation.xml"
            presentation_rels_xml = "ppt/_rels/presentation.xml.rels"
            if presentation_xml not in archive.namelist() or presentation_rels_xml not in archive.namelist():
                return document
            presentation_root = ET.fromstring(archive.read(presentation_xml))
            presentation_rels_root = ET.fromstring(archive.read(presentation_rels_xml))
            slide_targets: Dict[str, str] = {}
            for rel in presentation_rels_root:
                if local_name(rel.tag) != "Relationship":
                    continue
                rel_id = rel.attrib.get("Id")
                target = rel.attrib.get("Target")
                if rel_id and target:
                    slide_targets[rel_id] = normalize_relationship_target("ppt", target)
            slide_id_list = presentation_root.find(pptx_tag("sldIdLst"))
            if slide_id_list is None:
                return document
            for index, slide_ref in enumerate(slide_id_list.findall(pptx_tag("sldId")), start=1):
                slide_stable_id = str(slide_ref.attrib.get("id") or index)
                rel_id = slide_ref.attrib.get(f"{{{DOC_REL_NS}}}id")
                slide_path = slide_targets.get(rel_id or "")
                if not slide_path or slide_path not in archive.namelist():
                    continue
                slide_root = ET.fromstring(archive.read(slide_path))
                extra: Dict[str, object] = {}
                transition = xml_transition_payload(slide_root)
                if transition:
                    extra["transition"] = transition
                timeline = xml_timeline_payload(slide_root)
                if timeline:
                    extra["timeline"] = timeline
                if extra:
                    sidecar_by_id[f"slide:{slide_stable_id}"] = extra
    except Exception:
        return document

    if not sidecar_by_id:
        return document
    next_document = dict(document)
    next_slides = []
    for slide in slides:
        if not isinstance(slide, dict):
            next_slides.append(slide)
            continue
        merged = dict(slide)
        extra = sidecar_by_id.get(str(slide.get("id") or ""))
        if extra:
            merged.update(extra)
        next_slides.append(merged)
    next_document["slides"] = next_slides
    return next_document


def read_pptx_document_python(path: str, metadata: Dict[str, object]) -> Dict[str, object]:
    presentation = PptxPresentation(path)
    slides: List[Dict[str, object]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        objects = []
        for object_index, shape in enumerate(slide.shapes, start=1):
            obj = presentation_shape_object(shape, z_index=object_index)
            if obj is not None:
                objects.append(obj)
        slide_payload: Dict[str, object] = {
            "id": f"slide:{getattr(slide, 'slide_id', index)}",
            "objects": objects,
            "notes": presentation_notes_from_slide(slide),
        }
        layout_ref = slide_layout_ref(slide)
        if layout_ref:
            slide_payload["layout_ref"] = layout_ref
        master_ref = slide_master_ref(slide)
        if master_ref:
            slide_payload["master_ref"] = master_ref
        slides.append(normalize_presentation_slide_payload(slide_payload, index, derive_summary=True))
    return {
        "kind": "presentation",
        "format": "pptx",
        "path": path,
        **metadata,
        "warning": PRESENTATION_AUTOMATION_WARNING,
        "slides": slides,
    }


def read_pptx_document(path: str) -> Dict[str, object]:
    metadata = path_metadata(path)
    if not metadata["exists"]:
        return {
            "kind": "presentation",
            "format": "pptx",
            "path": path,
            **metadata,
            "warning": PRESENTATION_AUTOMATION_WARNING,
            "slides": [],
        }
    if python_pptx_available():
        try:
            return augment_presentation_document_from_xml(path, read_pptx_document_python(path, metadata))
        except Exception:
            pass
    return augment_presentation_document_from_xml(path, read_pptx_document_xml(path, metadata))


def choose_pptx_slide_layout(presentation, wants_body: bool, desired_layout_ref: Optional[str] = None):
    layouts = list(presentation.slide_layouts)
    if not layouts:
        raise RuntimeError("The presentation template does not provide any slide layouts.")
    preferred = None
    title_only = None
    desired = normalize_layout_ref(desired_layout_ref)
    for layout in layouts:
        layout_ref = normalize_layout_ref(getattr(layout, "name", None))
        placeholder_names = {pptx_placeholder_name(shape) for shape in layout.placeholders}
        has_title = bool(placeholder_names & {"TITLE", "CENTER_TITLE"})
        has_body = bool(placeholder_names & {"BODY", "OBJECT", "SUBTITLE"})
        if desired and layout_ref == desired:
            return layout
        if wants_body and has_title and has_body:
            return layout
        if has_title and title_only is None:
            title_only = layout
        if preferred is None:
            preferred = layout
    if wants_body:
        return title_only or preferred
    return title_only or preferred


def remove_pptx_slide(presentation, index: int) -> None:
    slide_id_list = presentation.slides._sldIdLst  # type: ignore[attr-defined]
    slide_id = slide_id_list[index]
    rel_id = slide_id.rId
    presentation.part.drop_rel(rel_id)
    del slide_id_list[index]


def clear_text_frame(text_frame) -> None:
    text_frame.clear()


def populate_text_frame(text_frame, paragraphs: List[object]) -> None:
    clear_text_frame(text_frame)
    normalized = body_items_from_outline(paragraphs)
    if not normalized:
        return
    first = text_frame.paragraphs[0]
    first.text = str(normalized[0].get("text") or "")
    first.level = int(normalized[0].get("level") or 0)
    for item in normalized[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = str(item.get("text") or "")
        paragraph.level = int(item.get("level") or 0)


def find_title_shape(slide):
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        return title_shape
    for shape in slide.shapes:
        if pptx_placeholder_name(shape) in {"TITLE", "CENTER_TITLE"} and getattr(shape, "has_text_frame", False):
            return shape
    return None


def find_body_shape(slide):
    for shape in slide.shapes:
        if pptx_placeholder_name(shape) in {"BODY", "OBJECT", "SUBTITLE"} and getattr(shape, "has_text_frame", False):
            return shape
    return None


def add_textbox_lines(slide, left, top, width, height, lines: List[str]):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    populate_text_frame(
        textbox.text_frame,
        [{"text": line, "level": 0} for line in lines if str(line or "").strip()],
    )
    return textbox


def remove_shape(shape) -> None:
    element = getattr(shape, "_element", None)
    if element is None:
        return
    parent = element.getparent()
    if parent is None:
        return
    parent.remove(element)


def reorder_slide_shapes(slide, shapes: List[object]) -> None:
    sp_tree = getattr(slide.shapes, "_spTree", None)
    if sp_tree is None:
        return
    for shape in shapes:
        element = getattr(shape, "_element", None)
        if element is None:
            continue
        try:
            sp_tree.remove(element)
        except Exception:
            pass
        sp_tree.append(element)


def apply_frame_to_shape(shape, frame: object) -> None:
    normalized = normalize_presentation_frame(frame)
    if not normalized:
        return
    for source, attr in (("x", "left"), ("y", "top"), ("w", "width"), ("h", "height")):
        try:
            setattr(shape, attr, int(normalized.get(source) or 0))
        except Exception:
            continue


def rgb_color_from_hex(color: Optional[str]):
    normalized = normalize_hex_color(color)
    if not normalized or RGBColor is None:
        return None
    value = normalized[1:]
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def apply_color_format_payload(color_format, payload: object) -> None:
    if not isinstance(payload, dict) or color_format is None:
        return
    color = rgb_color_from_hex(str(payload.get("color") or ""))
    if color is None:
        return
    try:
        color_format.rgb = color
    except Exception:
        return


def apply_shape_style(shape, raw_style: object) -> None:
    style = normalize_presentation_style(raw_style)
    if not style:
        return
    fill_payload = style.get("fill")
    if isinstance(fill_payload, dict):
        try:
            fill = shape.fill
            fill_type = str(fill_payload.get("type") or "").strip().lower()
            if fill_type in {"none", "background"}:
                fill.background()
            else:
                fill.solid()
                apply_color_format_payload(getattr(fill, "fore_color", None), fill_payload)
        except Exception:
            pass
    stroke_payload = style.get("stroke")
    if isinstance(stroke_payload, dict):
        try:
            line = shape.line
            apply_color_format_payload(getattr(line, "color", None), stroke_payload)
            if stroke_payload.get("width") is not None:
                line.width = int(stroke_payload.get("width") or 0)
        except Exception:
            pass


def apply_picture_crop(shape, raw_crop: object) -> None:
    crop = normalize_presentation_crop(raw_crop)
    if not crop:
        return
    for source, attr in (
        ("left", "crop_left"),
        ("top", "crop_top"),
        ("right", "crop_right"),
        ("bottom", "crop_bottom"),
    ):
        if source not in crop:
            continue
        try:
            setattr(shape, attr, float(crop[source]))
        except Exception:
            continue


def populate_chart_shape(shape, payload: Dict[str, object]) -> None:
    chart = getattr(shape, "chart", None)
    if chart is None:
        return
    title = str(payload.get("title") or "").strip()
    if title:
        try:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        except Exception:
            pass
    categories = payload.get("categories")
    series = normalize_chart_series(payload.get("series"))
    if not isinstance(categories, list) or not series or CategoryChartData is None:
        return
    try:
        chart_data = CategoryChartData()
        chart_data.categories = [str(value or "") for value in categories]
        for index, series_entry in enumerate(series, start=1):
            values = series_entry.get("values")
            if not isinstance(values, list):
                continue
            chart_data.add_series(str(series_entry.get("name") or f"Series {index}"), values)
        chart.replace_data(chart_data)
    except Exception:
        return


def apply_shape_object_properties(shape, payload: Dict[str, object]) -> None:
    if payload.get("rotation") is not None:
        try:
            shape.rotation = float(payload.get("rotation") or 0)
        except Exception:
            pass
    if payload.get("name"):
        try:
            shape.name = str(payload.get("name") or "")
        except Exception:
            pass
    alt_text = str(payload.get("alt_text") or "").strip()
    if alt_text:
        try:
            shape.alt_text = alt_text
        except Exception:
            pass
    apply_shape_style(shape, payload.get("style"))
    apply_picture_crop(shape, payload.get("crop"))
    populate_chart_shape(shape, payload)


def current_shape_kind(shape) -> str:
    if getattr(shape, "has_chart", False):
        return "chart"
    if getattr(shape, "has_table", False):
        return "table"
    shape_type_name = str(getattr(getattr(shape, "shape_type", None), "name", "") or "").upper()
    if shape_type_name == "PICTURE":
        return "image"
    if getattr(shape, "has_text_frame", False):
        if pptx_placeholder_name(shape) in PPTX_TITLE_PLACEHOLDER_TYPES:
            return "title"
        if pptx_placeholder_name(shape) in PPTX_BODY_PLACEHOLDER_TYPES:
            return "text_box"
        return "shape"
    return "shape"


def shape_requires_replacement(shape, payload: Dict[str, object]) -> bool:
    desired_kind = str(payload.get("kind") or "").strip()
    current_kind = current_shape_kind(shape)
    if desired_kind == "image":
        return current_kind != "image" or bool(payload.get("image_ref"))
    if desired_kind == "table":
        return current_kind != "table"
    if desired_kind == "chart":
        return current_kind != "chart"
    desired_auto_shape = str(payload.get("auto_shape") or "").strip().upper()
    if desired_auto_shape and current_kind == "shape":
        current_auto_shape = str(pptx_auto_shape_name(shape) or "").upper()
        return current_auto_shape != desired_auto_shape
    return False


def shape_matches_payload(shape, payload: Dict[str, object], used_ids: set[str]):
    shape_id = f"shape:{getattr(shape, 'shape_id', '') or ''}"
    if shape_id in used_ids:
        return False
    placeholder_kind = pptx_placeholder_name(shape)
    placeholder_index = pptx_placeholder_index(shape)
    desired_placeholder = str(payload.get("placeholder_kind") or "").upper()
    desired_placeholder_index = payload.get("placeholder_index")
    desired_kind = str(payload.get("kind") or "")
    if desired_placeholder_index is not None and placeholder_index == desired_placeholder_index:
        return True
    if desired_placeholder and placeholder_kind == desired_placeholder:
        return True
    if desired_kind == "title" and placeholder_kind in PPTX_TITLE_PLACEHOLDER_TYPES:
        return True
    return False


def safe_resolve_workspace_asset(raw_path: object) -> Optional[str]:
    text = str(raw_path or "").strip()
    if not text:
        return None
    try:
        resolved = resolve_workspace_path(text)
    except Exception:
        return None
    return resolved if os.path.exists(resolved) else None


def populate_table_shape(table, cells: List[List[str]]) -> None:
    for row_index, row in enumerate(cells):
        if row_index >= len(getattr(table, "rows", []) or []):
            break
        for col_index, value in enumerate(row):
            if col_index >= len(getattr(table, "columns", []) or []):
                break
            try:
                table.cell(row_index, col_index).text = str(value or "")
            except Exception:
                continue


def create_shape_from_payload(slide, payload: Dict[str, object]):
    kind = str(payload.get("kind") or "")
    frame = normalize_presentation_frame(payload.get("frame")) or emu_frame_dict(
        Inches(0.9),
        Inches(1.6),
        Inches(8.0),
        Inches(1.5),
    )
    if kind == "image":
        image_path = safe_resolve_workspace_asset(payload.get("image_ref"))
        if image_path:
            kwargs = {}
            if int(frame.get("w") or 0) > 0:
                kwargs["width"] = int(frame["w"])
            if int(frame.get("h") or 0) > 0:
                kwargs["height"] = int(frame["h"])
            return slide.shapes.add_picture(
                image_path,
                int(frame.get("x") or 0),
                int(frame.get("y") or 0),
                **kwargs,
            )
        return None
    if kind == "table":
        table_payload = payload.get("table") if isinstance(payload.get("table"), dict) else {}
        rows = max(int(table_payload.get("rows") or 0), 1)
        cols = max(int(table_payload.get("cols") or 0), 1)
        shape = slide.shapes.add_table(
            rows,
            cols,
            int(frame.get("x") or 0),
            int(frame.get("y") or 0),
            int(frame.get("w") or 0),
            int(frame.get("h") or 0),
        )
        cells = normalize_string_matrix(table_payload.get("cells"))
        if cells:
            populate_table_shape(shape.table, cells)
        return shape
    if kind == "shape":
        auto_shape_name = str(payload.get("auto_shape") or "").strip().upper()
        auto_shape_enum = getattr(MSO_AUTO_SHAPE_TYPE, auto_shape_name, None) if MSO_AUTO_SHAPE_TYPE is not None and auto_shape_name else None
        if auto_shape_enum is not None:
            shape = slide.shapes.add_shape(
                auto_shape_enum,
                int(frame.get("x") or 0),
                int(frame.get("y") or 0),
                int(frame.get("w") or 0),
                int(frame.get("h") or 0),
            )
            if getattr(shape, "has_text_frame", False):
                text_body = payload.get("body")
                if "text" in payload and "body" not in payload:
                    text_body = {"kind": "text", "text": str(payload.get("text") or "")}
                if text_body is not None:
                    populate_text_frame(shape.text_frame, presentation_paragraphs_from_body(text_body))
            return shape
    if kind in {"title", "text_box", "shape"} or "text" in payload or "body" in payload:
        textbox = slide.shapes.add_textbox(
            int(frame.get("x") or 0),
            int(frame.get("y") or 0),
            int(frame.get("w") or 0),
            int(frame.get("h") or 0),
        )
        text_body = payload.get("body")
        if "text" in payload and "body" not in payload:
            text_body = {"kind": "text", "text": str(payload.get("text") or "")}
        populate_text_frame(textbox.text_frame, presentation_paragraphs_from_body(text_body))
        return textbox
    return None


def update_notes_slide(slide, notes: object) -> None:
    normalized_notes = normalize_presentation_notes(notes)
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return
    for shape in notes_slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        placeholder_kind = pptx_placeholder_name(shape)
        if placeholder_kind in PPTX_NOTES_IGNORED_PLACEHOLDER_TYPES:
            continue
        populate_text_frame(
            shape.text_frame,
            [{"text": line, "level": 0} for line in normalized_notes],
        )
        return


def replace_shape_with_payload(slide, shape, payload: Dict[str, object]):
    fallback_frame = emu_frame_dict(
        getattr(shape, "left", 0),
        getattr(shape, "top", 0),
        getattr(shape, "width", 0),
        getattr(shape, "height", 0),
    )
    next_payload = dict(payload)
    if "frame" not in next_payload:
        next_payload["frame"] = fallback_frame
    remove_shape(shape)
    return create_shape_from_payload(slide, next_payload)


def update_slide_from_payload(slide, slide_payload: Dict[str, object]) -> None:
    existing_shapes = {
        f"shape:{getattr(shape, 'shape_id', '') or shape.name}": shape for shape in slide.shapes
    }
    used_ids: set[str] = set()
    touched_title = False
    touched_body = False
    ordered_shapes: List[Tuple[int, object]] = []
    touched_z_order = False
    explicit_remove_ids = {
        str(value).strip()
        for value in (slide_payload.get("remove_object_ids") or [])
        if str(value or "").strip()
    }

    for obj in slide_payload.get("objects") or []:
        object_id = str(obj.get("id") or "")
        shape = existing_shapes.get(object_id)
        if shape is None:
            for candidate in slide.shapes:
                if shape_matches_payload(candidate, obj, used_ids):
                    shape = candidate
                    break
        if shape is None:
            shape = create_shape_from_payload(slide, obj)
        elif shape_requires_replacement(shape, obj):
            shape = replace_shape_with_payload(slide, shape, obj)
        if shape is None:
            continue
        used_ids.add(f"shape:{getattr(shape, 'shape_id', '') or shape.name}")
        apply_frame_to_shape(shape, obj.get("frame"))
        apply_shape_object_properties(shape, obj)
        if getattr(shape, "has_table", False):
            table_payload = obj.get("table") if isinstance(obj.get("table"), dict) else {}
            cells = normalize_string_matrix(table_payload.get("cells"))
            if cells:
                populate_table_shape(shape.table, cells)
        if getattr(shape, "has_text_frame", False):
            body = obj.get("body")
            if "text" in obj and "body" not in obj:
                body = {"kind": "text", "text": str(obj.get("text") or "")}
            if body is not None:
                populate_text_frame(shape.text_frame, presentation_paragraphs_from_body(body))
        if obj.get("z_index") is not None:
            touched_z_order = True
        try:
            z_index = int(obj.get("z_index") or len(ordered_shapes) + 1)
        except (TypeError, ValueError):
            z_index = len(ordered_shapes) + 1
        ordered_shapes.append((z_index, shape))
        placeholder_kind = str(obj.get("placeholder_kind") or "").upper()
        kind = str(obj.get("kind") or "")
        if kind == "title" or placeholder_kind in PPTX_TITLE_PLACEHOLDER_TYPES:
            touched_title = True
        if placeholder_kind in PPTX_BODY_PLACEHOLDER_TYPES:
            touched_body = True

    for object_id in explicit_remove_ids:
        shape = existing_shapes.get(object_id)
        if shape is not None:
            remove_shape(shape)

    if slide_payload.get("prune_missing_objects") is True:
        for object_id, shape in existing_shapes.items():
            if object_id not in used_ids and object_id not in explicit_remove_ids:
                remove_shape(shape)

    if "title" in slide_payload and not touched_title:
        title_shape = find_title_shape(slide)
        if title_shape is not None:
            populate_text_frame(
                title_shape.text_frame,
                presentation_paragraphs_from_body({"kind": "text", "text": str(slide_payload.get("title") or "")}),
            )
        elif str(slide_payload.get("title") or "").strip():
            add_textbox_lines(
                slide,
                Inches(0.75),
                Inches(0.5),
                Inches(8.5),
                Inches(0.9),
                [str(slide_payload.get("title") or "")],
            )

    if "body" in slide_payload and not touched_body:
        body_shape = find_body_shape(slide)
        paragraphs = presentation_paragraphs_from_body(slide_payload.get("body"))
        if body_shape is not None:
            populate_text_frame(body_shape.text_frame, paragraphs)
        elif paragraphs:
            add_textbox_lines(
                slide,
                Inches(0.9),
                Inches(1.6),
                Inches(8.0),
                Inches(4.5),
                [str(item.get("text") or "") for item in paragraphs],
            )

    if "notes" in slide_payload:
        update_notes_slide(slide, slide_payload.get("notes"))

    if ordered_shapes and (slide_payload.get("reorder_objects") is True or touched_z_order):
        ordered_shapes.sort(key=lambda item: item[0])
        reorder_slide_shapes(slide, [shape for _, shape in ordered_shapes])


def write_pptx_document(path: str, slides: List[Dict[str, object]], expected_etag: Optional[str]) -> Dict[str, object]:
    assert_expected_etag(path, expected_etag)
    normalized_slides = [
        normalize_presentation_slide_payload(slide, index, derive_summary=False)
        for index, slide in enumerate(slides, start=1)
    ]
    if not normalized_slides:
        raise RuntimeError("Presentation AIO must contain at least one slide.")
    if not python_pptx_available():
        raise RuntimeError(
            "Structured .pptx automation requires python-pptx. Rebuild the OpenClaw runtime image."
        )

    presentation = PptxPresentation(path) if os.path.exists(path) else PptxPresentation()
    existing_slides = {
        f"slide:{getattr(slide, 'slide_id', index)}": slide
        for index, slide in enumerate(presentation.slides, start=1)
    }

    for slide_payload in normalized_slides:
        slide_id = str(slide_payload.get("id") or "")
        ppt_slide = existing_slides.get(slide_id)
        if ppt_slide is None:
            body = slide_payload.get("body") or {"kind": "text", "text": ""}
            ppt_slide = presentation.slides.add_slide(
                choose_pptx_slide_layout(
                    presentation,
                    bool(presentation_paragraphs_from_body(body)),
                    slide_payload.get("layout_ref"),
                )
            )
        update_slide_from_payload(ppt_slide, slide_payload)

    temp_fd, temp_path = tempfile.mkstemp(prefix=".entropic-office-", dir=posixpath.dirname(path))
    os.close(temp_fd)
    try:
        presentation.save(temp_path)
        os.replace(temp_path, path)
    finally:
        try:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
        except OSError:
            pass
    return read_pptx_document(path)


def spreadsheet_analysis_payload(document: Dict[str, object]) -> Dict[str, object]:
    sheets = [sheet for sheet in document.get("sheets") or [] if isinstance(sheet, dict)]
    cells = [
        cell
        for sheet in sheets
        for cell in sheet.get("cells") or []
        if isinstance(cell, dict)
    ]
    return {
        "source_backend": default_backend_name("spreadsheet", str(document.get("format") or "")),
        "stats": {
            "worksheet_count": len(sheets),
            "non_empty_cell_count": len(cells),
            "formula_cell_count": sum(1 for cell in cells if str(cell.get("kind") or "") == "formula" or cell.get("formula")),
            "merged_range_count": sum(len([entry for entry in sheet.get("mergedRanges") or [] if isinstance(entry, dict)]) for sheet in sheets),
            "freeze_pane_count": sum(1 for sheet in sheets if isinstance(sheet.get("freezePane"), dict)),
        },
        "worksheets": [
            {
                "name": str(sheet.get("name") or ""),
                "non_empty_cells": len([cell for cell in sheet.get("cells") or [] if isinstance(cell, dict)]),
                "formula_cells": len(
                    [
                        cell
                        for cell in sheet.get("cells") or []
                        if isinstance(cell, dict)
                        and (str(cell.get("kind") or "") == "formula" or cell.get("formula"))
                    ]
                ),
                "merged_ranges": len([entry for entry in sheet.get("mergedRanges") or [] if isinstance(entry, dict)]),
                "has_freeze_pane": isinstance(sheet.get("freezePane"), dict),
            }
            for sheet in sheets
        ],
    }


def spreadsheet_projection_payload(document: Dict[str, object]) -> Dict[str, object]:
    sheets = [sheet for sheet in document.get("sheets") or [] if isinstance(sheet, dict)]
    return {
        "summary-view": {
            "kind": "worksheet-summary",
            "worksheets": [
                {
                    "name": str(sheet.get("name") or ""),
                    "extent": {
                        "rows": int(sheet.get("rowCount") or 0),
                        "cols": int(sheet.get("colCount") or 0),
                    },
                    "freeze_pane": sheet.get("freezePane"),
                    "merged_ranges": [
                        str(entry.get("ref") or "")
                        for entry in sheet.get("mergedRanges") or []
                        if isinstance(entry, dict)
                    ],
                }
                for sheet in sheets
            ],
        },
        "formula-view": {
            "kind": "formula-summary",
            "cells": [
                {
                    "sheet": str(sheet.get("name") or ""),
                    "ref": str(cell.get("ref") or ""),
                    "formula": str(cell.get("formula") or ""),
                }
                for sheet in sheets
                for cell in sheet.get("cells") or []
                if isinstance(cell, dict) and cell.get("formula")
            ],
        },
        "layout-view": {
            "kind": "sheet-layout",
            "worksheets": [
                {
                    "name": str(sheet.get("name") or ""),
                    "rows": [entry for entry in sheet.get("rows") or [] if isinstance(entry, dict)],
                    "columns": [entry for entry in sheet.get("columns") or [] if isinstance(entry, dict)],
                    "filter_ref": str(sheet.get("filterRef") or ""),
                }
                for sheet in sheets
            ],
        },
    }


def spreadsheet_artifact_payload(document: Dict[str, object]) -> Dict[str, object]:
    format_name = str(document.get("format") or "").strip()
    return {
        "native": {
            "backend": default_backend_name("spreadsheet", format_name),
            "format": format_name,
            "role": "grid-metadata",
        }
    }


def presentation_analysis_payload(document: Dict[str, object]) -> Dict[str, object]:
    slides = [slide for slide in document.get("slides") or [] if isinstance(slide, dict)]
    objects = [
        obj
        for slide in slides
        for obj in slide.get("objects") or []
        if isinstance(obj, dict)
    ]
    return {
        "source_backend": default_backend_name("presentation", str(document.get("format") or "")),
        "stats": {
            "slide_count": len(slides),
            "object_count": len(objects),
            "notes_count": sum(len(normalize_presentation_notes(slide.get("notes"))) for slide in slides),
            "image_count": sum(1 for obj in objects if str(obj.get("kind") or "") == "image"),
            "styled_object_count": sum(1 for obj in objects if isinstance(obj.get("style"), dict)),
            "cropped_image_count": sum(1 for obj in objects if str(obj.get("kind") or "") == "image" and isinstance(obj.get("crop"), dict)),
        },
        "object_kinds": {
            "title": sum(1 for obj in objects if str(obj.get("kind") or "") == "title"),
            "text_box": sum(1 for obj in objects if str(obj.get("kind") or "") == "text_box"),
            "table": sum(1 for obj in objects if str(obj.get("kind") or "") == "table"),
            "chart": sum(1 for obj in objects if str(obj.get("kind") or "") == "chart"),
            "image": sum(1 for obj in objects if str(obj.get("kind") or "") == "image"),
        },
    }


def presentation_projection_payload(document: Dict[str, object]) -> Dict[str, object]:
    slides = [slide for slide in document.get("slides") or [] if isinstance(slide, dict)]
    return {
        "outline-view": {
            "kind": "slide-outline",
            "slides": [
                {
                    "id": str(slide.get("id") or ""),
                    "index": int(slide.get("index") or slide_index),
                    "title": str(slide.get("title") or ""),
                    "notes_preview": normalize_presentation_notes(slide.get("notes"))[:2],
                }
                for slide_index, slide in enumerate(slides, start=1)
            ],
        },
        "scene-view": {
            "kind": "scene-summary",
            "slides": [
                {
                    "id": str(slide.get("id") or ""),
                    "object_ids": [
                        str(obj.get("id") or "")
                        for obj in slide.get("objects") or []
                        if isinstance(obj, dict)
                    ],
                    "z_order": [
                        {
                            "id": str(obj.get("id") or ""),
                            "z_index": int(obj.get("z_index") or object_index),
                        }
                        for object_index, obj in enumerate(
                            [obj for obj in slide.get("objects") or [] if isinstance(obj, dict)],
                            start=1,
                        )
                    ],
                }
                for slide in slides
            ],
        },
    }


def presentation_artifact_payload(document: Dict[str, object]) -> Dict[str, object]:
    format_name = str(document.get("format") or "").strip()
    return {
        "native": {
            "backend": default_backend_name("presentation", format_name),
            "format": format_name,
            "role": "scene-metadata",
            "supports_raw_style_ops": python_pptx_available(),
        }
    }


def spreadsheet_document_to_aio(document: Dict[str, object]) -> Dict[str, object]:
    format_name = str(document.get("format") or "")
    backend_name = str(document.get("backend_name") or default_backend_name("spreadsheet", format_name))
    workbook_id = f"workbook:{slugify_identifier(basename_without_extension(str(document.get('path') or 'workbook')), 'workbook')}"
    worksheets = []
    for index, sheet in enumerate(document.get("sheets") or [], start=1):
        if not isinstance(sheet, dict):
            continue
        sheet_name = str(sheet.get("name") or f"Sheet{index}").strip() or f"Sheet{index}"
        worksheet_slug = slugify_identifier(sheet_name, f"sheet{index}")
        cells = []
        for raw_cell in sheet.get("cells") or []:
            if not isinstance(raw_cell, dict):
                continue
            ref = str(raw_cell.get("ref") or "").strip().upper()
            row = int(raw_cell.get("row") or 0)
            col = int(raw_cell.get("col") or 0)
            if not ref:
                if row <= 0 or col <= 0:
                    continue
                ref = cell_ref(row, col)
            elif row <= 0 or col <= 0:
                row, col = parse_cell_ref(ref)
            formula = str(raw_cell.get("formula") or "").strip()
            value = str(raw_cell.get("value") or "")
            display = str(raw_cell.get("display") or value)
            kind = str(raw_cell.get("kind") or ("formula" if formula else "string"))
            entry = {
                "id": f"cell:{worksheet_slug}:{ref}",
                "kind": "cell",
                "ref": ref,
                "row": row,
                "col": col,
            }
            if formula:
                entry["formula"] = formula
                if value:
                    entry["value"] = value
                if display and display != value:
                    entry["display"] = display
            else:
                entry["value"] = value
                entry["value_kind"] = kind
                if display and display != value:
                    entry["display"] = display
            cells.append(entry)
        worksheets.append(
            {
                "id": f"worksheet:{worksheet_slug}",
                "kind": "worksheet",
                "name": sheet_name,
                "extent": {
                    "rows": int(sheet.get("rowCount") or 0),
                    "cols": int(sheet.get("colCount") or 0),
                },
                "cells": cells,
            }
        )
        worksheet_entry = worksheets[-1]
        rows = normalize_sheet_dimension_rows(sheet.get("rows"))
        if rows:
            worksheet_entry["rows"] = rows
        columns = normalize_sheet_dimension_columns(sheet.get("columns"))
        if columns:
            worksheet_entry["columns"] = columns
        merged_ranges = normalize_sheet_merged_ranges(sheet.get("mergedRanges"))
        if merged_ranges:
            worksheet_entry["merged_ranges"] = [
                {
                    "id": str(entry.get("id") or f"range:{worksheet_slug}:{entry.get('ref')}"),
                    "kind": "merged_range",
                    "ref": str(entry.get("ref") or ""),
                }
                for entry in merged_ranges
            ]
        freeze_pane = normalize_sheet_freeze_pane(sheet.get("freezePane"))
        if freeze_pane:
            worksheet_entry["freeze_pane"] = freeze_pane
        filter_ref = str(sheet.get("filterRef") or "").strip().upper()
        if filter_ref:
            worksheet_entry["filter_ref"] = filter_ref
    return aio_envelope(
        "spreadsheet",
        document,
        {
            "id": workbook_id,
            "kind": "workbook",
            "worksheets": worksheets,
        },
        capability_set="current_entropic_office",
        analysis=spreadsheet_analysis_payload(document),
        projections=spreadsheet_projection_payload(document),
        artifacts=spreadsheet_artifact_payload(document),
        backend_name=backend_name,
    )


def document_document_to_aio(document: Dict[str, object]) -> Dict[str, object]:
    format_name = str(document.get("format") or "")
    backend_name = str(document.get("backend_name") or default_backend_name("document", format_name))
    document_id = f"document:{slugify_identifier(basename_without_extension(str(document.get('path') or 'document')), 'document')}"
    blocks = [
        block
        for block in document.get("blocks") or []
        if isinstance(block, dict)
    ]
    if not blocks:
        for index, paragraph in enumerate(document.get("paragraphs") or [], start=1):
            blocks.append(
                {
                    "id": f"block:{index}",
                    "kind": "paragraph",
                    "index": index,
                    "text": str(paragraph),
                }
            )
    return aio_envelope(
        "document",
        document,
        {
            "id": document_id,
            "kind": "document",
            "blocks": blocks,
        },
        capability_set="current_entropic_office",
        analysis=merge_nested_dict(
            document_analysis_payload(blocks, backend_name),
            document.get("analysis") if isinstance(document.get("analysis"), dict) else None,
        ),
        annotations=[item for item in document.get("annotations") or [] if isinstance(item, dict)] or collect_document_annotations(blocks),
        projections=merge_nested_dict(
            document_projection_payload(blocks),
            document.get("projections") if isinstance(document.get("projections"), dict) else None,
        ),
        artifacts=merge_nested_dict(
            document_artifact_payload(format_name, backend_name, 0),
            document.get("artifacts") if isinstance(document.get("artifacts"), dict) else None,
        ),
        backend_name=backend_name,
    )


def presentation_document_to_aio(document: Dict[str, object]) -> Dict[str, object]:
    format_name = str(document.get("format") or "")
    backend_name = str(document.get("backend_name") or default_backend_name("presentation", format_name))
    deck_id = f"deck:{slugify_identifier(basename_without_extension(str(document.get('path') or 'deck')), 'deck')}"
    slides = []
    for index, slide in enumerate(document.get("slides") or [], start=1):
        normalized = normalize_presentation_slide_payload(slide, index, derive_summary=True)
        slide_entry: Dict[str, object] = {
            "id": str(normalized.get("id") or f"slide:{index}"),
            "kind": "slide",
            "index": index,
        }
        if normalized.get("layout_ref"):
            slide_entry["layout_ref"] = str(normalized.get("layout_ref"))
        if normalized.get("master_ref"):
            slide_entry["master_ref"] = str(normalized.get("master_ref"))
        title = str(normalized.get("title") or "").strip()
        if title:
            slide_entry["title"] = title
        body = normalized.get("body")
        if body is not None and presentation_paragraphs_from_body(body):
            slide_entry["body"] = body
        if normalized.get("objects"):
            slide_entry["objects"] = normalized.get("objects")
        if normalized.get("transition"):
            slide_entry["transition"] = normalized.get("transition")
        if normalized.get("timeline"):
            slide_entry["timeline"] = normalized.get("timeline")
        notes = normalize_presentation_notes(normalized.get("notes"))
        if notes:
            slide_entry["notes"] = notes
        slides.append(slide_entry)
    return aio_envelope(
        "presentation",
        document,
        {
            "id": deck_id,
            "kind": "deck",
            "slides": slides,
        },
        capability_set="current_entropic_office",
        analysis=presentation_analysis_payload(document),
        projections=presentation_projection_payload(document),
        artifacts=presentation_artifact_payload(document),
        backend_name=backend_name,
    )


def document_to_aio(document: Dict[str, object]) -> Dict[str, object]:
    kind = str(document.get("kind") or "").strip()
    if kind == "spreadsheet":
        return spreadsheet_document_to_aio(document)
    if kind == "document":
        return document_document_to_aio(document)
    if kind == "presentation":
        return presentation_document_to_aio(document)
    raise RuntimeError(f"Unsupported AIO conversion for `{kind or 'unknown'}`.")


def attach_aio_projection(document: Dict[str, object]) -> Dict[str, object]:
    result = dict(document)
    result["aio"] = document_to_aio(document)
    return result


def spreadsheet_payload_from_aio(payload: Dict[str, object]) -> Dict[str, object]:
    ensure_aio_payload(payload)
    if str(payload.get("kind") or "").strip() != "spreadsheet":
        raise RuntimeError("Expected a spreadsheet AIO object.")
    object_payload = payload.get("object")
    if not isinstance(object_payload, dict):
        raise RuntimeError("AIO spreadsheet object is missing.")
    worksheets = object_payload.get("worksheets")
    if not isinstance(worksheets, list):
        raise RuntimeError("AIO spreadsheet object must contain `worksheets`.")
    sheets = []
    for index, worksheet in enumerate(worksheets, start=1):
        if not isinstance(worksheet, dict):
            continue
        name = str(worksheet.get("name") or f"Sheet{index}").strip() or f"Sheet{index}"
        cells = []
        for raw_cell in worksheet.get("cells") or []:
            if not isinstance(raw_cell, dict):
                continue
            ref = str(raw_cell.get("ref") or "").strip().upper()
            row = int(raw_cell.get("row") or 0)
            col = int(raw_cell.get("col") or 0)
            if not ref:
                if row <= 0 or col <= 0:
                    continue
                ref = cell_ref(row, col)
            elif row <= 0 or col <= 0:
                row, col = parse_cell_ref(ref)
            formula = str(raw_cell.get("formula") or "").strip()
            value = str(raw_cell.get("value") or "")
            display = str(raw_cell.get("display") or value)
            value_kind = str(raw_cell.get("value_kind") or raw_cell.get("kind") or "").strip()
            cell_entry = {
                "ref": ref,
                "row": row,
                "col": col,
                "display": display,
            }
            if formula:
                cell_entry["formula"] = formula
                cell_entry["value"] = value
            else:
                cell_entry["value"] = value
                if value_kind:
                    cell_entry["kind"] = value_kind
            cells.append(cell_entry)
        sheet_entry: Dict[str, object] = {"name": name, "cells": cells}
        extent = worksheet.get("extent")
        if isinstance(extent, dict):
            if extent.get("rows") is not None:
                sheet_entry["rowCount"] = int(extent.get("rows") or 0)
            if extent.get("cols") is not None:
                sheet_entry["colCount"] = int(extent.get("cols") or 0)
        rows = normalize_sheet_dimension_rows(worksheet.get("rows"))
        if "rows" in worksheet or rows:
            sheet_entry["rows"] = rows
        columns = normalize_sheet_dimension_columns(worksheet.get("columns"))
        if "columns" in worksheet or columns:
            sheet_entry["columns"] = columns
        merged_ranges = normalize_sheet_merged_ranges(
            worksheet.get("merged_ranges") if "merged_ranges" in worksheet else worksheet.get("mergedRanges")
        )
        if "merged_ranges" in worksheet or "mergedRanges" in worksheet or merged_ranges:
            sheet_entry["mergedRanges"] = merged_ranges
        freeze_pane = normalize_sheet_freeze_pane(
            worksheet.get("freeze_pane") if "freeze_pane" in worksheet else worksheet.get("freezePane")
        )
        if "freeze_pane" in worksheet or "freezePane" in worksheet or freeze_pane:
            sheet_entry["freezePane"] = freeze_pane
        filter_ref = str(worksheet.get("filter_ref") or worksheet.get("filterRef") or "").strip().upper() or None
        if "filter_ref" in worksheet or "filterRef" in worksheet or filter_ref:
            sheet_entry["filterRef"] = filter_ref
        sheets.append(sheet_entry)
    return {
        "expectedEtag": extract_aio_source_etag(payload),
        "sheets": sheets,
    }


def document_payload_from_aio(payload: Dict[str, object]) -> Dict[str, object]:
    ensure_aio_payload(payload)
    if str(payload.get("kind") or "").strip() != "document":
        raise RuntimeError("Expected a document AIO object.")
    object_payload = payload.get("object")
    if not isinstance(object_payload, dict):
        raise RuntimeError("AIO document object is missing.")
    blocks = object_payload.get("blocks")
    if not isinstance(blocks, list):
        raise RuntimeError("AIO document object must contain `blocks`.")
    normalized_blocks = [block for block in blocks if isinstance(block, dict)]
    paragraphs = document_blocks_to_paragraphs(normalized_blocks)
    return {
        "expectedEtag": extract_aio_source_etag(payload),
        "blocks": normalized_blocks,
        "paragraphs": paragraphs,
    }


def presentation_payload_from_aio(payload: Dict[str, object]) -> Dict[str, object]:
    ensure_aio_payload(payload)
    if str(payload.get("kind") or "").strip() != "presentation":
        raise RuntimeError("Expected a presentation AIO object.")
    object_payload = payload.get("object")
    if not isinstance(object_payload, dict):
        raise RuntimeError("AIO presentation object is missing.")
    slides = object_payload.get("slides")
    if not isinstance(slides, list):
        raise RuntimeError("AIO presentation object must contain `slides`.")
    return {
        "expectedEtag": extract_aio_source_etag(payload),
        "slides": [
            normalize_presentation_slide_payload(slide, index, derive_summary=False)
            for index, slide in enumerate(slides, start=1)
        ],
    }


def inspect_aio(path: str) -> Dict[str, object]:
    extension = split_extension(path)
    if extension == ".csv":
        return spreadsheet_document_to_aio(read_csv_document(path))
    if extension == ".xlsx":
        return spreadsheet_document_to_aio(read_xlsx_document(path))
    if extension == ".docx":
        return document_document_to_aio(read_docx_document(path))
    if extension in PANDOC_DOCUMENT_EXTENSIONS:
        return document_document_to_aio(read_pandoc_document(path))
    if extension in DOCLING_DOCUMENT_EXTENSIONS:
        return document_document_to_aio(read_docling_document(path))
    if extension == ".pptx":
        return presentation_document_to_aio(read_pptx_document(path))
    if extension == ".xls":
        raise RuntimeError("Legacy .xls files are not supported yet. Save or convert the workbook as .xlsx.")
    raise RuntimeError("Unsupported format for AIO inspection.")


def apply_aio(path: str, payload: Dict[str, object]) -> Dict[str, object]:
    ensure_aio_payload(payload)
    payload_source_path = extract_aio_source_path(payload)
    if payload_source_path:
        try:
            normalized_payload_source_path = resolve_workspace_path(payload_source_path)
        except ValueError as error:
            raise RuntimeError(str(error))
        if posixpath.normpath(normalized_payload_source_path) != posixpath.normpath(path):
            raise RuntimeError("AIO payload source path does not match the requested path.")
    extension = split_extension(path)
    if extension in {".csv", ".xlsx"}:
        return save_spreadsheet(path, spreadsheet_payload_from_aio(payload))
    if extension == ".docx":
        return save_document(path, document_payload_from_aio(payload))
    if extension in PANDOC_DOCUMENT_EXTENSIONS:
        return save_document(path, document_payload_from_aio(payload))
    if extension in DOCLING_DOCUMENT_EXTENSIONS:
        raise RuntimeError("PDF and image AIO application is not supported yet. These formats are read-only.")
    if extension == ".pptx":
        return save_presentation(path, payload)
    if extension == ".xls":
        raise RuntimeError("Legacy .xls files are not supported yet. Save or convert the workbook as .xlsx.")
    raise RuntimeError("Unsupported format for AIO application.")


def inspect_spreadsheet(path: str) -> Dict[str, object]:
    extension = split_extension(path)
    if extension == ".csv":
        return attach_aio_projection(read_csv_document(path))
    if extension == ".xlsx":
        return attach_aio_projection(read_xlsx_document(path))
    if extension == ".xls":
        raise RuntimeError("Legacy .xls files are not supported yet. Save or convert the workbook as .xlsx.")
    raise RuntimeError("Unsupported spreadsheet format. Use .xlsx or .csv.")


def save_spreadsheet(path: str, payload: Dict[str, object]) -> Dict[str, object]:
    if isinstance(payload, dict) and payload.get("spec") == AIO_SPEC:
        payload = spreadsheet_payload_from_aio(payload)
    extension = split_extension(path)
    sheets = payload.get("sheets") or []
    expected_etag = payload.get("expectedEtag")
    if extension == ".csv":
        return attach_aio_projection(write_csv_document(path, sheets, expected_etag))
    if extension == ".xlsx":
        return attach_aio_projection(write_xlsx_document(path, sheets, expected_etag))
    if extension == ".xls":
        raise RuntimeError("Legacy .xls files are not supported yet. Save or convert the workbook as .xlsx.")
    raise RuntimeError("Unsupported spreadsheet format. Use .xlsx or .csv.")


def normalize_spreadsheet(path: str) -> Dict[str, object]:
    extension = split_extension(path)
    if extension == ".csv":
        return attach_aio_projection(read_csv_document(path))
    if extension == ".xlsx":
        document = read_xlsx_document(path)
        if document.get("exists") and legacy_minimal_xlsx(path):
            return attach_aio_projection(write_xlsx_document(path, document.get("sheets") or [], document.get("etag")))
        return attach_aio_projection(document)
    if extension == ".xls":
        raise RuntimeError("Legacy .xls files are not supported yet. Save or convert the workbook as .xlsx.")
    raise RuntimeError("Unsupported spreadsheet format. Use .xlsx or .csv.")


def inspect_document(path: str) -> Dict[str, object]:
    extension = split_extension(path)
    if extension == ".docx":
        return attach_aio_projection(read_docx_document(path))
    if extension in PANDOC_DOCUMENT_EXTENSIONS:
        return attach_aio_projection(read_pandoc_document(path))
    if extension in DOCLING_DOCUMENT_EXTENSIONS:
        return attach_aio_projection(read_docling_document(path))
    raise RuntimeError("Unsupported document format. Use .docx, markdown/text formats, PDF, or supported images.")


def save_document(path: str, payload: Dict[str, object]) -> Dict[str, object]:
    if isinstance(payload, dict) and payload.get("spec") == AIO_SPEC:
        payload = document_payload_from_aio(payload)
    extension = split_extension(path)
    paragraphs = payload.get("paragraphs") or []
    expected_etag = payload.get("expectedEtag")
    blocks = [block for block in payload.get("blocks") or [] if isinstance(block, dict)]
    if extension == ".docx":
        if not blocks:
            blocks = [
                {
                    "id": f"block:{index}",
                    "kind": "paragraph",
                    "text": str(value),
                }
                for index, value in enumerate(paragraphs, start=1)
            ]
        return attach_aio_projection(write_docx_document(path, blocks, expected_etag))
    if extension in PANDOC_DOCUMENT_EXTENSIONS:
        if not blocks:
            blocks = [
                {
                    "id": f"block:{index}",
                    "kind": "paragraph",
                    "text": str(value),
                }
                for index, value in enumerate(paragraphs, start=1)
            ]
        return write_pandoc_document(path, blocks, expected_etag)
    if extension in DOCLING_DOCUMENT_EXTENSIONS:
        raise RuntimeError("PDF and image documents are read-only in entropic-office.")
    raise RuntimeError("Unsupported document format. Use .docx or a supported text document format.")


def save_presentation(path: str, payload: Dict[str, object]) -> Dict[str, object]:
    if isinstance(payload, dict) and payload.get("spec") == AIO_SPEC:
        payload = presentation_payload_from_aio(payload)
    extension = split_extension(path)
    slides = payload.get("slides") or []
    expected_etag = payload.get("expectedEtag")
    if extension == ".pptx":
        return attach_aio_projection(write_pptx_document(path, slides, expected_etag))
    raise RuntimeError("Unsupported presentation format. Use .pptx.")


def todo_spreadsheet_payload(items: Iterable[str]) -> Dict[str, object]:
    rows = [
        {"ref": "A1", "row": 1, "col": 1, "value": "Task", "kind": "string"},
        {"ref": "B1", "row": 1, "col": 2, "value": "Status", "kind": "string"},
    ]
    for index, item in enumerate(items, start=2):
        rows.append({"ref": cell_ref(index, 1), "row": index, "col": 1, "value": str(item), "kind": "string"})
        rows.append({"ref": cell_ref(index, 2), "row": index, "col": 2, "value": "Todo", "kind": "string"})
    return {"sheets": [{"name": "Sheet1", "cells": rows}]}


def blank_spreadsheet_payload() -> Dict[str, object]:
    return {"sheets": [{"name": "Sheet1", "cells": []}]}


def blank_document_payload(lines: Iterable[str]) -> Dict[str, object]:
    return {"paragraphs": [str(value) for value in lines]}


def request_desktop_action(payload: Dict[str, object]) -> Dict[str, object]:
    if not isinstance(payload, dict):
        raise RuntimeError("Desktop action payload must be a JSON object.")
    action_type = payload.get("type")
    allowed_types = {
        "open_workspace_file",
        "open_workspace_folder",
        "focus_window",
        "close_window",
    }
    if action_type not in allowed_types:
        raise RuntimeError("Unsupported desktop action type for the sandbox-to-desktop bridge.")

    action = dict(payload)
    if action_type in {"open_workspace_file", "open_workspace_folder"}:
        path = workspace_relative_path(str(action.get("path") or ""))
        if action_type == "open_workspace_file" and not path:
            raise RuntimeError("Workspace file path is required.")
        action["path"] = path

    os.makedirs(DESKTOP_ACTION_QUEUE_DIR, exist_ok=True)
    final_name = f"{int(time.time() * 1000)}-{os.getpid()}-{secrets.token_hex(8)}.json"
    final_path = os.path.join(DESKTOP_ACTION_QUEUE_DIR, final_name)
    fd, temp_path = tempfile.mkstemp(prefix=".desktop-action-", suffix=".tmp", dir=DESKTOP_ACTION_QUEUE_DIR)
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as handle:
            json.dump(action, handle, ensure_ascii=False, separators=(",", ":"))
            handle.write("\n")
        os.replace(temp_path, final_path)
    except Exception:
        try:
            os.unlink(temp_path)
        except OSError:
            pass
        raise
    return {"queued": True, "action": action}


def cli_usage() -> str:
    return """Usage:
  # Preferred Office/AIO workflow for .xlsx / .docx / .pptx and text documents
  # 1. Inspect the file:
  entropic-office api inspect-aio <path>
  # 2. Edit only the returned object member.
  # 3. Apply the full AIO payload back to the same workspace path:
  entropic-office api apply-aio <path>
  #
  # The object member is the canonical editable state. analysis, annotations,
  # projections, and artifacts are context only; do not write edits there.
  #
  # Spreadsheets use object.worksheets. Preserve formulas, dates, number-like
  # values, freeze panes, filters, dimensions, and merged ranges where possible.
  # Prefer formulas and totals over static-only spreadsheet output.
  #
  # Documents use object.blocks. Preserve headings, paragraphs, lists, tables,
  # links, and common inline marks where the backend exposes them.
  #
  # Presentations use object.slides with layout refs, title/body summaries,
  # objects[], frame, style, image_ref, table, chart, ordering, and notes.
  #
  # Markdown / HTML / rst / org / asciidoc / latex / txt route through the
  # Pandoc text adapter when available. PDF and image inspection route through
  # the Docling adapter and are currently read-only.
  #
  # After creating or editing an Office file, return a workspace-relative link,
  # or queue a safe desktop open request:
  entropic-office desktop open <path>
  # Advanced low-risk action queue API:
  printf '{"type":"open_workspace_file","path":"file.xlsx"}' | entropic-office api request-desktop-action

  # Legacy compatibility helpers
  entropic-office spreadsheet new <path>
  entropic-office spreadsheet todo <path> <item> [<item> ...]
  entropic-office document new <path>
  entropic-office document lines <path> <line> [<line> ...]

  # Low-level legacy APIs
  entropic-office api inspect-spreadsheet <path>
  entropic-office api normalize-spreadsheet <path>
  entropic-office api save-spreadsheet <path>
  entropic-office api inspect-document <path>
  entropic-office api save-document <path>
"""


def emit_json(payload: Dict[str, object]) -> None:
    sys.stdout.write(json.dumps(payload, ensure_ascii=False))
    sys.stdout.write("\n")


def run_api(argv: List[str]) -> Dict[str, object]:
    if len(argv) >= 3 and argv[2] == "request-desktop-action":
        if len(argv) >= 4:
            payload = json.loads(argv[3])
        else:
            payload = json.load(sys.stdin)
        return request_desktop_action(payload)
    if len(argv) < 4:
        raise RuntimeError(cli_usage().strip())
    command = argv[2]
    path = resolve_workspace_path(argv[3])
    if command == "inspect-spreadsheet":
        return inspect_spreadsheet(path)
    if command == "normalize-spreadsheet":
        return normalize_spreadsheet(path)
    if command == "save-spreadsheet":
        payload = json.load(sys.stdin)
        return save_spreadsheet(path, payload)
    if command == "inspect-document":
        return inspect_document(path)
    if command == "save-document":
        payload = json.load(sys.stdin)
        return save_document(path, payload)
    if command == "inspect-aio":
        return inspect_aio(path)
    if command == "apply-aio":
        payload = json.load(sys.stdin)
        return apply_aio(path, payload)
    raise RuntimeError(cli_usage().strip())


def run_cli(argv: List[str]) -> Dict[str, object]:
    if len(argv) < 4:
        raise RuntimeError(cli_usage().strip())
    category = argv[1]
    command = argv[2]
    if category == "desktop" and command == "open":
        return request_desktop_action({"type": "open_workspace_file", "path": argv[3]})
    path = resolve_workspace_path(argv[3])
    if category == "spreadsheet" and command == "new":
        return save_spreadsheet(path, blank_spreadsheet_payload())
    if category == "spreadsheet" and command == "todo":
        items = argv[4:] or [f"Item {index}" for index in range(1, 11)]
        return save_spreadsheet(path, todo_spreadsheet_payload(items))
    if category == "document" and command == "new":
        return save_document(path, blank_document_payload([]))
    if category == "document" and command == "lines":
        return save_document(path, blank_document_payload(argv[4:]))
    raise RuntimeError(cli_usage().strip())


def main() -> int:
    try:
        if len(sys.argv) < 2:
            raise RuntimeError(cli_usage().strip())
        if sys.argv[1] == "api":
            result = run_api(sys.argv)
        else:
            result = run_cli(sys.argv)
        emit_json({"ok": True, "result": result})
        return 0
    except Exception as error:
        emit_json({"ok": False, "error": str(error)})
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
